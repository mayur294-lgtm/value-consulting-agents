"""frontline-builders component evaluator (objective, code-only — no LLM).

Covers the three Frontline 2026 deliverable builders (ticket #200):
`tools/frontline_slides_pptx.py` (BackbaseSlidesPresenter), `tools/
frontline_2026_presenter.py` (Frontline2026Presenter) — both PPTX — and
`tools/frontline_2026_html.py` (Frontline2026HTML). All three produce
CLIENT-FACING deliverables; a builder that emits a broken file fails silently
until a consultant opens it in front of a client. Every check GENERATES a
real artifact from an in-memory fixture deck (a handful of add_* calls — no
new golden fixture file needed) and inspects the REAL output, matching the
executable-tier bar `rubrics/component/roi_excel_generator.py` sets:

  - pptx_slides_file_opens_with_expected_slides / pptx_presenter_file_opens_with_expected_slides
        generate a real .pptx via the builder, reopen it with python-pptx,
        assert the slide count and read back real slide text. Mutation:
        drop one add_*() call from the fixture deck and confirm the slide
        count no longer matches — proves the assertion is load-bearing, not
        vacuously true for any deck.
  - pptx_slides_brand_tokens_match / pptx_presenter_brand_tokens_match
        the builder's own color constants (the literal surface every shape's
        fill/font color is drawn from) must equal the canonical hex in
        presentations/frontline-2026/design-tokens.json. Mutation:
        monkeypatch one constant to an off-palette hex, confirm the check
        reddens, restore it and confirm it goes green again.
  - html_self_contained_no_cdn
        the same self_contained_no_cdn contract `rubrics/deliverable/decks.py`
        hard-fails on for finished decks, applied to the BUILDER's raw
        output: only fonts.googleapis.com / fonts.gstatic.com references are
        allowed. Mutation: inject a non-font external <script src> and
        confirm it reddens.
  - html_brand_tokens_root_match
        the generated HTML's `:root { --navy: ...; }` custom-property block
        — the single place Frontline2026HTML defines its palette — must
        equal design-tokens.json. Mutation: substitute an off-palette hex
        into the generated `:root` text and confirm it reddens.

ONE ROW FOR ALL THREE BUILDERS, not three separate rows (design Open Item
3). Reasoning: all three are one design-system surface (same token source,
same "produces a client deliverable" failure mode) authored and reviewed
together, not three independently-lifecycled components — splitting them
would 3x the registry/threshold/mutation bookkeeping for what is one
ticket's unit of work, the same shape roi-excel-generator uses (one row, six
checks, one fixture concern). Each builder still gets its OWN uniquely-named
checks so a regression in exactly one builder is attributable from the check
name alone, without conflating a PPTX bug with an HTML bug.

PYTHON-PPTX DEPENDENCY (#181 SKIP-as-FAIL): evals/requirements.txt does not
pin python-pptx (only PyYAML — see its own header comment: "Core runner uses
only the stdlib + PyYAML... run in CI with zero Langfuse/Anthropic setup").
The two PPTX checks therefore guard their `import pptx` and, if it is
missing, return a normal FAILING CheckResult (hard_fail=True, NOT
skipped=True) naming exactly what's missing and how to fix it — never a
silently-passing or `skipped=True` result, which #181 explicitly excludes
from scoring and would let an unverifiable builder through as if it were
verified clean. The HTML checks have no such dependency and always run.

target: path to presentations/frontline-2026/design-tokens.json — the
canonical brand-token source every check here loads from. A missing/invalid
tokens file fails every check immediately (no tokens, no way to verify
brand parity — that is a real failure, not something to skip past).
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from rubrics.base import CheckResult, repo_root
from rubrics._harness import bool_check, run_in_tmpdir

_SLIDES_EXPECTED_COUNT = 4
_PRESENTER_EXPECTED_COUNT = 4

_EXTERNAL_REF_RE = re.compile(r'(?:src|href)\s*=\s*["\'](https?://[^"\']+)["\']', re.I)
_ALLOWED_EXTERNAL_RE = re.compile(r"fonts\.(googleapis|gstatic)\.com", re.I)

_SLIDES_TOKEN_MAP = {
    "NAVY": "navy", "BLUE": "blue", "CYAN": "cyan", "RED": "red",
    "WHITE": "white", "OFF_WHITE": "off_white", "GREEN": "green",
    "BORDER": "border", "BLUE_DARK": "blue_dark", "LIGHT_BLUE": "blue_light",
    "TEXT_MUTED": "text_muted",
}

_PRESENTER_TOKEN_MAP = {
    "PRIMARY_NAVY": "navy", "ACTION_BLUE": "blue", "SURFACE_WHITE": "white",
    "BG_GRAY": "off_white", "TEXT_MAIN": "text_main", "TEXT_MUTED": "text_muted",
    "SUCCESS_GREEN": "green", "SEMANTIC_RED": "red",
}

_HTML_ROOT_TOKEN_MAP = {
    "--navy": "navy", "--action-blue": "blue", "--white": "white",
    "--bg-gray": "off_white", "--text-main": "text_main",
    "--text-muted": "text_muted", "--green": "green",
}


def _bc(name: str, ok: bool, *, detail: str, exercised: str, hard_fail: bool = True) -> CheckResult:
    result = bool_check(name, ok, detail=detail, hard_fail=hard_fail)
    result.exercised = exercised
    return result


def _ensure_repo_on_path() -> Path:
    root = repo_root()
    if str(root) not in sys.path:
        sys.path.insert(0, str(root))
    return root


def _load_tokens(tokens_path: str) -> dict:
    data = json.loads(Path(tokens_path).read_text())
    return {name: spec["hex"].lstrip("#").lower() for name, spec in data["colors"].items()}


def _pptx_available() -> tuple[bool, str]:
    try:
        import pptx  # noqa: F401
        return True, ""
    except ImportError as e:
        return False, str(e)


def _require_pptx(name: str, exercised: str) -> CheckResult | None:
    """#181 SKIP-as-FAIL: python-pptx missing -> a normal FAIL, never skipped=True."""
    available, err = _pptx_available()
    if available:
        return None
    return _bc(name, False,
               detail=f"python-pptx not installed ({err}) — a client-facing PPTX builder cannot be "
                      f"verified without it, so this MUST fail the gate rather than pass silently. "
                      f"Install via the pin already used elsewhere in this repo: python-pptx==1.0.2 "
                      f"(see requirements.txt); evals/requirements.txt does not currently carry it.",
               exercised=exercised)


def _rgb_hex(rgbcolor) -> str:
    return str(rgbcolor).lower()


# ─────────────────────────────────────────────────────────────────────
# frontline_slides_pptx.BackbaseSlidesPresenter
# ─────────────────────────────────────────────────────────────────────

def _build_slides_deck(cls, *, thank_you: bool = True):
    p = cls(title="Eval Fixture Deck")
    p.add_cover_color_block(label="EVAL", title="Fixture Cover", date="2026")
    p.add_toc(label="AGENDA", title="Contents", items=["Intro", "Body", "Close"])
    p.add_content_standard(theme="light", label="CONTENT", title="Content Slide",
                            subtitle="sub", body="body text")
    if thank_you:
        p.add_thank_you()
    return p


def _pptx_slides_file_opens_with_expected_slides(td: Path, target: str) -> CheckResult:
    name = "pptx_slides_file_opens_with_expected_slides"
    exercised = "tools.frontline_slides_pptx.BackbaseSlidesPresenter -> save() -> python-pptx reopen"
    guard = _require_pptx(name, exercised)
    if guard:
        return guard
    _ensure_repo_on_path()
    from pptx import Presentation as ReadPresentation
    from tools.frontline_slides_pptx import BackbaseSlidesPresenter

    out = td / "eval_slides.pptx"
    try:
        _build_slides_deck(BackbaseSlidesPresenter).save(str(out))
        reopened = ReadPresentation(str(out))
    except Exception as e:  # noqa: BLE001 - a generator crash is a clean check failure
        return _bc(name, False, detail=f"generation/reopen raised: {type(e).__name__}: {e}", exercised=exercised)

    slide_count = len(reopened.slides)
    cover_texts = [sh.text_frame.text for sh in reopened.slides[0].shapes
                   if sh.has_text_frame and sh.text_frame.text]
    cover_ok = any("Fixture Cover" in t for t in cover_texts)
    if not (out.exists() and slide_count == _SLIDES_EXPECTED_COUNT and cover_ok):
        return _bc(name, False,
                   detail=f"expected {_SLIDES_EXPECTED_COUNT} slides with cover title readable back; "
                          f"got {slide_count} slides, cover_texts={cover_texts}",
                   exercised=exercised)

    # Mutation proof: drop one add_*() call — the reopened slide count must
    # no longer match, proving the count assertion is load-bearing.
    out_mut = td / "eval_slides_mutated.pptx"
    _build_slides_deck(BackbaseSlidesPresenter, thank_you=False).save(str(out_mut))
    mutated_count = len(ReadPresentation(str(out_mut)).slides)
    if mutated_count == _SLIDES_EXPECTED_COUNT:
        return _bc(name, False,
                   detail=f"mutation proof failed: dropping add_thank_you() still produced "
                          f"{mutated_count} slides",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"{out.name}: {slide_count} slides opened via python-pptx, cover title read back "
                      f"({cover_texts}); dropping one add_*() call correctly changed the count to "
                      f"{mutated_count}",
               exercised=exercised)


def _pptx_slides_brand_tokens_match(td: Path, target: str) -> CheckResult:
    name = "pptx_slides_brand_tokens_match"
    exercised = "tools.frontline_slides_pptx.BackbaseSlidesPresenter color constants vs design-tokens.json"
    guard = _require_pptx(name, exercised)
    if guard:
        return guard
    _ensure_repo_on_path()
    from pptx.dml.color import RGBColor
    from tools.frontline_slides_pptx import BackbaseSlidesPresenter as Cls

    tokens = _load_tokens(target)
    mismatches = []
    for attr, token_name in _SLIDES_TOKEN_MAP.items():
        expected = tokens.get(token_name)
        actual = _rgb_hex(getattr(Cls, attr))
        if expected is None:
            mismatches.append(f"token '{token_name}' missing from design-tokens.json")
        elif actual != expected:
            mismatches.append(f"{attr}=#{actual} != tokens.{token_name}=#{expected}")
    if mismatches:
        return _bc(name, False, detail="; ".join(mismatches), exercised=exercised)

    # Mutation proof: monkeypatch NAVY to an off-palette hex, confirm it
    # reddens, then restore and confirm it goes green again.
    original_navy = Cls.NAVY
    try:
        Cls.NAVY = RGBColor(0x12, 0x34, 0x56)
        mutated_matches = _rgb_hex(Cls.NAVY) == tokens["navy"]
    finally:
        Cls.NAVY = original_navy
    if mutated_matches:
        return _bc(name, False, detail="mutation proof failed: off-palette NAVY (#123456) was not detected",
                   exercised=exercised)
    restored_matches = _rgb_hex(Cls.NAVY) == tokens["navy"]
    if not restored_matches:
        return _bc(name, False, detail="restore failed: NAVY did not return to on-palette after mutation",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"{len(_SLIDES_TOKEN_MAP)} color constant(s) match design-tokens.json exactly "
                      f"({sorted(_SLIDES_TOKEN_MAP)}); off-palette NAVY mutation correctly reddened and "
                      f"restore correctly greened",
               exercised=exercised)


# ─────────────────────────────────────────────────────────────────────
# frontline_2026_presenter.Frontline2026Presenter
# ─────────────────────────────────────────────────────────────────────

def _build_presenter_deck(cls, output_path: str, *, divider: bool = True):
    p = cls(output_path)
    p.add_cover_slide("EVAL", "Fixture Cover", "2026")
    p.add_agenda_slide("EVAL", "Agenda", "Eval Co", ["Intro", "Body", "Close"])
    p.add_content_slide("Content Slide", "sub", ["point 1"])
    if divider:
        p.add_section_divider("EVAL", "Divider")
    return p


def _pptx_presenter_file_opens_with_expected_slides(td: Path, target: str) -> CheckResult:
    name = "pptx_presenter_file_opens_with_expected_slides"
    exercised = "tools.frontline_2026_presenter.Frontline2026Presenter -> save() -> python-pptx reopen"
    guard = _require_pptx(name, exercised)
    if guard:
        return guard
    _ensure_repo_on_path()
    from pptx import Presentation as ReadPresentation
    from tools.frontline_2026_presenter import Frontline2026Presenter

    out = td / "eval_presenter.pptx"
    try:
        _build_presenter_deck(Frontline2026Presenter, str(out)).save(str(out))
        reopened = ReadPresentation(str(out))
    except Exception as e:  # noqa: BLE001 - a generator crash is a clean check failure
        return _bc(name, False, detail=f"generation/reopen raised: {type(e).__name__}: {e}", exercised=exercised)

    slide_count = len(reopened.slides)
    cover_texts = [sh.text_frame.text for sh in reopened.slides[0].shapes
                   if sh.has_text_frame and sh.text_frame.text]
    cover_ok = any("Fixture Cover" in t for t in cover_texts)
    if not (out.exists() and slide_count == _PRESENTER_EXPECTED_COUNT and cover_ok):
        return _bc(name, False,
                   detail=f"expected {_PRESENTER_EXPECTED_COUNT} slides with cover title readable back; "
                          f"got {slide_count} slides, cover_texts={cover_texts}",
                   exercised=exercised)

    out_mut = td / "eval_presenter_mutated.pptx"
    _build_presenter_deck(Frontline2026Presenter, str(out_mut), divider=False).save(str(out_mut))
    mutated_count = len(ReadPresentation(str(out_mut)).slides)
    if mutated_count == _PRESENTER_EXPECTED_COUNT:
        return _bc(name, False,
                   detail=f"mutation proof failed: dropping add_section_divider() still produced "
                          f"{mutated_count} slides",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"{out.name}: {slide_count} slides opened via python-pptx, cover title read back "
                      f"({cover_texts}); dropping one add_*() call correctly changed the count to "
                      f"{mutated_count}",
               exercised=exercised)


def _pptx_presenter_brand_tokens_match(td: Path, target: str) -> CheckResult:
    name = "pptx_presenter_brand_tokens_match"
    exercised = "tools.frontline_2026_presenter module-level color constants vs design-tokens.json"
    guard = _require_pptx(name, exercised)
    if guard:
        return guard
    _ensure_repo_on_path()
    import importlib
    from pptx.dml.color import RGBColor
    presenter_mod = importlib.import_module("tools.frontline_2026_presenter")

    tokens = _load_tokens(target)
    mismatches = []
    for attr, token_name in _PRESENTER_TOKEN_MAP.items():
        expected = tokens.get(token_name)
        actual = _rgb_hex(getattr(presenter_mod, attr))
        if expected is None:
            mismatches.append(f"token '{token_name}' missing from design-tokens.json")
        elif actual != expected:
            mismatches.append(f"{attr}=#{actual} != tokens.{token_name}=#{expected}")
    if mismatches:
        return _bc(name, False, detail="; ".join(mismatches), exercised=exercised)

    original_muted = presenter_mod.TEXT_MUTED
    try:
        presenter_mod.TEXT_MUTED = RGBColor(0x99, 0x00, 0x99)
        mutated_matches = _rgb_hex(presenter_mod.TEXT_MUTED) == tokens["text_muted"]
    finally:
        presenter_mod.TEXT_MUTED = original_muted
    if mutated_matches:
        return _bc(name, False,
                   detail="mutation proof failed: off-palette TEXT_MUTED (#990099) was not detected",
                   exercised=exercised)
    restored_matches = _rgb_hex(presenter_mod.TEXT_MUTED) == tokens["text_muted"]
    if not restored_matches:
        return _bc(name, False,
                   detail="restore failed: TEXT_MUTED did not return to on-palette after mutation",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"{len(_PRESENTER_TOKEN_MAP)} color constant(s) match design-tokens.json exactly "
                      f"({sorted(_PRESENTER_TOKEN_MAP)}); off-palette TEXT_MUTED mutation correctly "
                      f"reddened and restore correctly greened",
               exercised=exercised)


# ─────────────────────────────────────────────────────────────────────
# frontline_2026_html.Frontline2026HTML
# ─────────────────────────────────────────────────────────────────────

def _build_html_deck(cls):
    b = cls(title="Eval Fixture Deck")
    b.add_cover("EVAL", "Fixture Cover", "2026")
    b.add_section_divider("EVAL", "Divider", "tagline")
    b.add_content("Content Slide", "sub", ["point 1", "point 2"])
    return b


def _html_self_contained_no_cdn(td: Path, target: str) -> CheckResult:
    name = "html_self_contained_no_cdn"
    exercised = "tools.frontline_2026_html.Frontline2026HTML(...).render()"
    _ensure_repo_on_path()
    try:
        from tools.frontline_2026_html import Frontline2026HTML
    except ImportError as e:
        return _bc(name, False, detail=f"import failed: {e}", exercised=exercised)

    html = _build_html_deck(Frontline2026HTML).render()
    refs = _EXTERNAL_REF_RE.findall(html)
    bad = [u for u in refs if not _ALLOWED_EXTERNAL_RE.search(u)]
    if bad:
        return _bc(name, False, detail=f"{len(bad)} non-Google-Fonts external reference(s): {bad}",
                   exercised=exercised)
    if not refs:
        return _bc(name, False,
                   detail="no external references found at all — fixture deck may not be exercising "
                          "the font-link code path; treating as inconclusive rather than a false pass",
                   exercised=exercised)

    # Mutation proof: inject a non-font external <script src>, confirm it reddens.
    mutated = html.replace("</head>", '<script src="https://evil.cdn.example.com/x.js"></script></head>', 1)
    mutated_bad = [u for u in _EXTERNAL_REF_RE.findall(mutated) if not _ALLOWED_EXTERNAL_RE.search(u)]
    if not mutated_bad:
        return _bc(name, False,
                   detail="mutation proof failed: injected https://evil.cdn.example.com/x.js was not detected",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"self-contained: {len(refs)} external reference(s), all Google Fonts "
                      f"({refs}); injected non-font <script src> mutation correctly caught",
               exercised=exercised)


def _parse_root_vars(html: str) -> dict:
    m = re.search(r":root\s*\{([^}]*)\}", html, re.S)
    if not m:
        return {}
    return {var_m.group(1): var_m.group(2).lstrip("#").lower()
            for var_m in re.finditer(r"(--[a-z-]+)\s*:\s*(#[0-9a-fA-F]{6})", m.group(1))}


def _html_brand_tokens_root_match(td: Path, target: str) -> CheckResult:
    name = "html_brand_tokens_root_match"
    exercised = "tools.frontline_2026_html.Frontline2026HTML(...).render() :root block vs design-tokens.json"
    _ensure_repo_on_path()
    try:
        from tools.frontline_2026_html import Frontline2026HTML
    except ImportError as e:
        return _bc(name, False, detail=f"import failed: {e}", exercised=exercised)

    tokens = _load_tokens(target)
    html = _build_html_deck(Frontline2026HTML).render()
    root_vars = _parse_root_vars(html)
    mismatches = []
    for css_var, token_name in _HTML_ROOT_TOKEN_MAP.items():
        expected = tokens.get(token_name)
        actual = root_vars.get(css_var)
        if actual is None:
            mismatches.append(f"{css_var} missing from generated :root block")
        elif expected is None:
            mismatches.append(f"token '{token_name}' missing from design-tokens.json")
        elif actual != expected:
            mismatches.append(f"{css_var}=#{actual} != tokens.{token_name}=#{expected}")
    if mismatches:
        return _bc(name, False, detail="; ".join(mismatches), exercised=exercised)

    # Mutation proof: substitute an off-palette hex for --navy in the
    # GENERATED text (in-memory only — no source file touched), re-parse.
    mutated_html = re.sub(r"(--navy:\s*)#041326", r"\g<1>#123456", html, count=1)
    if mutated_html == html:
        return _bc(name, False,
                   detail="mutation setup failed: could not locate '--navy: #041326' in the generated :root",
                   exercised=exercised)
    mutated_vars = _parse_root_vars(mutated_html)
    if mutated_vars.get("--navy") == tokens["navy"]:
        return _bc(name, False, detail="mutation proof failed: off-palette --navy (#123456) was not detected",
                   exercised=exercised)

    return _bc(name, True,
               detail=f"{len(_HTML_ROOT_TOKEN_MAP)} :root custom propert(y/ies) match design-tokens.json "
                      f"exactly ({sorted(_HTML_ROOT_TOKEN_MAP)}); off-palette --navy mutation correctly "
                      f"reddened",
               exercised=exercised)


CHECKS = {
    "pptx_slides_file_opens_with_expected_slides": _pptx_slides_file_opens_with_expected_slides,
    "pptx_slides_brand_tokens_match": _pptx_slides_brand_tokens_match,
    "pptx_presenter_file_opens_with_expected_slides": _pptx_presenter_file_opens_with_expected_slides,
    "pptx_presenter_brand_tokens_match": _pptx_presenter_brand_tokens_match,
    "html_self_contained_no_cdn": _html_self_contained_no_cdn,
    "html_brand_tokens_root_match": _html_brand_tokens_root_match,
}


def evaluate(target: str) -> list[CheckResult]:
    """target: path to presentations/frontline-2026/design-tokens.json — the
    canonical brand-token source every check loads from directly."""
    p = Path(target)
    if not p.exists():
        return [CheckResult(name, 0.0, False, hard_fail=True,
                            detail=f"design tokens file not found: {target}") for name in CHECKS]
    try:
        json.loads(p.read_text())
    except json.JSONDecodeError as e:
        return [CheckResult(name, 0.0, False, hard_fail=True, detail=f"invalid JSON: {e}") for name in CHECKS]
    return [run_in_tmpdir(fn, target) for fn in CHECKS.values()]
