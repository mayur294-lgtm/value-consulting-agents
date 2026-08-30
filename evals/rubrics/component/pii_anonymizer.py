"""pii-anonymizer component evaluator — deterministic regression coverage for
the Presidio PII engine (ticket #161, PRD v6 §5, solution-design-v6.md D9).

`scripts/pii/engine.py` (Presidio detection + reversible pseudonymisation) and
`scripts/anonymize_transcript.py` (its public facade) are what every agent's
input text passes through before it ever reaches the model — this is the
acceptance gate for the whole v6 cycle. threshold: 1.00 in the registry
(deliberate deviation from the house 0.80 default) — a privacy control is
pass/fail, not "mostly correct." No `judge:` entries — every check here is
deterministic and free.

SEQUENCING (read before adding a check here)
  #161 runs 7th in the v6 build order, deliberately: the gate exists before
  more code lands on top of it. Six checks from the PRD's original 16-check
  list were NOT authored then because their subject did not exist yet. Four
  have since landed — `document_formats_converted_and_scrubbed` with #162,
  `image_input_produces_sidecar_and_redacted_copy` with #163,
  `image_unreadable_script_refuses_and_writes_nothing` with #173 (closing the
  non-Latin-script leak #163 measured and documented but did not fix), and
  `guard_fails_closed_on_inputs_path` with #164 — and are checks 11-14
  below. The rest still wait on their ticket:

    document_formats_converted_and_scrubbed   -> #162 — LANDED (check 11)
    image_input_produces_sidecar_and_redacted_copy -> #163 — LANDED (check 12)
    image_unreadable_script_refuses_and_writes_nothing -> #173 — LANDED (check 13)
    guard_fails_closed_on_inputs_path         -> #164 — LANDED (check 14)
    internal_domain_email_redacted_no_over_detection -> #181 — LANDED (check 15)
    nested_outputs_deanonymized               -> #165 — LANDED (check 16)
    xlsx_outputs_deanonymized                 -> #165 — LANDED (check 17)
    workspace_paths_contain_no_client_identifiers -> #166 — LANDED (check 18)
    mcp_query_client_name_blocked             -> already covered by the
                                                  separate `mcp-query-guard`
                                                  registry row (15 checks) —
                                                  intentionally NOT duplicated
                                                  here.

  A check for a module that does not exist yet must never be added as a
  silent pass or a SKIP — this repo has already shipped two gates that scored
  1.000 while certifying nothing (the path-2-only eval gate, and the
  two-line mcp-query-guard fixture pre-711b56c). The clean option, taken
  here, is simply not writing those checks until their ticket lands.

  Check 14, `guard_fails_closed_on_inputs_path`, differs in kind from the
  other 13: it does not import `scripts/pii/engine.py` at all. By design
  (.design/solution-design-v6.md D13), `.claude/hooks/anonymize-guard.py`
  was rewritten to NOT use the Presidio engine — a synchronous PreToolUse
  hook cannot pay a ~0.7-1.1s spaCy/Presidio cold start on every Read/Bash
  call. It invokes the real hook SCRIPT as a subprocess with a synthesized
  PreToolUse payload on stdin, exactly as `mcp_query_guard.py`'s rubric does
  for its own hook — never imports the hook module and monkeypatches its
  internals, since that would only prove the Python function behaves, not
  that the process-level contract (stdout JSON shape, exit code) holds.
  Since #196 both rows reach that subprocess through ONE implementation,
  `rubrics._harness.run_hook_subprocess`, instead of two copies that had
  already begun to diverge.

  The 18 checks below all exercise code that exists TODAY: `scripts/pii/
  engine.py`, `scripts/pii/denylist.py` (via the engine's deny-list
  recognizer), `scripts/pii/ingest.py` (#162, #163, #173),
  `scripts/anonymize_transcript.py`'s facade, `.claude/hooks/
  anonymize-guard.py` (#164), `scripts/artifact_boundary.py`'s
  `deanonymize_dir` exit gate (#165, checks 16-17), and
  `scripts/pii/identity.py`'s neutral workspace (#166, check 18).

  Check 18, `workspace_paths_contain_no_client_identifiers` (#166), is the
  only check here about PATHS rather than CONTENT — and the two are not
  substitutes. `compose_prompt` renders `engagement_dir` / `outputs_dir` /
  `transcript_path` into the invocation prompt as VALUES and `run_agent`
  sets `cwd` to the same client-named directory, so a transcript that
  scores 1.00 on every content check above still arrives with the client's
  name in the Runtime Parameters table (.design/solution-design-v6.md D6).
  Like check 14 it does not import the engine; unlike check 14 it imports
  its subject directly (`pii.identity` is stdlib-only, so there is no
  interpreter or fail-open reason to hold it at a subprocess boundary).

  Check 15, `internal_domain_email_redacted_no_over_detection` (#181),
  closes a production leak: Presidio's built-in `EMAIL_ADDRESS` recognizer
  validates the matched domain against REAL, REGISTERED TLDs via
  `tldextract`, so an address on a bank's own internal domain
  (`.internal`, `.corp`, `.local`, …) or an RFC 2606 reserved name
  (`.test`, `.example`, `.invalid`) was never recognised as an email at all
  and reached anonymised output in cleartext — found on a real DOCX where
  the person's name redacted and the internal-domain email next to it did
  not. The fix (`engine.py`'s `_InternalDomainEmailRecognizer`) detects on
  SHAPE instead of TLD registration; this check has two halves in ONE
  gate-bitable assertion: internal-domain addresses ARE redacted and
  restore byte-identical (the leak), and a set of plausible non-email
  `@`-bearing strings that are NOT shaped like `local-part@host.tld` — an
  npm-style `pkg@1.2.3` (numeric final label), a Python decorator, a
  bare `@` in prose — are NOT redacted (the over-detection guard — see the
  "denylist blocked the word 'all'" lesson in this file's own history).
  An `scp user@host.example:/path`-style remote path is DELIBERATELY
  excluded from that negative set: `user@host.example` is genuinely
  shape-identical to a real internal-domain email and the fix redacts it,
  which is documented as the correct, conservative call in engine.py's
  own "INTERNAL-DOMAIN EMAILS" note, not asserted as a false positive here.

  Checks 16-17, `nested_outputs_deanonymized` and `xlsx_outputs_
  deanonymized` (#165, salvaged from the closed PR #129 / issue #125),
  cover `scripts/artifact_boundary.py`'s `deanonymize_dir` — the EXIT gate
  where real client names re-enter deliverables, previously top-level-only
  (`outputs_dir.iterdir()`) and blind to generated `.xlsx` ROI models.
  `nested_outputs_deanonymized` proves the recursive walk (`rglob`) reaches
  a nested file, that dotfiles/dot-directories (`.anon_*`, `.pii_mapping.
  json`) and `interim*` stay excluded, AND both hazards #165 calls out by
  name: a legacy `[CLIENT]` placeholder restores while the UNRELATED
  `[CLIENT]_Business_Case_Questionnaire.xlsx` filename/prose template token
  (five other components use it) is left byte-identical, and `[CLIENT-ABBR]`
  is not corrupted by `[CLIENT]` (longest-placeholder-first, #159's
  guarantee confirmed on this path too). `xlsx_outputs_deanonymized` proves
  cell values, formula strings, and sheet titles restore; an unchanged
  workbook is NOT re-saved; and an unopenable workbook (or a missing
  `openpyxl`) is reported in `unrestored` with `client_ready: false` —
  never a silent skip, which is the exact defect (backlog.md) this ticket
  exists to close.

WHY THE VENV INTERPRETER
  `scripts/pii/engine.py` does `import presidio_analyzer` at module level, and
  Presidio needs Python 3.10-3.13 — the system `python3` here is 3.9.6 and
  cannot import it at all (see engine.py's own module docstring, "THE
  INTERPRETER SPLIT").

  Two different interpreter questions live in this file; do not conflate
  them. (a) The HOOK subprocesses (checks 14 and 19) run under whatever
  `.claude/settings.json` registers — bare `python3` — resolved by
  `rubrics._harness.registered_interpreter()`, never `sys.executable`
  (#192/backlog :116). That is deliberate: `anonymize-guard.py` is
  stdlib-only precisely so it can run under a consultant's 3.9.6.
  (b) This rubric ITSELF IMPORTS `scripts/pii/engine.py` directly, in
  process, because several checks (distinct-placeholder bijectivity, shared
  entity_mapping identity across two files, mapping-file bytes on disk) need
  to inspect engine internals that a subprocess boundary would hide. That
  means whatever interpreter runs `evals/run_experiment.py` for THIS
  component must itself be able to import Presidio — i.e. `.venv/bin/python`,
  not bare `python3`. CI resolves this the same way: `actions/setup-python`
  pins 3.11 and the workflow installs `requirements.txt` (which carries the
  Presidio/spaCy pins) before running the suite, so `sys.executable` in CI is
  already Presidio-capable — no separate venv step needed there. Locally,
  run this file with `.venv/bin/python evals/run_experiment.py --component
  pii-anonymizer`, not the system `python3`.

FIXTURE
  `evals/goldens/pii_roundtrip_fixture.md` — entirely invented, synthetic
  multi-stakeholder transcript (see its own header comment). It deliberately
  spans FOUR document shapes because PERSON detection is shape-dependent, not
  just fixture-dependent (see engine.py's "MEASURED DETECTION LIMITS" and
  #159's 30-name x 5-shape measurement): a markdown table row, an
  attendee-bullet list, a speaker-label line, and prose. A prose-only fixture
  scores 30/30 on PERSON and would certify that gap as passing — this one
  does not, on purpose (see `_person_detection_by_shape` and
  KNOWN_TABLE_SHAPE_MISS below).

  Two more fixtures are synthesised INLINE, in a tempdir, never touching the
  repo: `_cross_transcript_merge_collision_free` (two small transcripts
  sharing one email address, per #160's regression) and
  `_mapping_files_chmod_600_and_cleaned` (the same two-transcript shape,
  inspecting the mapping files `anonymize_transcript_file` writes to disk).
  Nothing here ever mutates the committed golden — it is read-only.

  #162 adds a THIRD inline fixture set: one PDF, DOCX, PPTX, XLSX and CSV
  built programmatically in a tempdir by `_build_document_fixtures`, all
  carrying the same planted PII. No binary is ever committed — a committed
  .docx or .xlsx is opaque to review and, if ever built from real content,
  is exactly the contamination the synthetic-quarantine cycle exists to
  stop. Four of the five carry a person's name inside a REAL TABLE
  (a Word table, a PowerPoint table, a worksheet, CSV rows), because that
  is the shape #162 was most likely to get wrong — see
  `_document_formats_converted_and_scrubbed`.

  #163 adds a FOURTH inline fixture set, on the same no-committed-binaries
  rule: a screenshot drawn with Pillow (`_build_screenshot`), plus a DOCX and
  a PPTX that embed it. The screenshot lays its person name out in COLUMNS,
  because that is the shape tesseract's own page segmentation shreds before
  the detector ever runs — a screenshot of prose would certify the easy case.
  It renders with `ImageFont.load_default(size=...)`, which ships inside
  Pillow, so it OCRs identically on a laptop and on the CI runner; a system
  font path would pass locally and quietly degrade in CI. The check asserts
  on the redacted copy's PIXELS — per OCR word, and on the BACKGROUND between
  the letters, since black fill over black glyphs is indistinguishable from
  the glyphs — so "a file was written" cannot pass it.

  #173 adds a FIFTH inline fixture: a screenshot carrying only
  `UNREADABLE_SCRIPT_NAME`, a name in a script the local `eng`-only OCR
  cannot read (module docstring's own measured Sinhala example), plus a
  second, genuinely textless screenshot (shapes, no text) to prove the two
  failure modes — "unreadable" vs "nothing there" — stay distinguishable.
  Same font-portability rule as #163's screenshot: `_image_font`, not a
  system font, so it OCRs identically on a laptop and in CI.

CLOSED IN #209 — the markdown-table PERSON miss, and the false reason for it
  `evals/goldens/pii_roundtrip_fixture.md`'s Stakeholder Directory table
  carries "Aisha Rahman" in a table cell because the PERSON recogniser does
  not detect it there. Until #209 `MUST_NOT_LEAK` deliberately EXCLUDED that
  one value, so this component scored 1.000 at threshold 1.00 while its own
  `detail` line printed `table=MISSED ('Aisha Rahman')` on every run. Visible
  is better than hidden, but an assertion set trimmed to fit the failure is
  still a threshold-shaped exclusion, and it is now gone: the value is in
  `MUST_NOT_LEAK` and the leak is closed.

  ⚠️ The documented reason for tolerating it was measurably FALSE, and the
  same false sentence had been copied into `engine.py` (twice) and
  `.prd/backlog.md`. It claimed spaCy tags the table-cell name ORGANIZATION,
  so catching it would mean enabling ORGANIZATION and stripping
  Backbase/Temenos. Measured on this fixture with `en_core_web_lg`: spaCy
  tags it **EVENT**, and enabling ORGANIZATION does not catch it at all while
  it does strip "Backbase" and "Salesforce" — full cost, zero benefit. There
  was no trade-off to weigh. ORGANIZATION stays disabled for that reason.

  THE FIX is the one engine.py always called "the real fix": STAKEHOLDER
  names on the engagement deny-list, where a term is model- and
  shape-independent and therefore fires in a table cell, a bullet, prose and
  a speaker label alike. `denylist.extract_stakeholder_terms` mines them from
  the engagement's own CLIENT_PROFILE.md / engagement_intake.md /
  ENGAGEMENT_CONTEXT.md.

  HOW THIS RUBRIC PROVES IT RATHER THAN ASSUMING IT — `resolved_deny_terms()`
  does not hard-code the stakeholder name. It seeds a fixture engagement with
  ordinary CLIENT_PROFILE.md / engagement_intake.md documents and resolves the
  deny-list through `denylist.resolve_engagement_deny_list`, the production
  entry point. Hard-coding "Aisha Rahman" into a list here would close the
  assertion without closing the leak, and would pass unchanged if the
  extraction were reverted. Resolving it means reverting the extraction turns
  `no_raw_pii_in_anonymized_output` red, which is the mutation proof.

  WHY THE SEEDED INTAKE NAMES ONLY THE EXECUTIVE SPONSOR — deliberately, and
  not to flatter the result. "Marcus Chen" and "Priya Iyer" stay OFF the
  deny-list so that Presidio's PERSON NER remains independently measured:
  `distinct_values_distinct_placeholders` asserts >= 2 distinct PERSON values,
  and `no_raw_pii_in_anonymized_output` reports which detector caught each
  shape. Seed every name and both of those go vacuous — the rubric would stop
  measuring NER at all and could not tell a working model from a dead one. It
  is also the realistic case: intake is filled in before discovery, so it
  carries the sponsor, and the rest of the attendees surface in the transcript.

  The doc- and image-format checks keep the NARROW, client-only deny list on
  purpose — it is `_new_session`'s default, and that default is load-bearing:
  their whole point is that `pii.ingest`'s record-per-row table rendering
  makes the name detectable by PERSON. Hand them the stakeholder term and
  they would pass even if that rendering regressed to pipe tables. Only the
  two checks whose subject IS deny-list coverage
  (`round_trip_byte_identical`, `no_raw_pii_in_anonymized_output`) pass
  `deny_terms=resolved_deny_terms()`.
"""
from __future__ import annotations

import hashlib
import io
import json
import re
import stat
import sys
import tempfile
from pathlib import Path
from typing import Optional

from rubrics.base import CheckResult, repo_root
from rubrics._harness import (
    HookRun,
    bool_check,
    build_fixture_engagement,
    check_runs_under_registered_interpreter,
    fault_injection_skip,
    inject_fault,
    pretooluse_payload,
    run_hook_subprocess,
)

# Make `scripts/` importable as a package root, exactly as
# scripts/anonymize_transcript.py does for itself.
_SCRIPTS_DIR = repo_root() / "scripts"
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

FIXTURE_REL_PATH = Path("evals") / "goldens" / "pii_roundtrip_fixture.md"

# The fixture's client identity (see the fixture's own header comment) —
# an obviously-synthetic institution name, not a real bank.
CLIENT_FULL = "Zzzplaceholder Meridian Holdings"
CLIENT_SHORT = "Meridian"

# The client-identity-ONLY deny list. Still the default for every check whose
# subject is something other than deny-list coverage — in particular the
# document- and image-format checks, which have to keep proving that
# `pii.ingest`'s record-per-row rendering makes a tabular name PERSON-
# detectable, and would pass vacuously if handed the stakeholder term.
CLIENT_DENY_TERMS = [CLIENT_FULL, CLIENT_SHORT]
DENY_TERMS = CLIENT_DENY_TERMS  # back-compat alias for existing call sites

# The fixture engagement's stakeholder, seeded into the deny-list SOURCE
# documents below (never into a hard-coded term list) — see the module
# docstring for why only one of the fixture's three people is seeded.
SEEDED_STAKEHOLDER = "Aisha Rahman"
SEEDED_CLIENT_DIR = "zzzplaceholderclient"

_RESOLVED_DENY_TERMS: Optional[list] = None


def _seed_denylist_source_engagement(root: Path):
    """An ordinary engagement tree — the two documents `denylist.py` reads,
    in the shapes `templates/client_profile.md` and
    `templates/inputs/engagement_intake.md` actually produce. Nothing here
    is special-cased for the engine: it is what a consultant's filled-in
    intake looks like.

    CLIENT_PROFILE.md is GENERATED, by `_harness.default_client_profile_text`
    via the `stakeholder=` argument, and that is the point: this fixture's
    claim is that the document is ORDINARY, so deriving it from the shared
    template shape is more honest than a hand-written copy that can drift away
    from `templates/client_profile.md` without anyone noticing. The intake's
    stakeholder TABLE is passed verbatim instead — a table row is the document
    shape #209 turns on, so it is the subject of the check, not scaffolding.
    """
    return build_fixture_engagement(
        root,
        slug=SEEDED_CLIENT_DIR,
        client_name=CLIENT_FULL,
        short_name=CLIENT_SHORT,
        stakeholder=SEEDED_STAKEHOLDER,
        engagement="2026-08_onboarding_assessment",
        documents={
            "inputs/engagement_intake.md": (
                "# Engagement Intake\n\n"
                f"- **Client Name:** {CLIENT_FULL}\n\n"
                "## Stakeholders Interviewed\n\n"
                "| Name/Role | Department | Perspective |\n"
                "|-----------|------------|-------------|\n"
                f"| {SEEDED_STAKEHOLDER} — Chief Financial Officer | Finance | Strategic |\n"
            ),
        },
    )


def resolved_deny_terms() -> list:
    """The deny-list PRODUCTION would resolve for the fixture engagement.

    Deliberately not a literal list. `denylist.resolve_engagement_deny_list`
    is the function `orchestrate.py` and `anonymize_transcript.py` call, so
    running it here is what makes "the stakeholder name is redacted" a
    statement about the shipped extraction rather than about this file.

    `FixtureEngagement.resolved_deny_terms()` runs
    `_harness.assert_deny_list_non_empty` before it returns, defaulting the
    required members to the fixture's own client name and stakeholder — so the
    vacuity guard is now STRUCTURAL (the only supported way to get a resolved
    deny-list out of a fixture goes through it) rather than a call this
    function has to remember to make. A silently-empty deny-list is the exact
    way this repo has twice shipped a gate scoring 1.000 while certifying
    nothing.

    SCOPE: this is the RESOLVED, stakeholder-inclusive list. Only the two
    checks whose subject IS deny-list coverage take it — see `_new_session`
    and the module docstring's closing paragraph for why the document- and
    image-format checks must keep the narrow, client-only default (#209).
    """
    global _RESOLVED_DENY_TERMS
    if _RESOLVED_DENY_TERMS is not None:
        return _RESOLVED_DENY_TERMS

    with tempfile.TemporaryDirectory(prefix="pii_eval_denylist_") as td:
        _RESOLVED_DENY_TERMS = _seed_denylist_source_engagement(
            Path(td)).resolved_deny_terms()
    return _RESOLVED_DENY_TERMS


# Raw values planted in the fixture that MUST NEVER appear in anonymised
# output. "Aisha Rahman" is in this list as of #209: the table-shape PERSON
# miss is closed by the stakeholder deny-list path, so excluding it would now
# be excluding a value that does not leak. The CLIENT_ENTITY's own case
# variants stay out — they are asserted separately by
# _client_name_redacted_via_denylist.
MUST_NOT_LEAK = [
    "priya.iyer@zzzplaceholdermeridian.com",
    "marcus.chen@zzzplaceholdermeridian.com",
    "ops.desk@zzzplaceholdermeridian.com",
    "(555) 201-4477",
    "555.201.9981",
    "923-45-6781",     # SSN
    "8834021177",      # account number 1
    "5521873390",      # account number 2
    "Marcus Chen",
    "Priya Iyer",
    "Aisha Rahman",    # table shape — closed via the stakeholder deny-list
]

# Person names planted in the fixture, one per document shape, used by
# _no_raw_pii_in_anonymized_output's per-shape reporting. Every shape must
# now be DETECTED; the report also names WHICH detector caught each one, so
# "the NER improved" can never be silently confused with "the deny-list
# covered for it".
PERSON_BY_SHAPE = {
    "prose": "Priya Iyer",
    "attendee_bullet": "Marcus Chen",
    "speaker_label": "Priya Iyer",
    "table": "Aisha Rahman",
}

def _fixture_path(target: str) -> Path:
    if target:
        p = Path(target)
        if p.is_file():
            return p
    return repo_root() / FIXTURE_REL_PATH


def _fixture_text(target: str) -> str:
    return _fixture_path(target).read_text(encoding="utf-8")


def _engine():
    from pii import engine as _e  # noqa: PLC0415 - only this rubric needs Presidio
    return _e


def _facade():
    import anonymize_transcript as _at  # noqa: PLC0415
    return _at


def _boundary():
    import artifact_boundary as _ab  # noqa: PLC0415 - ticket #165 checks only
    return _ab


def _denylist():
    from pii import denylist as _d  # noqa: PLC0415
    return _d


def _new_session(engine, *, entity_mapping: Optional[dict] = None, deny_terms=None):
    """Default deny list is CLIENT-ONLY, and that default is load-bearing.

    The document- and image-format checks rely on it: they prove that
    `pii.ingest`'s record-per-row table rendering makes a tabular person name
    PERSON-detectable, and a stakeholder deny term would redact that name
    regardless of the rendering, so those checks would pass a regression.
    Only the two checks whose subject IS deny-list coverage pass
    `deny_terms=resolved_deny_terms()`. Do not "unify" this default.
    """
    return engine.PIISession(
        deny_terms if deny_terms is not None else CLIENT_DENY_TERMS,
        entity_mapping=entity_mapping,
        warn_on_empty=False,
    )


# --- checks ------------------------------------------------------------

def _round_trip_byte_identical(target: str) -> CheckResult:
    """anonymize(fixture) -> deanonymize(...) must restore the ORIGINAL bytes
    exactly, through BOTH maintained copies of deanonymize_text: the engine's
    own (`pii.engine.deanonymize_text`) and the pure-stdlib facade copy
    (`scripts/anonymize_transcript.deanonymize_text`) that
    `artifact_boundary.deanonymize_dir` actually calls in production. The
    facade module docstring itself warns these are two independently
    maintained copies with "no drift_check for this pair yet" — so this
    check also catches the two silently diverging, not just one of them
    breaking round-trip.

    Runs on the PRODUCTION-resolved deny list (`resolved_deny_terms()`), not
    the client-only one, so the stakeholder terms #209 added are covered
    here too: a new deny term that redacted correctly but did not restore
    byte-identically would be a data-loss bug, and this is the check that
    would catch it.
    """
    name = "round_trip_byte_identical"
    engine = _engine()
    at = _facade()
    original = _fixture_text(target)

    session = _new_session(engine, deny_terms=resolved_deny_terms())
    anonymized = session.anonymize(original)
    mapping = session.mapping_file_dict()

    restored_engine = engine.deanonymize_text(anonymized, mapping)
    restored_facade = at.deanonymize_text(anonymized, mapping)

    ok = (restored_engine == original) and (restored_facade == original)
    return bool_check(name, ok, detail=(
        f"engine round-trip identical={restored_engine == original}; "
        f"facade round-trip identical={restored_facade == original}; "
        f"len(original)={len(original)} len(anonymized)={len(anonymized)}"
    ))


def _distinct_values_distinct_placeholders(target: str) -> CheckResult:
    """Every DISTINCT raw value that got detected must map to its OWN
    placeholder — no two different original values collapsing onto the same
    `<ENTITY_N>` token. Checked structurally (the mapping is injective, per
    entity type) AND against the known, hand-verified raw values planted in
    the fixture (3 distinct emails, 2 distinct phones, 2 distinct account-ish
    numbers, 1 SSN, >=2 distinct CLIENT strings, 2 distinct PERSON names)."""
    name = "distinct_values_distinct_placeholders"
    engine = _engine()
    session = _new_session(engine)
    session.anonymize(_fixture_text(target))
    mapping = session.entity_mapping

    problems = []
    for etype, values in mapping.items():
        placeholders = list(values.values())
        if len(placeholders) != len(set(placeholders)):
            problems.append(f"{etype}: {len(placeholders)} values but only "
                             f"{len(set(placeholders))} distinct placeholders")

    expected_min = {
        "EMAIL_ADDRESS": 3, "PHONE_NUMBER": 2, "US_SSN": 1,
        "CLIENT": 2, "PERSON": 2,
    }
    for etype, minimum in expected_min.items():
        got = len(mapping.get(etype, {}))
        if got < minimum:
            problems.append(f"{etype}: expected >= {minimum} distinct values, got {got}")

    ok = not problems
    return bool_check(name, ok, detail=(
        "; ".join(problems) if problems else
        f"all {sum(len(v) for v in mapping.values())} distinct values across "
        f"{len(mapping)} entity types map 1:1 to distinct placeholders"
    ))


def _repeated_value_reuses_placeholder(target: str) -> CheckResult:
    """The fixture repeats one email address (priya.iyer@...) three times.
    All three occurrences in the anonymized text must be the SAME
    placeholder, and the mapping must hold exactly one entry for that value
    (not three)."""
    name = "repeated_value_reuses_placeholder"
    engine = _engine()
    session = _new_session(engine)
    original = _fixture_text(target)
    repeated_value = "priya.iyer@zzzplaceholdermeridian.com"
    assert original.count(repeated_value) >= 3, "fixture regression: repeated email no longer repeats >=3x"

    anonymized = session.anonymize(original)
    mapping = session.entity_mapping
    placeholder = mapping.get("EMAIL_ADDRESS", {}).get(repeated_value)

    ok = (
        placeholder is not None
        and anonymized.count(placeholder) >= 3
        and repeated_value not in anonymized
    )
    return bool_check(name, ok, detail=(
        f"placeholder={placeholder!r} occurrences_in_output="
        f"{anonymized.count(placeholder) if placeholder else 'n/a'} "
        f"(expected >=3, one mapping entry, raw value gone)"
    ))


def _no_raw_pii_in_anonymized_output(target: str) -> CheckResult:
    """The gate-bites-mandatory check (D9 / PRD §5): none of the fixture's
    deterministic PII values (emails, phones, SSN, account-ish numbers) or
    reliably-detected PERSON names may survive into the anonymized text.
    Reverting the Presidio detector (disabling the entities passed to
    analyze(), or the deny-list recognizer) must make this fail — see the
    PR description's gate-bites transcript.

    Gating scope now includes the table-shape name ("Aisha Rahman"), which
    #209 closed via the stakeholder deny-list path. It was excluded before,
    and that exclusion — not the threshold — was what kept this component at
    1.000 while its own detail line printed `table=MISSED` every run.

    The session runs on `resolved_deny_terms()`, resolved through
    `denylist.resolve_engagement_deny_list` from seeded CLIENT_PROFILE.md /
    engagement_intake.md documents. That is what makes this a mutation-
    provable gate: revert `extract_stakeholder_terms` and the resolved list
    loses the name, the name survives into the output, and this check goes
    red. A hard-coded deny term would have survived that revert untouched.

    `detail`/`evidence` still carry the full per-shape measurement, and now
    also name the DETECTOR for each shape (PERSON = Presidio NER, CLIENT =
    the deny-list). Keeping the two distinguishable matters: without it, a
    future NER regression would be invisible behind deny-list coverage.
    """
    name = "no_raw_pii_in_anonymized_output"
    engine = _engine()
    session = _new_session(engine, deny_terms=resolved_deny_terms())
    anonymized = session.anonymize(_fixture_text(target))

    leaked = [v for v in MUST_NOT_LEAK if v in anonymized]

    # Per-shape measurement. A name is "detected for its shape" if it no
    # longer appears raw anywhere post-anonymization; each PERSON_BY_SHAPE
    # name is unique to its shape's sentence in the fixture, so this
    # attribution is sound even though the flat mapping itself doesn't carry
    # shape provenance. The entity type comes from the mapping, so the report
    # says which detector actually caught it.
    def _detector(value: str) -> str:
        for etype, values in session.entity_mapping.items():
            if value in values:
                return etype
        return "none"

    per_shape = {
        shape: (person not in anonymized, _detector(person))
        for shape, person in PERSON_BY_SHAPE.items()
    }
    missed = sorted({PERSON_BY_SHAPE[s] for s, (hit, _) in per_shape.items() if not hit})

    ok = not leaked
    detail = (
        f"leaked={leaked!r}; per-shape person detection: " +
        ", ".join(
            f"{s}={'DETECTED via ' + etype if hit else 'MISSED'} "
            f"({PERSON_BY_SHAPE[s]!r})"
            for s, (hit, etype) in per_shape.items()
        ) +
        " (CLIENT = engagement deny-list, PERSON = Presidio NER; every shape "
        "is gated via MUST_NOT_LEAK)"
    )
    return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True, detail=detail,
                        evidence=[f"leaked raw value: {v}" for v in leaked]
                                 + [f"undetected in shape: {v}" for v in missed])


def _cross_transcript_merge_collision_free(target: str) -> CheckResult:  # noqa: ARG001
    """#160 regression guard — the ONE thing keeping cross-transcript
    numbering collision-free. Exercises TWO transcripts sharing one email
    address, run through `anonymize_transcript_file` with a SHARED
    `entity_mapping` dict — exactly what `orchestrate.py`'s `step_discovery`
    does. Before the #160 fix, each transcript numbered independently, so
    transcript A's `<EMAIL_ADDRESS_1>` and transcript B's
    `<EMAIL_ADDRESS_1>` were DIFFERENT values — a later merge bound one
    placeholder to two originals. Asserts all four properties the ticket
    calls out: distinct values -> distinct placeholders across both files;
    the SHARED value -> the SAME placeholder in both; no placeholder maps to
    two different values; both files restore byte-identical from the final
    combined mapping alone.
    """
    name = "cross_transcript_merge_collision_free"
    at = _facade()
    with tempfile.TemporaryDirectory(prefix="pii_eval_cross_") as td:
        root = Path(td)
        engagement_dir = root / "zzzplaceholder_engagement"
        inputs_dir = engagement_dir / "inputs"
        inputs_dir.mkdir(parents=True)

        shared_email = "shared.contact@zzzplaceholdermeridian.com"
        transcript_a = inputs_dir / "transcript_a.md"
        transcript_b = inputs_dir / "transcript_b.md"
        transcript_a.write_text(
            f"Contact for this workstream: {shared_email}. "
            f"Local lead: Naledi Dube can also help.\n"
        )
        transcript_b.write_text(
            f"Escalation goes to {shared_email}. "
            f"Secondary contact: someone.else@zzzplaceholdermeridian.com.\n"
        )

        shared_mapping: dict = {}
        anon_a, _map_a = at.anonymize_transcript_file(
            transcript_a, engagement_dir, output_dir=inputs_dir, entity_mapping=shared_mapping)
        anon_b, map_b = at.anonymize_transcript_file(
            transcript_b, engagement_dir, output_dir=inputs_dir, entity_mapping=shared_mapping)

        text_a = anon_a.read_text()
        text_b = anon_b.read_text()
        final_mapping = json.loads(map_b.read_text())

        # 1. distinct values -> distinct placeholders, across BOTH files combined.
        email_map = shared_mapping.get("EMAIL_ADDRESS", {})
        distinct_ok = len(set(email_map.values())) == len(email_map)

        # 2. the shared value gets the SAME placeholder in both files' text.
        shared_placeholder = email_map.get(shared_email)
        same_placeholder_ok = (
            shared_placeholder is not None
            and shared_placeholder in text_a
            and shared_placeholder in text_b
        )

        # 3. no placeholder maps to two different originals (across ALL types).
        no_collision_ok = True
        for etype, values in shared_mapping.items():
            inverse: dict = {}
            for original, placeholder in values.items():
                if inverse.setdefault(placeholder, original) != original:
                    no_collision_ok = False

        # 4. both files restore byte-identical from the FINAL combined mapping
        #    alone (no per-transcript mapping needed).
        restored_a = at.deanonymize_text(text_a, final_mapping)
        restored_b = at.deanonymize_text(text_b, final_mapping)
        restore_ok = (
            restored_a == transcript_a.read_text()
            and restored_b == transcript_b.read_text()
        )

        ok = distinct_ok and same_placeholder_ok and no_collision_ok and restore_ok
        return bool_check(name, ok, detail=(
            f"distinct_ok={distinct_ok} same_placeholder_ok={same_placeholder_ok} "
            f"(placeholder={shared_placeholder!r}) no_collision_ok={no_collision_ok} "
            f"restore_ok={restore_ok}"
        ))


def _mapping_files_chmod_600_and_cleaned(target: str) -> CheckResult:  # noqa: ARG001
    """`anonymize_transcript_file` writes a per-transcript
    `.anon_mapping_<stem>.json` file, chmod 0600, because it carries real
    PII (PRD v6 §9). This check asserts that mode AND the precondition that
    makes it SAFE for a caller to delete every per-transcript mapping file
    once a later one has been written with the shared `entity_mapping`:
    each earlier per-transcript mapping must be a strict SUBSET of the
    final, most-recently-written one, and the final one alone must be
    sufficient to restore every transcript's anonymized text byte-for-byte
    — i.e. discarding the earlier files loses no information.

    SCOPE NOTE: this checks the FACADE's contract (`anonymize_transcript.
    anonymize_transcript_file` + the shared `entity_mapping` param), which is
    what exists today. `orchestrate.py`'s `step_discovery` does NOT
    currently delete the per-transcript files it writes (verified by
    inspection: the mapping_path returned by each call is discarded, never
    unlinked) despite PRD v6 §9 asserting "per-transcript mappings are
    deleted once the combined mapping is written" — that deletion call is a
    real, separate gap, not exercised here because it lives in orchestrate.py
    and PR2 (this ticket's scope, per solution-design-v6.md's Build Sequence)
    is the engine + facade, not orchestrate.py's pipeline wiring. Flagged as
    a finding for a follow-up ticket rather than silently asserted or
    silently dropped.
    """
    name = "mapping_files_chmod_600_and_cleaned"
    at = _facade()
    with tempfile.TemporaryDirectory(prefix="pii_eval_mapclean_") as td:
        root = Path(td)
        engagement_dir = root / "zzzplaceholder_engagement"
        inputs_dir = engagement_dir / "inputs"
        inputs_dir.mkdir(parents=True)

        transcript_a = inputs_dir / "transcript_a.md"
        transcript_b = inputs_dir / "transcript_b.md"
        transcript_a.write_text("Lead: Naledi Dube. Contact: a.person@zzzplaceholdermeridian.com.\n")
        transcript_b.write_text("Follow-up: a.person@zzzplaceholdermeridian.com and b.person@zzzplaceholdermeridian.com.\n")

        shared_mapping: dict = {}
        anon_a, map_a = at.anonymize_transcript_file(
            transcript_a, engagement_dir, output_dir=inputs_dir, entity_mapping=shared_mapping)
        anon_b, map_b = at.anonymize_transcript_file(
            transcript_b, engagement_dir, output_dir=inputs_dir, entity_mapping=shared_mapping)

        mode_a = stat.S_IMODE(map_a.stat().st_mode)
        mode_b = stat.S_IMODE(map_b.stat().st_mode)
        mode_ok = (mode_a == 0o600) and (mode_b == 0o600)

        snapshot_a = json.loads(map_a.read_text()).get("entities", {})
        snapshot_b = json.loads(map_b.read_text()).get("entities", {})

        # snapshot_a must be a strict subset of snapshot_b (every value ->
        # placeholder pair present in A is present, IDENTICALLY, in B) — the
        # property that makes deleting A safe once B is on disk.
        subset_ok = True
        for etype, values in snapshot_a.items():
            for original, placeholder in values.items():
                if snapshot_b.get(etype, {}).get(original) != placeholder:
                    subset_ok = False

        # Discarding map_a entirely and restoring BOTH transcripts from
        # map_b alone must still work byte-identical.
        restored_a = at.deanonymize_file(anon_a, map_b)
        restored_b = at.deanonymize_file(anon_b, map_b)
        cleanup_safe_ok = (
            restored_a == transcript_a.read_text()
            and restored_b == transcript_b.read_text()
        )

        ok = mode_ok and subset_ok and cleanup_safe_ok
        return bool_check(name, ok, detail=(
            f"mode_a={oct(mode_a)} mode_b={oct(mode_b)} mode_ok={mode_ok} "
            f"subset_ok={subset_ok} cleanup_safe_ok={cleanup_safe_ok} — NOTE: "
            f"orchestrate.py does not currently delete per-transcript mapping "
            f"files; this check proves it WOULD be safe to, see docstring"
        ))


def _client_name_redacted_via_denylist(target: str) -> CheckResult:
    """D3: the deny-list, not Presidio NER, is the PRIMARY client-identity
    detector. Asserts the client's full legal name, its short form, AND the
    short form embedded in the client's own portal domain
    (https://portal.meridian.com/...) are all redacted via the CLIENT
    entity — none of Presidio's built-in recognizers know what "the client"
    is; only the deny-list does."""
    name = "client_name_redacted_via_denylist"
    engine = _engine()
    session = _new_session(engine)
    original = _fixture_text(target)
    assert CLIENT_FULL in original and "portal.meridian.com" in original, \
        "fixture regression: client full name / portal domain no longer present"

    anonymized = session.anonymize(original)
    mapping = session.entity_mapping.get("CLIENT", {})

    full_name_redacted = CLIENT_FULL not in anonymized and CLIENT_FULL in mapping
    domain_redacted = "portal.meridian.com" not in anonymized
    has_client_entities = len(mapping) >= 2  # full name + at least one short-form case variant

    ok = full_name_redacted and domain_redacted and has_client_entities
    return bool_check(name, ok, detail=(
        f"full_name_redacted={full_name_redacted} domain_redacted={domain_redacted} "
        f"distinct_CLIENT_values={len(mapping)} (expected >=2)"
    ))


def _allowlist_prevents_generic_overredaction(target: str) -> CheckResult:
    """The engine's allow-list (engine.py's `_build_allow_list`, built from
    `denylist.GENERIC_STOPLIST`) exists to stop generic banking words
    ("First", "National", "Trust", "Capital", "Pacific", "Union", "State", …)
    being redacted when they appear on their own, not as part of the
    client's registered name.

    Two assertions, because on the CURRENT entity configuration (no
    ORGANIZATION/LOCATION recognizer enabled — see engine.py's
    DEFAULT_ENTITIES comment) generic prose words don't trigger any OTHER
    recognizer either, so an end-to-end "redact without the allow-list"
    demonstration is not reproducible against today's model/config — that is
    reported here, not hidden:

      1. STRUCTURAL (the gate-bites-verified assertion): the allow-list is
         actually built and wired into the session — every
         GENERIC_STOPLIST word not already a deny term appears as an anchored
         allow-list pattern.
      2. BEHAVIOURAL (a regression guard, not gate-bitable via the allow-list
         alone today, but real): a paragraph of generic banking peer-comparison
         language stays fully unredacted end-to-end.
    """
    name = "allowlist_prevents_generic_overredaction"
    engine = _engine()
    from pii import denylist  # noqa: PLC0415

    session = _new_session(engine)
    allow_list = session._allow_list  # noqa: SLF001 - inspecting the real wired state, not a copy

    stoplist_words = sorted(denylist.GENERIC_STOPLIST)
    expected_patterns = {r"^" + w + r"$" for w in stoplist_words}
    structural_ok = expected_patterns.issubset(set(allow_list)) and len(allow_list) > 0

    generic_text = ("The team benchmarked results against First National, Pacific Trust, "
                     "Capital Union, and State Savings as generic regional peers.")
    anonymized = session.anonymize(generic_text)
    behavioural_ok = anonymized == generic_text

    ok = structural_ok and behavioural_ok
    return bool_check(name, ok, detail=(
        f"structural_ok={structural_ok} (allow_list has {len(allow_list)} patterns) "
        f"behavioural_ok={behavioural_ok}"
    ))


def _empty_entity_list_warns(target: str) -> CheckResult:  # noqa: ARG001
    """Loud, non-blocking warning when the deny-list resolves empty
    (EMPTY_DENY_LIST_WARNING) — the exact failure mode PRD v6 §1 opens with
    ("a live audit had person names AND the client name reach the API in
    plaintext with no warning at all"). Must warn AND must not block: the
    session still works and generic entities (email/phone/SSN) still get
    redacted."""
    name = "empty_entity_list_warns"
    engine = _engine()
    import io

    warn_stream = io.StringIO()
    session = engine.PIISession([], warn_stream=warn_stream, warn_on_empty=True)
    warned = engine.EMPTY_DENY_LIST_WARNING.strip() in warn_stream.getvalue()

    quiet_stream = io.StringIO()
    _configured_session = engine.PIISession(DENY_TERMS, warn_stream=quiet_stream, warn_on_empty=True)
    stayed_quiet = quiet_stream.getvalue() == ""

    text = "Reach out at generic.person@zzzplaceholdermeridian.com for details."
    anonymized = session.anonymize(text)
    still_functional = "generic.person@zzzplaceholdermeridian.com" not in anonymized

    ok = warned and stayed_quiet and still_functional
    return bool_check(name, ok, detail=(
        f"warned_when_empty={warned} silent_when_configured={stayed_quiet} "
        f"still_redacts_generic_pii={still_functional}"
    ))


def _legacy_flat_mapping_still_restores(target: str) -> CheckResult:  # noqa: ARG001
    """Every accepted mapping shape (flatten_mapping's docstring: v2 nested,
    v2 bare, v1 legacy flat) must still restore — "mappings on disk are
    data, and a consultant with a six-month-old engagement must still be
    able to produce a client-ready deliverable." Exercises BOTH maintained
    copies (engine.py's and the facade's pure-stdlib copy) so a future
    change to one that isn't mirrored to the other is caught here rather
    than only in production."""
    name = "legacy_flat_mapping_still_restores"
    engine = _engine()
    at = _facade()

    v1_legacy = {
        "[CLIENT]": "Zzzplaceholder Legacy Holdings",
        "[EMAIL-REDACTED]": "legacy.contact@zzzplaceholdermeridian.com",
        "[PERSON-1]": "Legacy Testperson",
    }
    v1_text = "Client: [CLIENT]. Contact [PERSON-1] at [EMAIL-REDACTED]."
    v1_expected = "Client: Zzzplaceholder Legacy Holdings. Contact Legacy Testperson at legacy.contact@zzzplaceholdermeridian.com."

    v2_bare = {"PERSON": {"Bare Testperson": "<PERSON_1>"}}
    v2_bare_text = "Attendee: <PERSON_1>."
    v2_bare_expected = "Attendee: Bare Testperson."

    v2_nested = {"version": 2, "entities": {"CLIENT": {"Zzzplaceholder Nested Co": "<CLIENT_1>"}}}
    v2_nested_text = "Regarding <CLIENT_1>'s engagement."
    v2_nested_expected = "Regarding Zzzplaceholder Nested Co's engagement."

    cases = [
        ("v1_legacy", v1_text, v1_legacy, v1_expected),
        ("v2_bare", v2_bare_text, v2_bare, v2_bare_expected),
        ("v2_nested", v2_nested_text, v2_nested, v2_nested_expected),
    ]

    problems = []
    for label, text, mapping, expected in cases:
        got_engine = engine.deanonymize_text(text, mapping)
        got_facade = at.deanonymize_text(text, mapping)
        if got_engine != expected:
            problems.append(f"{label}: engine.deanonymize_text -> {got_engine!r} (expected {expected!r})")
        if got_facade != expected:
            problems.append(f"{label}: facade.deanonymize_text -> {got_facade!r} (expected {expected!r})")
        if got_engine != got_facade:
            problems.append(f"{label}: engine/facade DIVERGED — {got_engine!r} != {got_facade!r}")

    ok = not problems
    return bool_check(name, ok, detail="; ".join(problems) if problems else
                        "all 3 mapping shapes (v1 legacy, v2 bare, v2 nested) restore "
                        "correctly through both maintained copies")


# --- #165 artifact_boundary.deanonymize_dir: recursion, xlsx, hazards ------

def _nested_outputs_deanonymized(target: str) -> CheckResult:  # noqa: ARG001
    """Ticket #165 — `deanonymize_dir` must walk `outputs/` RECURSIVELY
    (`rglob('*')`, not top-level `iterdir()`), while still excluding
    dotfiles/dot-directories (`.anon_*`, `.pii_mapping.json`) and
    `interim*`-named files. GATE-BITES: reverting the recursive walk back to
    `iterdir()` makes this FAIL, because the nested file is never visited.

    Also folds in the two hazards #165 calls out by name, both exercised
    through the REAL `deanonymize_dir` (never a hand-rolled replace):

    Hazard 1 — legacy `[CLIENT]` vs. the `[CLIENT]_Business_Case_
    Questionnaire.xlsx` filename/prose template token used by five other
    components (roi-financial-modeler, generate-roi-questionnaire,
    generate-roi-excel, usecase-doc, prototype). A blind substring replace
    of an old engagement's `[CLIENT]` mapping would corrupt that token; this
    asserts the legacy placeholder DOES restore elsewhere in the same file
    while the template token is left byte-identical.

    Hazard 2 — longest-placeholder-first ordering (#159 verified this for
    the engine's own anonymize path; this confirms the restore path):
    `[CLIENT-ABBR]` must resolve to its own value, not get corrupted by
    `[CLIENT]` being processed first/overlapping.
    """
    name = "nested_outputs_deanonymized"
    ab = _boundary()

    with tempfile.TemporaryDirectory(prefix="pii_eval_nested_") as td:
        root = Path(td)
        engagement_dir = root / "zzzplaceholder_engagement"
        outputs_dir = engagement_dir / "outputs"
        nested_dir = outputs_dir / "subdir" / "deeper"
        nested_dir.mkdir(parents=True)

        mapping = {
            "[CLIENT]": "Zzzplaceholder Nested Holdings",
            "[CLIENT-ABBR]": "ZNH",
        }
        mapping_file = engagement_dir / ".pii_mapping.json"
        mapping_file.write_text(json.dumps(mapping))

        # Top-level file.
        top_file = outputs_dir / "summary.md"
        top_file.write_text("Prepared for [CLIENT].\n")

        # Nested file (2 levels deep) — the recursion witness.
        nested_file = nested_dir / "detail.html"
        nested_file.write_text("<p>Approved by [CLIENT-ABBR] on behalf of [CLIENT].</p>")

        # Hazard 1: a file with BOTH a genuine legacy placeholder occurrence
        # AND the unrelated filename/prose template token that starts with
        # the identical bracket text.
        hazard_file = outputs_dir / "cover_note.md"
        hazard_file.write_text(
            "This is [CLIENT]'s engagement.\n"
            "See attached: [CLIENT]_Business_Case_Questionnaire.xlsx\n"
        )

        # Must be excluded: interim* file and dotfile/dot-dir inside outputs/.
        interim_file = outputs_dir / "interim_draft.md"
        interim_file.write_text("draft notes: [CLIENT]")
        dotfile = outputs_dir / ".pii_mapping.json"
        dotfile.write_text(json.dumps(mapping))
        dotdir_file = outputs_dir / ".anon_cache" / "raw.md"
        dotdir_file.parent.mkdir(parents=True)
        dotdir_file.write_text("[CLIENT]")

        interim_before = interim_file.read_text()
        dotfile_before = dotfile.read_bytes()
        dotdir_before = dotdir_file.read_bytes()

        report = ab.deanonymize_dir(outputs_dir, mapping_file)

        problems = []
        if not report.get("client_ready"):
            problems.append(f"client_ready False: {report}")

        top_restored = top_file.read_text()
        if top_restored != "Prepared for Zzzplaceholder Nested Holdings.\n":
            problems.append(f"top-level file not restored correctly: {top_restored!r}")

        nested_restored = nested_file.read_text()
        expected_nested = "<p>Approved by ZNH on behalf of Zzzplaceholder Nested Holdings.</p>"
        if nested_restored != expected_nested:
            problems.append(f"NESTED FILE NOT RESTORED (recursion witness): {nested_restored!r} "
                             f"(expected {expected_nested!r})")

        hazard_restored = hazard_file.read_text()
        expected_hazard = (
            "This is Zzzplaceholder Nested Holdings's engagement.\n"
            "See attached: [CLIENT]_Business_Case_Questionnaire.xlsx\n"
        )
        if hazard_restored != expected_hazard:
            problems.append(f"HAZARD 1 (template-token collision) FAILED: {hazard_restored!r} "
                             f"(expected {expected_hazard!r})")

        if interim_file.read_text() != interim_before:
            problems.append("interim_* file was modified — must stay excluded")
        if dotfile.read_bytes() != dotfile_before:
            problems.append(".pii_mapping.json inside outputs/ was modified — must never be touched")
        if dotdir_file.read_bytes() != dotdir_before:
            problems.append("file inside a dot-directory was modified — must stay excluded")

        ok = not problems
        return bool_check(name, ok, detail="; ".join(problems) if problems else (
            f"recursion witness restored, hazard 1 (template token) preserved, "
            f"hazard 2 (longest-match ordering) correct, interim/dotfile/dot-dir "
            f"exclusions held; files_restored={report.get('files_restored')}"
        ))


def _xlsx_outputs_deanonymized(target: str) -> CheckResult:  # noqa: ARG001
    """Ticket #165 — `.xlsx` outputs must be restored: placeholders in
    string cell values (including formula strings) and sheet titles, saved
    IN PLACE only when something actually changed. GATE-BITES: reverting the
    `.xlsx` branch (e.g. skipping xlsx files entirely) makes this FAIL,
    because the workbook would come back with placeholders still in it.

    Also asserts the unavailable-engine / unopenable-workbook failure mode:
    `openpyxl` import failure must land the file in `unrestored` with the
    re-run command named, and `client_ready` must be False — never a silent
    skip (this module's own header: "must stay importable... Silently
    skipping xlsx is exactly the audit bug").
    """
    name = "xlsx_outputs_deanonymized"
    ab = _boundary()
    import openpyxl  # noqa: PLC0415 - this component's env always has it (requirements.txt)

    with tempfile.TemporaryDirectory(prefix="pii_eval_xlsx_") as td:
        root = Path(td)
        engagement_dir = root / "zzzplaceholder_engagement"
        outputs_dir = engagement_dir / "outputs"
        outputs_dir.mkdir(parents=True)

        mapping = {"version": 2, "entities": {
            "CLIENT": {"Zzzplaceholder Sheet Holdings": "<CLIENT_1>"},
        }}
        mapping_file = engagement_dir / ".pii_mapping.json"
        mapping_file.write_text(json.dumps(mapping))

        xlsx_path = outputs_dir / "ROI_Model.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "<CLIENT_1> Summary"
        ws["A1"] = "Prepared for <CLIENT_1>"
        ws["B2"] = '=CONCATENATE("Client: ", "<CLIENT_1>")'
        wb.save(xlsx_path)
        before_bytes = xlsx_path.read_bytes()

        report = ab.deanonymize_dir(outputs_dir, mapping_file)

        problems = []
        if not report.get("client_ready"):
            problems.append(f"client_ready False on a clean restore: {report}")
        if report.get("files_restored", 0) < 1:
            problems.append(f"files_restored did not count the xlsx: {report}")

        wb2 = openpyxl.load_workbook(xlsx_path)
        ws2 = wb2.active
        if ws2.title != "Zzzplaceholder Sheet Holdings Summary":
            problems.append(f"sheet title not restored: {ws2.title!r}")
        if ws2["A1"].value != "Prepared for Zzzplaceholder Sheet Holdings":
            problems.append(f"cell value not restored: {ws2['A1'].value!r}")
        expected_formula = '=CONCATENATE("Client: ", "Zzzplaceholder Sheet Holdings")'
        if ws2["B2"].value != expected_formula:
            problems.append(f"formula string not restored: {ws2['B2'].value!r} "
                             f"(expected {expected_formula!r})")

        # Re-run with nothing left to restore: the workbook must NOT be
        # re-saved (openpyxl round-trips are not byte-stable even when no
        # cell changes, so re-saving unconditionally would be detectable
        # here and is exactly what "only when something actually changed"
        # forbids).
        after_bytes = xlsx_path.read_bytes()
        report2 = ab.deanonymize_dir(outputs_dir, mapping_file)
        if report2.get("files_restored", 0) != 0:
            problems.append(f"second run on an already-restored workbook reported "
                             f"files_restored={report2.get('files_restored')}, expected 0")
        if xlsx_path.read_bytes() != after_bytes:
            problems.append("workbook was re-saved on a no-op second run "
                             "(checksum changed) — must save only when changed")

        # --- Negative path: openpyxl unavailable / workbook unopenable ----
        broken_dir = root / "broken_engagement" / "outputs"
        broken_dir.mkdir(parents=True)
        broken_mapping_file = root / "broken_engagement" / ".pii_mapping.json"
        broken_mapping_file.write_text(json.dumps(mapping))
        broken_xlsx = broken_dir / "Corrupt_Model.xlsx"
        broken_xlsx.write_bytes(b"not a real zip/xlsx at all")

        broken_report = ab.deanonymize_dir(broken_dir, broken_mapping_file)
        if broken_report.get("client_ready") is not False:
            problems.append(f"unopenable workbook did not fail client_ready: {broken_report}")
        unrestored = broken_report.get("unrestored", [])
        if not any("Corrupt_Model.xlsx" in u for u in unrestored):
            problems.append(f"unopenable workbook not named in 'unrestored': {unrestored!r}")

        ok = not problems
        return bool_check(name, ok, detail="; ".join(problems) if problems else (
            f"sheet title/cell value/formula string restored, unchanged workbook "
            f"not re-saved, unopenable workbook reported in 'unrestored' with "
            f"client_ready=False; files_restored={report.get('files_restored')}"
        ))


# --- #162 document ingest: fixtures, built programmatically ----------------
#
# Never committed as binaries. See the module docstring's FIXTURE note.

# The same planted PII in every format, so the check is a genuine
# cross-format comparison rather than five unrelated documents.
DOC_TITLE = f"{CLIENT_FULL} — Stakeholder Directory"
DOC_INTRO = (
    f"Prepared for {CLIENT_SHORT} as part of the digital onboarding review."
)
DOC_HEADERS = ["Name", "Role", "Email", "Phone", "Account No."]
DOC_ROWS = [
    ["Aisha Rahman", "Chief Financial Officer",
     "a.rahman@zzzplaceholdermeridian.com", "(555) 201-4477", "8834021177"],
    ["Marcus Chen", "Head of Digital Banking",
     "m.chen@zzzplaceholdermeridian.com", "555.201.9981", "5521873390"],
]

# Values that must not survive into any `.anon_` artifact. Job titles are
# deliberately excluded — they are not PII and redacting them would destroy
# the analysis (engine.py's DEFAULT_ENTITIES comment).
DOC_MUST_NOT_LEAK = [CLIENT_FULL] + [
    v for row in DOC_ROWS for v in row if v not in ("Chief Financial Officer",
                                                     "Head of Digital Banking")
]

# "Aisha Rahman" is planted in TABULAR data in four of the five formats on
# purpose: on `en_core_web_lg` this is the exact name #159 recorded as
# tagged ORGANIZATION inside a markdown table cell, i.e. the name that
# leaks if ingest renders tables as pipe tables. A fixture without it would
# certify the gap #162 exists to close as passing.
DOC_TABULAR_PERSON = "Aisha Rahman"

# Per-format structural marker the extract must carry, proving structure
# survived rather than being flattened into a wall of text.
DOC_STRUCTURE_MARKER = {
    "pdf": "## Page 1",
    "docx": "## Table 1",
    "pptx": "## Slide 2",
    "xlsx": "## Sheet: Contacts",
    "csv": "### Row 1",
}


def _build_minimal_pdf(pages, encrypted: bool = False) -> bytes:
    """A valid single-font PDF, hand-assembled — no PDF writer is installed
    and adding one just to build a fixture is not worth a dependency.

    `encrypted=True` attaches a standard-security-handler `/Encrypt`
    dictionary the empty password cannot open, which is how the
    password-protected path is exercised without shipping a real encrypted
    binary.
    """
    def esc(s: str) -> str:
        return s.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")

    n_pages = len(pages)
    kids = " ".join(f"{4 + 2 * i} 0 R" for i in range(n_pages))
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        f"<< /Type /Pages /Kids [{kids}] /Count {n_pages} >>".encode(),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    for lines in pages:
        stream = "BT /F1 11 Tf 13 TL 56 760 Td\n"
        stream += "".join(f"({esc(line)}) Tj T*\n" for line in lines)
        stream += "ET"
        body = stream.encode("latin-1", "replace")
        content_num = len(objs) + 2
        objs.append(
            (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
             f"/Resources << /Font << /F1 3 0 R >> >> "
             f"/Contents {content_num} 0 R >>").encode()
        )
        objs.append(b"<< /Length %d >>\nstream\n" % len(body) + body + b"\nendstream")

    encrypt_ref = ""
    extra_trailer = ""
    if encrypted:
        objs.append(b"<< /Filter /Standard /V 1 /R 2 /O <" + b"11" * 32 +
                     b"> /U <" + b"22" * 32 + b"> /P -1 >>")
        encrypt_ref = f" /Encrypt {len(objs)} 0 R"
        extra_trailer = " /ID [<%s> <%s>]" % ("33" * 16, "33" * 16)

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objs) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (f"trailer\n<< /Size {len(objs) + 1} /Root 1 0 R{encrypt_ref}"
            f"{extra_trailer} >>\nstartxref\n{xref}\n%%EOF\n").encode()
    return bytes(out)


def _build_document_fixtures(directory: Path) -> dict:
    """One document per supported format, same planted PII in each.

    Returns {format: path}. Tables are REAL tables (a Word table, a
    PowerPoint table, a worksheet, CSV rows) — not pre-rendered text — so
    the check exercises ingest's own rendering decision rather than a shape
    the fixture chose for it. The PDF is the exception and carries the
    contact block as label lines: a PDF has no table structure to re-render,
    so its layout is whatever the file says (documented seam in
    `scripts/pii/ingest.py`).
    """
    import csv as _csv  # noqa: PLC0415
    paths = {}

    # --- PDF -------------------------------------------------------------
    pdf_pages = [[DOC_TITLE, "", DOC_INTRO], ["Contacts", ""]]
    for row in DOC_ROWS:
        for header, value in zip(DOC_HEADERS, row):
            pdf_pages[1].append(f"{header}: {value}.")
        pdf_pages[1].append("")
    paths["pdf"] = directory / "stakeholders.pdf"
    paths["pdf"].write_bytes(_build_minimal_pdf(pdf_pages))

    # --- DOCX: headings + a real table ------------------------------------
    import docx  # noqa: PLC0415
    document = docx.Document()
    document.add_heading(DOC_TITLE, level=1)
    document.add_paragraph(DOC_INTRO)
    document.add_heading("Contacts", level=2)
    table = document.add_table(rows=1 + len(DOC_ROWS), cols=len(DOC_HEADERS))
    for j, header in enumerate(DOC_HEADERS):
        table.cell(0, j).text = header
    for i, row in enumerate(DOC_ROWS, start=1):
        for j, value in enumerate(row):
            table.cell(i, j).text = value
    paths["docx"] = directory / "stakeholders.docx"
    document.save(str(paths["docx"]))

    # --- PPTX: two slides, one carrying a real table ----------------------
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = DOC_TITLE
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(8.5), Inches(0.8))
    box.text_frame.text = DOC_INTRO
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide2.shapes.title.text = "Contacts"
    graphic = slide2.shapes.add_table(
        1 + len(DOC_ROWS), len(DOC_HEADERS),
        Inches(0.4), Inches(1.5), Inches(9), Inches(2),
    )
    for j, header in enumerate(DOC_HEADERS):
        graphic.table.cell(0, j).text = header
    for i, row in enumerate(DOC_ROWS, start=1):
        for j, value in enumerate(row):
            graphic.table.cell(i, j).text = value
    paths["pptx"] = directory / "stakeholders.pptx"
    presentation.save(str(paths["pptx"]))

    # --- XLSX: two sheets --------------------------------------------------
    import openpyxl  # noqa: PLC0415
    workbook = openpyxl.Workbook()
    overview = workbook.active
    overview.title = "Overview"
    overview["A1"] = DOC_TITLE
    overview["A2"] = DOC_INTRO
    contacts = workbook.create_sheet("Contacts")
    contacts.append(DOC_HEADERS)
    for row in DOC_ROWS:
        contacts.append(row)
    paths["xlsx"] = directory / "stakeholders.xlsx"
    workbook.save(str(paths["xlsx"]))

    # --- CSV ---------------------------------------------------------------
    paths["csv"] = directory / "stakeholders.csv"
    with paths["csv"].open("w", newline="", encoding="utf-8") as handle:
        writer = _csv.writer(handle)
        writer.writerow(DOC_HEADERS)
        for row in DOC_ROWS:
            writer.writerow(row)

    return paths


def _document_formats_converted_and_scrubbed(target: str) -> CheckResult:  # noqa: ARG001
    """#162: every text-bearing input format is converted to text, anonymised
    through the SAME engine transcripts use, and written as `.anon_<name>.md`.

    The gate the PRD's "5 of 5 formats covered, from 0 today" metric rests
    on. Nine properties, all of which must hold for every one of PDF, DOCX,
    PPTX, XLSX and CSV:

      1. `.anon_<name>.md` is written at the path `anon_path_for` predicts.
      2. NO planted raw value survives into it — client name, both person
         names, both emails, both phones, both account numbers.
      3. The person name planted in TABULAR data is redacted. Asserted
         separately and explicitly, because this is the shape a markdown
         pipe-table rendering silently fails on: "Aisha Rahman" in a table
         cell is tagged ORGANIZATION by `en_core_web_lg` (#159), which is
         not an enabled entity type. A fixture that only planted names in
         prose would pass while the real gap stayed open.
      4. Structure survived — the format's own boundary markers (pages,
         sheets, slides, tables) are present, so the extract is a usable
         document rather than a wall of text.
      5. Tabular data is NOT rendered as a markdown pipe table (the one
         shape measured to lose names — see ingest.py's measurement table).
      6. ONE scheme: the mapping built by ingest restores every planted
         value, through the same `deanonymize_text` transcripts use, and the
         placeholder for a value shared across formats is IDENTICAL in all
         of them.
      7. Deterministic: ingesting the same bytes twice produces a
         byte-identical `.anon_` artifact.
      8. Read-only: the source file's sha256 is unchanged.
      9. Failure is typed and never silent: an unsupported format, an
         image, a password-protected PDF and a document that yields no text
         each raise their own error naming the format.
    """
    name = "document_formats_converted_and_scrubbed"
    engine = _engine()
    from pii import ingest  # noqa: PLC0415

    problems: list = []
    evidence: list = []

    with tempfile.TemporaryDirectory(prefix="pii_eval_ingest_") as td:
        root = Path(td)
        fixtures = _build_document_fixtures(root)

        missing = [f for f in ("pdf", "docx", "pptx", "xlsx", "csv") if f not in fixtures]
        if missing:
            problems.append(f"fixture build produced no {missing}")

        # ONE session across every format — that is what makes a value seen
        # in a spreadsheet and in a deck share a placeholder.
        session = _new_session(engine)
        shared_placeholders: dict = {}

        for fmt in ("pdf", "docx", "pptx", "xlsx", "csv"):
            src = fixtures[fmt]
            digest_before = hashlib.sha256(src.read_bytes()).hexdigest()

            # Which planted values this particular document actually carries.
            # The CSV is a bare table with no title row, so it has no client
            # name to redact — asserting a value is restored when it was
            # never in the document is a fixture bug, not a finding. The
            # floor below stops that leniency hiding an empty extraction.
            extracted = ingest.extract_text(src)
            present = [v for v in DOC_MUST_NOT_LEAK if v in extracted]
            if len(present) < 8:
                problems.append(
                    f"{fmt}: extract carries only {len(present)} of the "
                    f"{len(DOC_MUST_NOT_LEAK)} planted values — the document did "
                    f"not convert properly, so the checks below prove nothing"
                )
            if DOC_TABULAR_PERSON not in extracted:
                problems.append(
                    f"{fmt}: the tabular person name never reached the extract — "
                    f"the fixture is not testing what it claims to"
                )

            result = ingest.ingest_file(src, session=session)
            expected_path = ingest.anon_path_for(src)

            # 1. output contract
            if result.anon_path != expected_path or not expected_path.is_file():
                problems.append(f"{fmt}: expected {expected_path.name}, got {result.anon_path.name}")
                continue
            anonymized = expected_path.read_text(encoding="utf-8")

            # 2. no planted raw value survives
            leaked = [v for v in present if v in anonymized]
            if leaked:
                problems.append(f"{fmt}: leaked {leaked!r}")
                evidence.extend(f"{fmt}: leaked raw value: {v}" for v in leaked)

            # 3. the TABULAR person name specifically
            if DOC_TABULAR_PERSON in anonymized:
                problems.append(
                    f"{fmt}: person name in tabular data ({DOC_TABULAR_PERSON!r}) "
                    f"survived — the rendering is producing a shape the detector misses"
                )

            # 4. structure survived
            marker = DOC_STRUCTURE_MARKER[fmt]
            if marker not in anonymized:
                problems.append(f"{fmt}: structure marker {marker!r} missing from extract")

            # 5. no pipe-table rendering
            if re.search(r"^\|\s*-{3,}", anonymized, re.MULTILINE):
                problems.append(f"{fmt}: tabular data rendered as a markdown pipe table")

            # 6. one scheme — the shared mapping restores everything
            restored = engine.deanonymize_text(anonymized, session.mapping_file_dict())
            unrestored = [v for v in present if v not in restored]
            if unrestored:
                problems.append(f"{fmt}: mapping did not restore {unrestored!r}")
            for value in present:
                for etype, values in session.entity_mapping.items():
                    if value in values:
                        seen = shared_placeholders.setdefault(value, values[value])
                        if seen != values[value]:
                            problems.append(
                                f"{fmt}: {value!r} has placeholder {values[value]} "
                                f"but {seen} elsewhere — schemes have diverged"
                            )
                        break

            # 7. determinism — same bytes in, byte-identical artifact out
            first_bytes = expected_path.read_bytes()
            replay_session = _new_session(engine)
            ingest.ingest_file(src, session=replay_session,
                                output_dir=root / ("replay_%s" % fmt))
            replay_bytes = ingest.anon_path_for(src, root / ("replay_%s" % fmt)).read_bytes()
            if replay_bytes != first_bytes:
                problems.append(f"{fmt}: second ingest of the same file differed byte-for-byte")

            # 8. the original was never modified
            if hashlib.sha256(src.read_bytes()).hexdigest() != digest_before:
                problems.append(f"{fmt}: SOURCE FILE WAS MODIFIED by ingest")

            evidence.append(
f"{fmt}: {src.name} -> {expected_path.name} "
                f"({result.extracted_chars} chars extracted, "
                f"{len(present)} planted values present, all redacted and restored)"
            )

        # 9. typed errors — never silent, never empty output
        unsupported = root / "deck.key"
        unsupported.write_bytes(b"not really a keynote file")
        try:
            ingest.extract_text(unsupported)
            problems.append(".key: no error raised for an unsupported format")
        except ingest.UnsupportedFormatError as exc:
            if "KEY" not in exc.message.upper():
                problems.append(f".key: error does not name the format: {exc.message!r}")
        except Exception as exc:  # noqa: BLE001
            problems.append(f".key: wrong error type {type(exc).__name__}")

        # Images are SUPPORTED since #163 — a truncated one is a damaged file,
        # not an unsupported format, and must still fail loudly rather than be
        # reported as "successfully anonymised, 0 bytes".
        image = root / "screenshot.png"
        image.write_bytes(b"\x89PNG\r\n\x1a\n")
        try:
            ingest.extract_text(image)
            problems.append(".png: no error raised for a damaged image")
        except ingest.ExtractionFailedError as exc:
            if "PNG" not in exc.message.upper():
                problems.append(f".png: error does not name the format: {exc.message!r}")
        except Exception as exc:  # noqa: BLE001
            problems.append(f".png: wrong error type {type(exc).__name__}")

        encrypted = root / "locked.pdf"
        encrypted.write_bytes(_build_minimal_pdf([[DOC_TITLE]], encrypted=True))
        try:
            ingest.extract_text(encrypted)
            problems.append("encrypted pdf: no error raised")
        except ingest.ExtractionFailedError:
            pass
        except Exception as exc:  # noqa: BLE001
            problems.append(f"encrypted pdf: wrong error type {type(exc).__name__}")

        blank = root / "blank.csv"
        blank.write_text("\n\n", encoding="utf-8")
        try:
            ingest.extract_text(blank)
            problems.append("textless file: no error raised — an empty scrub reported as success")
        except ingest.EmptyExtractionError:
            pass
        except Exception as exc:  # noqa: BLE001
            problems.append(f"textless file: wrong error type {type(exc).__name__}")

    ok = not problems
    detail = ("; ".join(problems) if problems else
              "5/5 formats (pdf, docx, pptx, xlsx, csv) converted, scrubbed, "
              "structure preserved, tabular person name redacted, deterministic, "
              "sources unmodified, one shared mapping restores all; 4 typed "
              "error paths (unsupported, image, encrypted pdf, textless) fire")
    return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True,
                        detail=detail, evidence=evidence)


# --- #163 image ingest: fixtures, built programmatically -------------------
#
# Never committed as binaries — same rule as #162's document fixtures, and the
# real images under `engagements/` are real client data, gitignored, and
# blocked by the guard. Everything here is drawn from scratch with Pillow.

# The screenshot's planted PII. Same client identity and the same tabular
# person name as the document fixtures, so the two checks are comparable.
IMG_PERSON = DOC_TABULAR_PERSON          # "Aisha Rahman"
IMG_ROLE = "Chief Financial Officer"     # the control: NOT PII, must survive
IMG_EMAIL = "a.rahman@zzzplaceholdermeridian.com"
IMG_PHONE = "(555) 201-4477"
IMG_ACCOUNT = "8834021177"

# Everything the sidecar must not carry and the redacted copy must not show.
IMG_MUST_NOT_LEAK = [CLIENT_FULL, IMG_PERSON, IMG_EMAIL, IMG_PHONE, IMG_ACCOUNT]


def _image_font(size: int):
    """Pillow's own scalable default face.

    Deliberately NOT a system font: `ImageFont.load_default(size=...)` ships
    inside Pillow, so this fixture renders and OCRs identically on a macOS
    laptop and on the CI runner. A system font path would make the check pass
    locally and fail — or worse, quietly degrade — in CI.
    """
    from PIL import ImageFont  # noqa: PLC0415
    return ImageFont.load_default(size=size)


def _build_screenshot(path: Path) -> dict:
    """A synthetic screenshot of a stakeholder table.

    The person name sits in a TABULAR/columnar layout on purpose — the shape
    #162 measured the detector failing on, and the shape tesseract's own page
    segmentation shreds into columns before the detector ever sees it. A
    screenshot of prose would certify the easy case.

    Returns {value: [per-WORD boxes]} — where each word of each planted string
    was drawn, so the check can assert on the PIXELS of those exact regions
    rather than on the mere existence of an output file. Per WORD, not per
    value, because redaction draws one box per OCR word: the space between two
    words carries no ink and is correctly left alone, so measuring the whole
    value's box would score a perfect redaction at ~0.90 and invite someone to
    "fix" it by loosening the threshold.
    """
    from PIL import Image, ImageDraw  # noqa: PLC0415

    font = _image_font(20)
    image = Image.new("RGB", (1500, 230), "white")
    draw = ImageDraw.Draw(image)
    boxes = {}

    def put(x, y, text, track=True):
        draw.text((x, y), text, font=font, fill="black")
        if not track:
            return
        word_boxes = []
        cursor = x
        for word in text.split(" "):
            if word:
                word_boxes.append(draw.textbbox((cursor, y), word, font=font))
            cursor += draw.textlength(word + " ", font=font)
        boxes[text] = word_boxes

    put(24, 18, CLIENT_FULL)
    draw.text((360, 18), "- Stakeholder Directory", font=font, fill="black")
    for x, header in ((24, "Name"), (300, "Role"), (640, "Email"), (1260, "Account")):
        draw.text((x, 70), header, font=font, fill="black")
    put(24, 112, IMG_PERSON)
    put(300, 112, IMG_ROLE)          # tracked as the CONTROL region
    put(640, 112, IMG_EMAIL)
    put(1260, 112, IMG_ACCOUNT)
    draw.text((24, 160), "Direct line:", font=font, fill="black")
    put(190, 160, IMG_PHONE)

    image.save(str(path), format="PNG")
    return boxes


def _regions_are_filled(original, redacted, word_boxes, *, threshold: float = 0.98) -> bool:
    """Is every WORD of a planted value actually covered over in the copy?"""
    return bool(word_boxes) and all(
        _region_is_filled(original, redacted, box, threshold=threshold)
        for box in word_boxes
    )


def _region_is_filled(original, redacted, box, *, threshold: float = 0.98) -> bool:
    """Is the word drawn in `box` actually covered over in the redacted copy?

    Asserts on PIXELS, not on a file existing — a straight copy, a resize, or
    a box drawn in the wrong place all fail here.

    The signal is the BACKGROUND, not the ink. Black fill over black glyphs is
    indistinguishable from the glyphs themselves, so "is the ink black?" would
    pass on an untouched image. What proves a fill landed is that the WHITE
    space between and around the letters is now black. `box` is also widened
    restricted to the rows that actually carried ink first: `ImageDraw.textbbox`
    reports the font's ascender-to-descender line box, which is taller than the
    glyphs and taller than the OCR word box, so measuring it raw would count
    untouched padding as a miss.
    """
    left, top, right, bottom = box
    before = original.convert("L").crop((left, top, right, bottom))
    after = redacted.convert("L").crop((left, top, right, bottom))
    width = before.width
    if width == 0 or before.height == 0:
        return False

    before_pixels = list(before.getdata())
    after_pixels = list(after.getdata())

    # Rows that carried ink in the original — the glyph band.
    ink_rows = [y for y in range(before.height)
                if any(before_pixels[y * width + x] < 128 for x in range(width))]
    if not ink_rows:
        return False

    background = 0
    covered = 0
    for y in ink_rows:
        for x in range(width):
            if before_pixels[y * width + x] >= 200:      # was background
                background += 1
                if after_pixels[y * width + x] < 16:     # is now fill
                    covered += 1
    if not background:
        return False
    return covered / float(background) >= threshold


def _regions_untouched(original, redacted, word_boxes) -> bool:
    """Are these boxes pixel-identical between the original and the copy?

    The over-redaction guard: a redactor that blanked the whole image would
    satisfy every "is it filled?" assertion above and destroy the artifact's
    entire reason for existing.
    """
    return all(
        original.convert("RGB").crop(box).tobytes()
        == redacted.convert("RGB").crop(box).tobytes()
        for box in word_boxes
    )


def _build_docx_with_image(directory: Path, image: Path) -> Path:
    """A Word document whose only PII is inside an embedded picture."""
    import docx  # noqa: PLC0415
    document = docx.Document()
    document.add_heading("Review pack", level=1)
    document.add_paragraph("The console screenshot below was supplied by the client.")
    document.add_picture(str(image))
    document.add_paragraph("End of pack.")
    target = directory / "pack.docx"
    document.save(str(target))
    return target


def _build_pptx_with_image(directory: Path, image: Path) -> Path:
    """A deck whose only PII is inside an embedded picture."""
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Console review"
    slide.shapes.add_picture(str(image), Inches(0.4), Inches(1.6), width=Inches(9))
    target = directory / "deck.pptx"
    presentation.save(str(target))
    return target


def _image_input_produces_sidecar_and_redacted_copy(target: str) -> CheckResult:  # noqa: ARG001
    """#163: an image produces BOTH an anonymised text sidecar and a redacted
    copy, from one local OCR pass, with the round-trip carried by the text.

    Eleven properties, deferred here by #161 and gated at 1.00:

      1. BOTH artifacts are written, at the paths `anon_path_for` and
         `redacted_image_path_for` predict.
      2. The OCR actually read the planted values — otherwise everything
         below would pass vacuously against an empty extract.
      3. No planted raw value survives into the sidecar, including the person
         name in TABULAR layout (the shape the detector misses and the shape
         tesseract shreds into columns).
      4. ONE scheme: the same session's mapping restores every planted value
         out of the sidecar, via the same `deanonymize_text` transcripts use.
      5. The redacted copy's PIXELS changed where the PII was — each planted
         value's drawn region is blanked. Not "a file was written": a straight
         copy fails property 5 outright.
      6. It is a redaction, not a wipe — the non-PII control region (a job
         title) is pixel-identical to the original.
      7. Originals are never modified (sha256 before/after).
      8. Deterministic: the same image twice gives byte-identical artifacts,
         both the sidecar and the PNG.
      9. Images embedded in a DOCX and a PPTX route through the IDENTICAL
         path — their text lands in the host document's sidecar and each
         picture gets its own redacted copy.
     10. The Flow C logo notice fires ONCE PER ENGAGEMENT, not once per image.
     11. A missing local OCR binary raises the typed `OCRUnavailableError`
         and REFUSES — no artifact is written and the raw image is not passed
         through unread.
    """
    name = "image_input_produces_sidecar_and_redacted_copy"
    engine = _engine()
    from PIL import Image  # noqa: PLC0415
    from pii import ingest  # noqa: PLC0415

    problems: list = []
    evidence: list = []

    with tempfile.TemporaryDirectory(prefix="pii_eval_image_") as td:
        root = Path(td)
        source = root / "shot.png"
        boxes = _build_screenshot(source)
        digest_before = hashlib.sha256(source.read_bytes()).hexdigest()

        ingest.reset_engagement_notices()

        # 2. the OCR has to have read the planted values first
        extracted = ingest.extract_text(source)
        present = [v for v in IMG_MUST_NOT_LEAK if v in extracted]
        if len(present) < len(IMG_MUST_NOT_LEAK):
            problems.append(
                "OCR read only %d of the %d planted values %r — the checks below "
                "would prove nothing" % (len(present), len(IMG_MUST_NOT_LEAK),
                                          [v for v in IMG_MUST_NOT_LEAK if v not in extracted])
            )
        if IMG_PERSON not in extracted:
            problems.append(
                "the person name in tabular layout never reached the OCR text — "
                "the fixture is not testing what it claims to"
            )

        session = _new_session(engine)
        notices = io.StringIO()
        result = ingest.ingest_file(
            source, session=session, engagement_dir=root, notice_stream=notices,
        )

        # 1. both artifacts, at the predicted paths
        expected_sidecar = ingest.anon_path_for(source)
        expected_redacted = ingest.redacted_image_path_for(source)
        if result.anon_path != expected_sidecar or not expected_sidecar.is_file():
            problems.append("sidecar: expected %s, got %s"
                            % (expected_sidecar.name, result.anon_path.name))
        if result.redacted_path != expected_redacted or not expected_redacted.is_file():
            problems.append("redacted copy: expected %s, got %s"
                            % (expected_redacted.name,
                               result.redacted_path.name if result.redacted_path else None))

        if expected_sidecar.is_file():
            sidecar = expected_sidecar.read_text(encoding="utf-8")

            # 3. nothing planted survives — the tabular name called out by name
            leaked = [v for v in present if v in sidecar]
            if leaked:
                problems.append("sidecar leaked %r" % leaked)
            if IMG_PERSON in sidecar:
                problems.append(
                    "the person name in tabular layout (%r) survived into the "
                    "sidecar — the OCR reflow is producing a shape the detector "
                    "misses" % IMG_PERSON
                )
            # An empty or contentless sidecar is not a pass: the non-PII
            # content has to be there, and so do the placeholders.
            if IMG_ROLE not in sidecar:
                problems.append(
                    "sidecar does not carry the non-PII content (%r) — it is "
                    "empty or the OCR text never reached it" % IMG_ROLE
                )
            if not re.search(r"<[A-Z][A-Z0-9_]*_\d+>", sidecar):
                problems.append("sidecar carries no placeholders at all")

            # 4. one scheme — the shared mapping restores everything
            restored = engine.deanonymize_text(sidecar, session.mapping_file_dict())
            unrestored = [v for v in present if v not in restored]
            if unrestored:
                problems.append("mapping did not restore %r from the sidecar" % unrestored)

        if expected_redacted.is_file():
            original_image = Image.open(str(source))
            redacted_image = Image.open(str(expected_redacted))

            # 5. the PIXELS changed, in the PII regions specifically
            if expected_redacted.read_bytes() == source.read_bytes():
                problems.append("the redacted copy is a byte-for-byte copy of the original")
            for value in (CLIENT_FULL, IMG_PERSON, IMG_EMAIL, IMG_PHONE, IMG_ACCOUNT):
                box = boxes.get(value)
                if box is None:
                    continue
                if not _regions_are_filled(original_image, redacted_image, box):
                    problems.append(
                        "the region where %r was drawn is still legible in the "
                        "redacted copy" % value
                    )

            # 6. redaction, not a wipe
            if not _regions_untouched(original_image, redacted_image, boxes[IMG_ROLE]):
                problems.append(
                    "the non-PII control region (%r) was altered — this is a wipe, "
                    "not a redaction" % IMG_ROLE
                )
            evidence.append(
                "shot.png -> %s + %s (%d region(s) blanked, %d planted values "
                "redacted and restored)" % (expected_sidecar.name,
                                            expected_redacted.name,
                                            result.regions_redacted, len(present))
            )

        # 7. the original was never modified
        if hashlib.sha256(source.read_bytes()).hexdigest() != digest_before:
            problems.append("SOURCE IMAGE WAS MODIFIED by ingest")

        # 8. determinism — both artifacts, byte for byte
        replay_dir = root / "replay"
        replay_session = _new_session(engine)
        ingest.ingest_file(source, session=replay_session, engagement_dir=root,
                           output_dir=replay_dir, notice_stream=io.StringIO())
        for label, first, second in (
            ("sidecar", expected_sidecar, ingest.anon_path_for(source, replay_dir)),
            ("redacted copy", expected_redacted,
             ingest.redacted_image_path_for(source, replay_dir)),
        ):
            if not (first.is_file() and second.is_file()):
                continue
            if first.read_bytes() != second.read_bytes():
                problems.append("%s: a second ingest of the same image differed "
                                "byte-for-byte" % label)

        # 9. embedded images route through the IDENTICAL path
        for fmt, builder in (("docx", _build_docx_with_image),
                             ("pptx", _build_pptx_with_image)):
            host = builder(root, source)
            host_digest = hashlib.sha256(host.read_bytes()).hexdigest()
            host_session = _new_session(engine)
            host_result = ingest.ingest_file(
                host, session=host_session, engagement_dir=root,
                notice_stream=io.StringIO(),
            )
            host_text = host_result.anon_path.read_text(encoding="utf-8")
            if ingest.IMAGE_SEAM_MARKER in host_text:
                problems.append("%s: the embedded picture was left unread" % fmt)
            if IMG_ROLE not in host_text:
                problems.append("%s: the embedded picture's text never reached "
                                "the document sidecar" % fmt)
            embedded_leaks = [v for v in IMG_MUST_NOT_LEAK if v in host_text]
            if embedded_leaks:
                problems.append("%s: embedded picture leaked %r" % (fmt, embedded_leaks))
            if not host_result.redacted_paths:
                problems.append("%s: no redacted copy was written for the "
                                "embedded picture" % fmt)
            else:
                copy = host_result.redacted_paths[0]
                if copy != ingest.redacted_image_path_for(host, index=1):
                    problems.append("%s: redacted copy named %s, expected %s"
                                    % (fmt, copy.name,
                                       ingest.redacted_image_path_for(host, index=1).name))
                if copy.is_file() and not _regions_are_filled(
                    Image.open(str(source)), Image.open(str(copy)), boxes[IMG_PERSON]
                ):
                    problems.append("%s: the embedded picture's copy still shows "
                                    "the person name" % fmt)
            if hashlib.sha256(host.read_bytes()).hexdigest() != host_digest:
                problems.append("%s: SOURCE DOCUMENT WAS MODIFIED by ingest" % fmt)
            evidence.append("%s: embedded picture OCR'd into %s, redacted copy %s"
                            % (fmt, host_result.anon_path.name,
                               host_result.redacted_paths[0].name
                               if host_result.redacted_paths else "MISSING"))

        # 10. the logo notice — once per ENGAGEMENT, not once per image
        ingest.reset_engagement_notices()
        notice_stream = io.StringIO()
        second_image = root / "shot2.png"
        _build_screenshot(second_image)
        flags = []
        for image_path in (source, second_image,
                           root / "pack.docx", root / "deck.pptx"):
            flags.append(ingest.ingest_file(
                image_path, session=_new_session(engine), engagement_dir=root,
                output_dir=root / "notice", notice_stream=notice_stream,
            ).logo_notice_shown)
        printed = notice_stream.getvalue()
        if printed.count("About screenshots") != 1:
            problems.append(
                "the logo notice printed %d times across 4 image ingests in one "
                "engagement — it must fire once per engagement, not per image"
                % printed.count("About screenshots")
            )
        if flags[:1] != [True] or any(flags[1:]):
            problems.append("logo-notice flags across one engagement were %r — "
                            "expected only the first ingest to report it" % flags)
        if ingest.LOGO_NOTICE.strip() not in printed:
            problems.append("the notice printed is not the Flow C copy verbatim")
        if "logo" not in printed.lower():
            problems.append("the notice does not state the logo limitation")

        # 11. no local OCR -> typed refusal, nothing written, image not read
        import pytesseract  # noqa: PLC0415
        real_cmd = pytesseract.pytesseract.tesseract_cmd
        refusal_dir = root / "refused"
        refusal_dir.mkdir()
        pytesseract.pytesseract.tesseract_cmd = str(root / "no-such-ocr-binary")
        try:
            try:
                ingest.ingest_file(source, session=_new_session(engine),
                                   engagement_dir=root, output_dir=refusal_dir,
                                   notice_stream=io.StringIO())
                problems.append("no OCR installed: the image was NOT refused")
            except ingest.OCRUnavailableError as exc:
                if not isinstance(exc, ingest.IngestError) or not exc.message:
                    problems.append("no OCR installed: error carries no message")
                prose = exc.message.lower()
                for banned in ("presidio", "spacy", "pytesseract", "pip install"):
                    if banned in prose:
                        problems.append(
                            "no OCR installed: message names %r — copy rule 1 "
                            "keeps tool names out of consultant-facing text"
                            % banned
                        )
            except Exception as exc:  # noqa: BLE001
                problems.append("no OCR installed: wrong error type %s"
                                % type(exc).__name__)
            if list(refusal_dir.iterdir()):
                problems.append("no OCR installed: artifacts were written anyway (%r)"
                                % [p.name for p in refusal_dir.iterdir()])
        finally:
            pytesseract.pytesseract.tesseract_cmd = real_cmd

    ok = not problems
    detail = ("; ".join(problems) if problems else
              "image -> sidecar + redacted copy; tabular person name, client "
              "name, email, phone and account redacted in the text and blanked "
              "in the pixels; non-PII region untouched; mapping restores all; "
              "deterministic; source unmodified; embedded docx/pptx pictures "
              "route through the same path; logo notice once per engagement; "
              "missing OCR refuses with a typed error and writes nothing")
    return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True,
                        detail=detail, evidence=evidence)


# A name in a script the local `eng`-only OCR cannot read at all — the
# module docstring's own measured Sinhala example ("නිමල් පෙරේරා"). Rendered
# with `_image_font`, the same Pillow-bundled, OS-independent face every
# other OCR fixture in this file uses (see `_image_font`'s own docstring for
# why that matters): the point measured here is that `eng`-only tesseract
# cannot read this script AT ALL, which holds regardless of font, and using
# the bundled face keeps this reproducing identically in CI.
UNREADABLE_SCRIPT_NAME = "නිමල් පෙරේරා"


def _build_unreadable_script_screenshot(path: Path) -> None:
    """A synthetic screenshot whose only content is a name in a script the
    local OCR cannot read — see `UNREADABLE_SCRIPT_NAME`."""
    from PIL import Image, ImageDraw  # noqa: PLC0415

    font = _image_font(20)
    image = Image.new("RGB", (500, 80), "white")
    draw = ImageDraw.Draw(image)
    draw.text((20, 20), UNREADABLE_SCRIPT_NAME, font=font, fill="black")
    image.save(str(path), format="PNG")


def _build_textless_screenshot(path: Path) -> None:
    """A synthetic screenshot with NO text at all — a chart/logo stand-in
    (shapes only). Used to prove the low-confidence refusal does not
    misfire on a genuinely empty image: that case must stay an ordinary
    `EmptyExtractionError`, not `OCRLowConfidenceError` — see the module
    docstring's own "must stay distinguishable" note."""
    from PIL import Image, ImageDraw  # noqa: PLC0415

    image = Image.new("RGB", (400, 200), "white")
    draw = ImageDraw.Draw(image)
    draw.ellipse([40, 40, 360, 160], outline="black", width=4)
    draw.rectangle([140, 80, 260, 120], outline="black", width=4)
    image.save(str(path), format="PNG")


def _image_unreadable_script_refuses_and_writes_nothing(target: str) -> CheckResult:  # noqa: ARG001
    """#173: an image whose text cannot be confidently read is REFUSED, not
    silently passed through unredacted. Closes the leak #163 measured and
    documented in ingest.py's module docstring ("NON-LATIN SCRIPT IN
    IMAGES"): a non-Latin-script name transliterates into noise the detector
    cannot recognise, so — before this ticket — no box was drawn and the real
    name stayed fully legible in the redacted copy: an artifact whose
    `.anon_` prefix asserted a scrub that never happened.

    Five properties, gated at 1.00:

      1. Ingesting the unreadable-script image raises `OCRLowConfidenceError`
         — a typed `IngestError`, never a bare exception, and never a
         silent "success".
      2. NO sidecar (`.anon_...md`) is written for it.
      3. NO redacted image copy (`.anon_...png`) is written for it — the
         specific leak this check exists to close.
      4. The refusal message is plain language, names no tool, and states
         the image will not be used (ux-design-v6.md Copy Rules).
      5. The refusal is NOT triggered by a genuinely textless image (a
         chart/logo stand-in): that case must still reach the ordinary
         `EmptyExtractionError` path, proving the two failure modes stay
         distinguishable rather than collapsing into one.
    """
    name = "image_unreadable_script_refuses_and_writes_nothing"
    from pii import ingest  # noqa: PLC0415

    problems: list = []
    evidence: list = []

    with tempfile.TemporaryDirectory(prefix="pii_eval_unreadable_") as td:
        root = Path(td)
        source = root / "shot.png"
        _build_unreadable_script_screenshot(source)
        digest_before = hashlib.sha256(source.read_bytes()).hexdigest()

        ingest.reset_engagement_notices()
        engine = _engine()

        # 1 + 4. the typed refusal, with a plain-language, tool-free message
        try:
            ingest.ingest_file(source, session=_new_session(engine),
                               engagement_dir=root, notice_stream=io.StringIO())
            problems.append(
                "an unreadable-script image was NOT refused — ingest_file "
                "returned normally instead of raising"
            )
        except ingest.OCRLowConfidenceError as exc:
            if not isinstance(exc, ingest.IngestError) or not exc.message:
                problems.append("refusal carries no message")
            else:
                prose = exc.message.lower()
                for banned in ("presidio", "spacy", "pytesseract", "tesseract",
                               "pip install", "ocr"):
                    if banned in prose:
                        problems.append(
                            "message names %r — copy rule 1 keeps tool names "
                            "out of consultant-facing text" % banned
                        )
                for required in ("could not be read", "must not be opened"):
                    if required not in prose:
                        problems.append("message is missing %r" % required)
                evidence.append("refusal message: %r" % exc.message)
        except Exception as exc:  # noqa: BLE001
            problems.append("wrong error type %s raised (expected "
                            "OCRLowConfidenceError)" % type(exc).__name__)

        # 2 + 3. no artifacts at all — the directory holds only the source
        written = sorted(p.name for p in root.iterdir() if p != source)
        if written:
            problems.append("artifacts were written despite the refusal: %r" % written)
        anon_sidecar = ingest.anon_path_for(source)
        if anon_sidecar.is_file():
            problems.append("a .anon_ sidecar was produced anyway")
        anon_image = ingest.redacted_image_path_for(source)
        if anon_image.is_file():
            problems.append(
                "a redacted .anon_ image copy was produced anyway — this is "
                "the exact #163 leak #173 exists to close"
            )
        evidence.append("directory listing after refusal: %r" %
                        sorted(p.name for p in root.iterdir()))

        if hashlib.sha256(source.read_bytes()).hexdigest() != digest_before:
            problems.append("SOURCE IMAGE WAS MODIFIED by ingest")

        # 5. a genuinely textless image must NOT take this refusal path
        textless = root / "chart.png"
        _build_textless_screenshot(textless)
        try:
            ingest.ingest_file(textless, session=_new_session(engine),
                               engagement_dir=root, notice_stream=io.StringIO())
            problems.append("a textless image was NOT refused at all — "
                            "expected EmptyExtractionError")
        except ingest.OCRLowConfidenceError:
            problems.append(
                "a textless image (no words recognised) raised "
                "OCRLowConfidenceError — it must fall through to "
                "EmptyExtractionError instead; the two failure modes have "
                "collapsed into one"
            )
        except ingest.EmptyExtractionError:
            evidence.append("textless image: correctly EmptyExtractionError, "
                            "not OCRLowConfidenceError")
        except Exception as exc:  # noqa: BLE001
            problems.append("textless image raised the wrong error type %s"
                            % type(exc).__name__)

    ok = not problems
    detail = ("; ".join(problems) if problems else
              "unreadable-script image raises OCRLowConfidenceError; no "
              "sidecar and no redacted copy written; source unmodified; "
              "message is plain-language, tool-free and states the image is "
              "blocked; a textless image still takes the ordinary "
              "EmptyExtractionError path, not this refusal")
    return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True,
                        detail=detail, evidence=evidence)


# --- guard_fails_closed_on_inputs_path (#164) -------------------------------
#
# `.claude/hooks/anonymize-guard.py` itself, invoked as a subprocess — see
# module docstring's "Check 14" note for why this does not go through
# `_engine()`/`_facade()` like every other check here.

ANONYMIZE_GUARD_HOOK_REL_PATH = Path(".claude") / "hooks" / "anonymize-guard.py"


def _anonymize_guard_hook_path() -> Path:
    return repo_root() / ANONYMIZE_GUARD_HOOK_REL_PATH


def _read_tool_payload(file_path: str) -> bytes:
    return pretooluse_payload("Read", {"file_path": file_path})


def _run_anonymize_guard(project_dir: Path, stdin_bytes: bytes) -> HookRun:
    """This row's ONE subprocess entry point for `anonymize-guard.py` — both
    the fail-closed check and `runs_under_registered_interpreter` go through
    it, so the latter observes the argv this row really spawns rather than a
    re-implementation written for it.

    Everything about the invocation (registered interpreter, argv recording,
    CLAUDE_PROJECT_DIR, timeout, stdout parsing) lives in
    `rubrics._harness.run_hook_subprocess`; `mcp_query_guard.py`'s `_run_hook`
    is the identical wrapper over the identical helper, which is what #196
    extracted. Settings.json registers this hook as bare `python3`,
    deliberately kept off the Presidio venv per solution-design-v6.md D13 —
    see the harness docstring for why that must not become `sys.executable`.
    """
    return run_hook_subprocess(_anonymize_guard_hook_path(), stdin_bytes,
                               project_dir=project_dir)


def _guard_fails_closed_on_inputs_path(target: str) -> CheckResult:  # noqa: ARG001
    """The #164 rewrite's central contract (.design/solution-design-v6.md
    D13): ANY unexpected failure while evaluating a path under
    engagements/*/inputs/ must DENY; the identical failure for a path
    OUTSIDE that scope must ALLOW — the guard is fail-closed on raw client
    material only, never globally (a globally fail-closed guard wedged
    every session once already, PR #82).

    Fault injection: chmod the `inputs/` directory (and, for the control
    case, an unrelated directory outside engagements/) to 000 so any
    stat()/exists() call inside raises a real, unmocked PermissionError —
    not a simulated fault. Both files exist and are genuinely unscrubbed
    (no `.anon_` sibling), so absent the fault they would both be denied
    for the ORDINARY reason (missing sibling); the fault must change the
    outside-inputs/ outcome to allow while the inputs/ outcome stays deny,
    proving the fail-closed/fail-open SPLIT, not just that denial happens
    somewhere.

    Skips (never false-passes or false-fails) when running as root, which
    bypasses directory permission bits entirely.
    """
    name = "guard_fails_closed_on_inputs_path"
    skip = fault_injection_skip(name, perms="directory")
    if skip is not None:
        return skip

    hook = _anonymize_guard_hook_path()
    if not hook.exists():
        return bool_check(name, False, detail=f"{hook} not found — cannot run the check")

    with tempfile.TemporaryDirectory(prefix="anonymize_guard_eval_") as td:
        root = Path(td)
        fixture = build_fixture_engagement(
            root,
            slug=SEEDED_CLIENT_DIR,
            client_name=CLIENT_FULL,
            short_name=CLIENT_SHORT,
            engagement="2026-01_test_engagement",
            # Genuinely unscrubbed: no `.anon_` sibling. Absent the fault BOTH
            # files would be denied for the ORDINARY reason (missing sibling),
            # which is what makes the fault the only variable below.
            documents={"inputs/fault_test.md": "placeholder content, no sibling\n"},
        )
        inputs_dir = fixture.inputs_dir
        raw_inputs_file = inputs_dir / "fault_test.md"

        outside_dir = root / "scratch_outside"
        outside_dir.mkdir(parents=True)
        outside_file = outside_dir / "fault_test.md"
        outside_file.write_text("placeholder content, no sibling\n", encoding="utf-8")

        # BOTH directories faulted in ONE block, and that is the check: the
        # contract is a SPLIT — the identical PermissionError must DENY inside
        # engagements/*/inputs/ and ALLOW outside it. Faulting only inputs/
        # would prove nothing more than "denial happens somewhere", and a
        # globally fail-closed guard wedged every session once already (#82).
        with inject_fault(inputs_dir, outside_dir):
            result_inputs = _run_anonymize_guard(root, _read_tool_payload(str(raw_inputs_file)))
            result_outside = _run_anonymize_guard(root, _read_tool_payload(str(outside_file)))

        ok = (
            result_inputs.returncode == 0 and result_inputs.denied
            and result_outside.returncode == 0 and not result_outside.denied
        )
        return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True, detail=(
            f"inputs/ (fault): rc={result_inputs.returncode} denied={result_inputs.denied} "
            f"outside inputs/ (fault): rc={result_outside.returncode} denied={result_outside.denied} "
            f"reason={result_inputs.reason[:120]!r}"
        ))


def _runs_under_registered_interpreter(_target: str) -> CheckResult:  # noqa: ARG001
    """#192/backlog :116 — `anonymize-guard.py` must be invoked under
    whatever interpreter `.claude/settings.json` actually registers for it
    (bare `python3`), never a silent fallback to `sys.executable`. Shared
    check for Python-hook rows — see rubrics/_harness.py.

    `spawn` deliberately goes through this module's own production
    `_run_anonymize_guard`, the same helper the fail-closed check uses, so
    the check observes the argv the row REALLY spawns rather than a
    re-implementation written for it. The probe path sits OUTSIDE
    engagements/*/inputs/ so the guard simply allows it; the decision is
    irrelevant here — only the argv matters."""
    hook = _anonymize_guard_hook_path()

    def _spawn() -> None:
        with tempfile.TemporaryDirectory(prefix="anonymize_guard_interp_") as td:
            probe = Path(td) / "scratch_outside" / "probe.md"
            probe.parent.mkdir(parents=True)
            probe.write_text("placeholder content\n", encoding="utf-8")
            _run_anonymize_guard(Path(td), _read_tool_payload(str(probe)))

    return check_runs_under_registered_interpreter(
        hook, hook_label="anonymize-guard.py", spawn=_spawn
    )


# --- #181 internal-domain email shape recognizer ---------------------------
#
# Deliberately does NOT reuse the committed golden fixture — the leak this
# closes is specifically about TLDs no golden fixture has ever contained
# (every existing fixture uses real-looking TLDs, which is exactly how 14
# checks at 1.00 and 4 gate-bites proofs all passed while this leaked).
# Both halves below run against the real engine, in a fresh session, with
# the standard synthetic client identity — never real client material.

# Addresses that were confirmed LEAKING before the #181 fix — internal and
# RFC 2606 reserved TLDs Presidio's built-in EmailRecognizer rejects via
# tldextract because they are not real, registered public-suffix entries.
INTERNAL_DOMAIN_EMAIL_CASES = [
    "j.smith@zzzplaceholderbank.internal",
    "j.smith@zzzplaceholderbank.corp",
    "j.smith@zzzplaceholderbank.local",
    "j.smith@zzzplaceholderbank.lan",
    "j.smith@zzzplaceholderbank.intranet",
    "p.nair@zzzplaceholderbank.test",
    "p.nair@zzzplaceholderbank.example",
]

# Real-TLD addresses that already worked before #181 — must show NO
# regression: still a single EMAIL_ADDRESS entity, still restores clean.
# The overlap-dedup guarantee (module docstring's "INTERNAL-DOMAIN EMAILS")
# is what this list actually exercises.
REAL_TLD_EMAIL_CASES = [
    "j.smith@zzzplaceholderbank.com",
    "j.smith@zzzplaceholderbank.co.uk",
]

# Plausible non-email strings containing "@" and a dot — the over-detection
# guard. Every one of these must survive UNCHANGED. Deliberately excludes
# the scp-style "user@host.example:/path" case — see module docstring: that
# string is genuinely shape-identical to a real internal-domain email and
# the fix redacts it on purpose, so asserting it here would be asserting
# the wrong thing, not a real over-detection bug.
MUST_NOT_OVER_DETECT = [
    "install pkg@1.2.3 from npm",              # version pin, numeric TLD-slot
    "@app.route('/users')\ndef list_users(): pass",  # decorator + sentence
    "@property\ndef area(self):\n    return self._radius ** 2",  # python decorator
    "Follow @jack_doe for updates.",            # Twitter-style handle
    "Reach me @ the office later.",             # bare @ in prose
]


def _internal_domain_email_redacted_no_over_detection(target: str) -> CheckResult:  # noqa: ARG001
    """#181 gate: internal-domain / reserved-TLD email addresses are
    redacted as EMAIL_ADDRESS and restore byte-identical (the leak); real-TLD
    addresses still redact as a single clean span, no regression from the
    overlap between the two recognizers; and a set of plausible non-email
    `@`-bearing strings are left completely unchanged (the over-detection
    guard).

    Gate-bites mandatory (see the PR description): reverting
    `_InternalDomainEmailRecognizer`'s registration in `_get_analyzer` must
    make the INTERNAL_DOMAIN_EMAIL_CASES half of this check fail, while
    leaving REAL_TLD_EMAIL_CASES and MUST_NOT_OVER_DETECT passing —
    demonstrating this check actually depends on the new recognizer, not on
    something else in the engine.
    """
    name = "internal_domain_email_redacted_no_over_detection"
    engine = _engine()
    problems = []

    # 1) internal/reserved-TLD addresses: redacted, round-trip byte-identical.
    for addr in INTERNAL_DOMAIN_EMAIL_CASES:
        session = _new_session(engine)
        text = f"Contact {addr} for details."
        anonymized = session.anonymize(text)
        if addr in anonymized:
            problems.append(f"LEAKED (not redacted): {addr!r}")
            continue
        restored = session.deanonymize(anonymized)
        if restored != text:
            problems.append(
                f"round-trip broken for {addr!r}: {restored!r} != {text!r}"
            )
        placeholders = session.entity_mapping.get("EMAIL_ADDRESS", {})
        if placeholders.get(addr, "").rpartition("_")[0] != "<EMAIL_ADDRESS":
            problems.append(
                f"{addr!r} redacted under an unexpected entity — mapping={placeholders!r}"
            )

    # 2) real-TLD addresses: still exactly one clean EMAIL_ADDRESS span each
    #    (the overlap-dedup guarantee — no duplicate/nested placeholder).
    for addr in REAL_TLD_EMAIL_CASES:
        session = _new_session(engine)
        text = f"Contact {addr} for details."
        results = session.analyze(text)
        email_spans = [r for r in results if r.entity_type == "EMAIL_ADDRESS"]
        if len(email_spans) != 1:
            problems.append(
                f"expected exactly 1 EMAIL_ADDRESS span for {addr!r}, got "
                f"{len(email_spans)}: {[(r.start, r.end, r.score) for r in email_spans]!r}"
            )
        anonymized = session.anonymize(text)
        if addr in anonymized:
            problems.append(f"regression — real-TLD address no longer redacted: {addr!r}")
        elif session.deanonymize(anonymized) != text:
            problems.append(f"round-trip broken for real-TLD address {addr!r}")

    # 3) plausible non-email strings: must survive completely unchanged.
    for text in MUST_NOT_OVER_DETECT:
        session = _new_session(engine)
        anonymized = session.anonymize(text)
        if anonymized != text:
            problems.append(f"OVER-DETECTED: {text!r} -> {anonymized!r}")

    ok = not problems
    return CheckResult(name, 1.0 if ok else 0.0, ok, hard_fail=True,
                        detail="; ".join(problems) if problems else (
                            f"{len(INTERNAL_DOMAIN_EMAIL_CASES)} internal/reserved-TLD "
                            f"addresses redacted + round-tripped clean; "
                            f"{len(REAL_TLD_EMAIL_CASES)} real-TLD addresses still a "
                            f"single clean span; {len(MUST_NOT_OVER_DETECT)} lookalike "
                            f"strings left unchanged"
                        ),
                        evidence=[f"issue: {p}" for p in problems])


# --- #166 neutral workspace ------------------------------------------------
#
# The path envelope, not the file contents. `compose_prompt` renders
# `engagement_dir` / `outputs_dir` / `transcript_path` into the invocation
# prompt as VALUES and `run_agent` sets `cwd` to the same client-named
# directory, so a perfectly scrubbed transcript still arrives with the
# client's name in the Runtime Parameters table
# (.design/solution-design-v6.md D6). `pii.identity.materialise_workspace`
# is the answer: a directory in which every path segment is generated by us.
#
# The fixture below is deliberately in the PRE-migration, client-named shape
# (`engagements/zzzplaceholderclient/...`). That is what #166 actually has to
# neutralise — opaque directories arrive with #168 — and it is what makes the
# check bite: if the workspace inherited ANY name from the engagement, from a
# raw input file, or from an `.anon_` artifact (whose filename embeds the raw
# client-named filename it came from), a deny-list term appears in a path
# segment and this fails.

# Synthetic client identity for the workspace fixture. Same
# obviously-placeholder register as every other fixture in this file — no
# fictional bank names, no real ones.
WORKSPACE_CLIENT_DIR = "zzzplaceholderclient"
WORKSPACE_CLIENT_NAME = "Zzzplaceholder Meridian Holdings"
# Planted in the RAW (unscrubbed) input file only. If it turns up anywhere in
# the workspace, a raw file was copied in — under any name, including a
# neutral one, which a path-segment check alone would miss.
WORKSPACE_RAW_MARKER = "RAW-CLIENT-MATERIAL-MUST-NOT-BE-COPIED"


def _identity():
    from pii import identity as _i  # noqa: PLC0415 - stdlib only, #166
    return _i


def _seed_workspace_engagement(root: Path):
    """One engagement whose every client-controlled name is a deny-list term:
    the client directory, a raw input file, the `.anon_` artifacts named after
    it, and an input SUBDIRECTORY.

    `slug=WORKSPACE_CLIENT_DIR` — a CLIENT-NAMED directory, the PRE-migration
    shape — is deliberate and must not be "modernised" to
    `build_fixture_engagement(slug=None)`'s opaque id. An opaque id is exactly
    what `materialise_workspace` is supposed to PRODUCE; handing it one to
    start with would make this check unable to tell a working neutraliser from
    a broken one.
    """
    return build_fixture_engagement(
        root,
        slug=WORKSPACE_CLIENT_DIR,
        client_name=WORKSPACE_CLIENT_NAME,
        short_name=None,
        engagement="2026-08_retail_assessment",
        subdirs=("outputs",),
        documents={
            # A deny-list SOURCE: it exists to hold the client's name, so it
            # must not enter a directory an agent runs in, scrubbed sibling
            # or not.
            "inputs/engagement_intake.md":
                "# Engagement Intake\n\n- **Client Name:** %s\n" % WORKSPACE_CLIENT_NAME,
            # Raw, unscrubbed client material + its scrubbed sibling.
            "inputs/Zzzplaceholder_Meridian_Annual_Report.pdf":
                WORKSPACE_RAW_MARKER + "\n",
            "inputs/.anon_Zzzplaceholder_Meridian_Annual_Report.pdf.md":
                "Annual report for <CLIENT_1>.\n",
            # A plain-text transcript (facade naming: `.anon_<name>`, no added `.md`).
            "inputs/.anon_transcript_1.md": "<PERSON_1> said hello.\n",
            # An image: OCR sidecar + redacted copy. One source, two artifacts.
            "inputs/.anon_Meridian_Screenshot.png.md":
                "OCR text mentioning <CLIENT_1>.\n",
            "inputs/.anon_Meridian_Screenshot.png.png": "redacted-png-bytes\n",
            # A foldered input — the FOLDER name is client-controlled too.
            "inputs/Meridian_Board_Pack/.anon_Meridian_Deck.pptx.md":
                "<CLIENT_1> board deck.\n",
        },
    )


def _workspace_paths_contain_no_client_identifiers(target: str) -> CheckResult:  # noqa: ARG001
    """#166: NO path segment inside a materialised workspace may match a
    deny-list term, and nothing but `.anon_` artifacts may be in it.

    Four assertions, one gate:

      1. The deny-list resolved for the fixture engagement is NON-EMPTY and
         contains the client's name. Asserted FIRST, and now STRUCTURALLY:
         `FixtureEngagement.resolved_deny_terms()` runs
         `_harness.assert_deny_list_non_empty` before it hands back a single
         term, so this check cannot obtain a deny-list that skipped the
         guard — it is no longer the first `if` in a function someone could
         reorder. With an empty deny-list every other assertion here passes
         vacuously, which is precisely how this repo has twice shipped a gate
         scoring 1.000 while certifying nothing. The AssertionError it raises
         surfaces as a RED check via `evaluate()`'s wrapper, naming the
         fixture.
      2. No segment of any path in the workspace — the workspace root
         included, not just the files under it — contains any deny term
         (case-insensitive substring, since `Meridian` hides inside
         `.anon_Meridian_Screenshot.png.md`).
      3. Every file in the workspace is an `.anon_` artifact, and the raw
         input's content marker appears NOWHERE. Checking names alone would
         pass a mutation that copied the raw PDF in under a neutral name.
      4. The image's OCR sidecar and its redacted copy still share one index,
         so the pair `pii.ingest` produces is not split apart.
    """
    name = "workspace_paths_contain_no_client_identifiers"
    identity = _identity()

    problems = []
    with tempfile.TemporaryDirectory(prefix="pii_workspace_eval_") as td:
        root = Path(td)
        fixture = _seed_workspace_engagement(root)

        # 1 — the deny-list must actually carry the identity, or this is
        # vacuous. Raises (-> a RED check naming the fixture) before any term
        # is returned; see the docstring.
        terms = fixture.resolved_deny_terms()

        engagement = fixture.engagement_dir
        workspace = identity.materialise_workspace(engagement)
        try:
            paths = [workspace.path] + sorted(workspace.path.rglob("*"))
            tree = [str(p) for p in paths]

            # 2 — no client identifier in any path segment.
            segments = set()
            for path in paths:
                segments.update(Path(path).parts)
            leaked = sorted(
                "%s (matches deny term %r)" % (segment, term)
                for segment in segments
                for term in terms
                if term.lower() in segment.lower()
            )
            problems.extend(leaked)

            files = [p for p in paths if p.is_file()]

            # 3 — only .anon_ artifacts, and no raw content under any name.
            non_anon = sorted(p.name for p in files if not p.name.startswith(".anon_"))
            if non_anon:
                problems.append("non-.anon_ file(s) in the workspace: %s" % non_anon)
            for path in files:
                if WORKSPACE_RAW_MARKER in path.read_text(encoding="utf-8", errors="replace"):
                    problems.append(
                        "raw client material copied into the workspace as %r "
                        "(neutral name, raw bytes)" % path.name
                    )

            # 4 — the sidecar/redacted-copy pair keeps one shared index.
            stems = {p.name[: -len(p.suffix)] for p in files if p.suffix}
            paired = sorted(
                stem for stem in stems
                if (workspace.inputs / (stem + ".md")).is_file()
                and (workspace.inputs / (stem + ".png")).is_file()
            )
            if len(paired) != 1:
                problems.append(
                    "expected exactly one .md/.png artifact pair sharing an index, got %s"
                    % (paired,)
                )
            if len(files) != 5:
                problems.append(
                    "expected 5 workspace artifacts (3 sidecars + 1 transcript + 1 "
                    "redacted image), got %d: %s" % (len(files), sorted(p.name for p in files))
                )
        finally:
            workspace.cleanup()

    ok = not problems
    return CheckResult(
        name, 1.0 if ok else 0.0, ok, hard_fail=True,
        detail=("; ".join(problems) if problems else (
            "%d deny term(s) (%s); %d workspace path segment(s) checked, none "
            "matched; only .anon_ artifacts present; sidecar/redacted pair intact"
            % (len(terms), ", ".join(repr(t) for t in terms), len(segments))
        )),
        evidence=(["issue: " + p for p in problems] if problems
                  else ["workspace tree: " + t for t in tree]),
    )



def _accented_latin_identity_extracted_whole(_target: str) -> CheckResult:
    """Accented Latin names survive extraction INTACT — the 2026-08-30 widening.

    Asserts through `denylist.resolve_engagement_deny_list`, the function
    `orchestrate.py` and `anonymize_transcript.py` actually call, so this is a
    statement about the shipped extractor rather than about this file.

    Measured before the widening, against the real extractors: a single-word
    accented CLIENT name was not merely mis-cased, it was LOST —
    "Länsförsäkringar" yielded only "kringar", "Bagócs" yielded nothing at all,
    and "Åland" yielded "land", which under-detects the client and emits a
    generic term that over-blocks. The stakeholder path failed separately:
    "José Ramírez" was dropped whole by `_PERSON_TOKEN_RE`.

    Two regexes, two failure modes, one check — because a fix applied to only
    one of them looks complete and still leaks the client, which is the more
    important of the two (solution-design-v6 D3).
    """
    name = "accented_latin_identity_extracted_whole"
    with tempfile.TemporaryDirectory(prefix="pii_accent_") as td:
        root = Path(td)
        client_dir = root / "engagements" / "zzzaccented"
        eng = client_dir / "2026-08_test"
        (eng / "inputs").mkdir(parents=True)
        (client_dir / "CLIENT_PROFILE.md").write_text(
            "# CLIENT_PROFILE\n\n"
            "- **Client Name:** Zzzbagócs\n"
            "- **Client Name:** Zzzlänsförsäkringar\n"
            "- **Primary Contact:** Zzzjosé Zzzramírez, CFO\n",
            encoding="utf-8")
        terms = {t.lower() for t in _denylist().resolve_engagement_deny_list(str(eng))}

    want_client = {"zzzbagócs", "zzzlänsförsäkringar"}
    want_person = "zzzjosé zzzramírez"
    missing_client = sorted(w for w in want_client if w not in terms)
    person_ok = want_person in terms
    # the pre-widening failure emitted ASCII fragments instead of the name
    fragments = sorted(t for t in terms if t in {"kringar", "dito", "land", "bagócs"[1:]})
    ok = not missing_client and person_ok and not fragments
    return bool_check(name, ok, detail=(
        f"missing_client={missing_client} person_extracted={person_ok} "
        f"ascii_fragments={fragments}"))


def evaluate(target: str) -> list:
    fixture = _fixture_path(target)
    if not fixture.exists():
        return [CheckResult(
            "fixture_present", 0.0, False, hard_fail=True,
            detail=f"{fixture} not found — cannot run any check",
        )]

    try:
        _engine()
    except Exception as exc:  # noqa: BLE001 - convert to a reportable failure, not a crash
        return [CheckResult(
            "presidio_importable", 0.0, False, hard_fail=True,
            detail=(
                f"could not import scripts/pii/engine.py ({type(exc).__name__}: {exc}). "
                f"This component requires the Presidio-capable interpreter — run via "
                f".venv/bin/python evals/run_experiment.py --component pii-anonymizer, "
                f"not the system python3. See this module's docstring."
            ),
        )]

    checks = [
        _accented_latin_identity_extracted_whole,
        _round_trip_byte_identical,
        _distinct_values_distinct_placeholders,
        _repeated_value_reuses_placeholder,
        _no_raw_pii_in_anonymized_output,
        _cross_transcript_merge_collision_free,
        _mapping_files_chmod_600_and_cleaned,
        _client_name_redacted_via_denylist,
        _allowlist_prevents_generic_overredaction,
        _empty_entity_list_warns,
        _legacy_flat_mapping_still_restores,
        _document_formats_converted_and_scrubbed,
        _image_input_produces_sidecar_and_redacted_copy,
        _image_unreadable_script_refuses_and_writes_nothing,
        _guard_fails_closed_on_inputs_path,
        _internal_domain_email_redacted_no_over_detection,
        _nested_outputs_deanonymized,
        _xlsx_outputs_deanonymized,
        _workspace_paths_contain_no_client_identifiers,
        _runs_under_registered_interpreter,
    ]
    results = []
    for fn in checks:
        try:
            results.append(fn(str(fixture)))
        except Exception as exc:  # noqa: BLE001 - a raising check must still report, not crash the suite
            results.append(CheckResult(
                fn.__name__.lstrip("_"), 0.0, False, hard_fail=True,
                detail=f"check raised {type(exc).__name__}: {exc}",
            ))
    return results
