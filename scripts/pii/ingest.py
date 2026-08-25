#!/usr/bin/env python3
"""
Document and image ingest — PDF / DOCX / PPTX / XLSX / CSV / images to
anonymised text, plus a redacted copy for every image.

WHAT THIS IS (ticket #162, .prd/prd-v6.md §3, .design/solution-design-v6.md D11)
  The PII gate used to see `.md/.txt/.vtt/.srt/.json/.csv/.log` only. Real
  engagement `inputs/` directories hold 43 PDFs, 10 spreadsheets, 8 Word
  documents, 2 decks — and 3 markdown files. Roughly 3 of 77 real files were
  in scope; annual reports, RFPs, client decks and pricing spreadsheets were
  read into context unscrubbed. This module converts the other 74.

  Extract text -> anonymise through `scripts/pii/engine.py` -> write a
  `.anon_` sibling. Agents read only the `.anon_` artifact; the raw document
  is never read by any agent, in any mode.

ONE ANONYMISATION SCHEME (this is the point of the cycle, not an aside)
  Every extractor returns TEXT, and every byte of that text goes through the
  same `PIISession` that transcripts go through. There is no second scheme,
  no per-format placeholder convention and no per-format mapping. Pass one
  session (or one shared `entity_mapping`) across a whole engagement and the
  same email in a PDF and in a spreadsheet gets the same placeholder — which
  is what `artifact_boundary.deanonymize_dir` depends on when it restores the
  deliverable from a single mapping.

  This is also why `presidio-structured` is NOT used for the XLSX/CSV path —
  see "D11 DECISION" below.

⚠️ THE RENDERING FORMAT IS A DETECTION DECISION, NOT A READABILITY ONE
  Read this before "tidying up" the table rendering into a markdown table.

  spaCy's `en_core_web_lg` misses person names in a SHAPE-DEPENDENT way
  (#159/#161; engine.py's "MEASURED DETECTION LIMITS"). In a markdown table
  cell `| Aisha Rahman | CFO |` it tags the name ORGANIZATION, which is not
  an enabled entity type (enabling it would strip "Backbase", "Temenos" and
  every other vendor name a deliverable depends on), so the name passes
  straight through the gate. The same name on a label line is detected.

  This module turns every spreadsheet row, every Word table and every
  PowerPoint table into text. A markdown pipe table is the natural, readable
  choice and it is exactly the one shape the detector fails on — a client's
  stakeholder or customer list in an XLSX would render as a pipe table and
  leak. So the format was chosen on measurement:

      30 synthetic person names in tabular data, same content, same engine.
      "run-on spans" counts detected entities whose span ran past the value
      into the next line; "labels lost" counts field labels swallowed by one
      of those spans (see `_terminate`):

        rendering                                  PERSON  run-on  labels lost
        ------------------------------------------ ------  ------  -----------
        markdown pipe table (4-col)                 28/30       0            0
        markdown pipe table (2-col)                 28/30       0            0
        tab-separated columns                       28/30       0            0
        one-line-per-row, labelled                  30/30       0            0
        record-per-row label lines (no terminator)  30/30       2            1
        record-per-row label lines + terminator     30/30       0            0

      Every pipe-table and TSV miss included "Aisha Rahman" — the same name
      #159 recorded as tagged ORGANIZATION in a table cell.

  CHOSEN: record-per-row label lines with a sentence terminator
  (`_render_records` + `_terminate`) — the bottom row. It detects every
  name, loses no field labels, is a shape `pii/denylist.py`'s label-line
  extractor already understands, and a record with one field per line is
  genuinely readable: a row's fields stay together and stay labelled, which
  a wall of text would not.

  NOTE — this applies to the formats where the text is CONSTRUCTED from
  structured data: DOCX/PPTX tables, XLSX sheets, CSV rows. A PDF has no
  structured table to re-render; its text is extracted as laid out, so a
  pipe-formatted table inside a PDF stays pipe-formatted and carries that
  shape's detection risk. That is a seam, not an oversight: reflowing a
  PDF's visual layout into records is a guess, and a wrong guess would
  scramble the document.

  If you change `_render_records`, re-run the measurement first. The eval
  check `document_formats_converted_and_scrubbed` plants a person name in
  tabular data in every format specifically so a regression here fails the
  gate rather than passing it.

D11 DECISION — `presidio-structured` 0.0.8 REJECTED for the XLSX/CSV path
  PRD §8 and design D11 asked for this to be settled on evidence, with the
  plain text path as the recorded fallback. Measured against 0.0.8:

  1. **It cannot use our operator, so it cannot share our mapping.**
     `PandasDataProcessor._generate_operator_mapping` constructs its own
     `OperatorsFactory()` as a local, with no injection point, so
     `InstanceCounterAnonymizer` ("entity_counter") is unregistered there:

         InvalidParamError: Invalid operator class 'entity_counter'.

     Using it would mean a SECOND anonymisation scheme with its own
     placeholders and its own mapping, for spreadsheets only — the exact
     fragmentation this cycle exists to remove, and it would break
     cross-file placeholder consistency.

  2. **Its column-level typing is wrong on our data.** It assigns ONE entity
     type per column from a sample. On a 4-column stakeholder sheet:

         StructuredAnalysis(entity_mapping={'Name': 'PERSON',
                                            'Email': 'URL',
                                            'Account': 'DATE_TIME'})

     Email typed as URL, a 10-digit account number typed as DATE_TIME.
     DATE_TIME and URL are both deliberately NOT anonymised here
     (engine.py's DEFAULT_ENTITIES: dates drive every ROI model; the URL
     recogniser partial-matches inside email local parts) — so those two
     columns would have been passed through in cleartext.

  FALLBACK TAKEN (as D11 prescribes): extract cell strings and run them
  through the plain text analyzer/anonymizer — the inverse of what
  `deanonymize_dir` already does for restore. `pandas` and
  `presidio-structured` are not imported by this module at all.

OUTPUT NAMING — `.anon_<full original filename>.<artifact ext>`
  `report.pdf` -> `.anon_report.pdf.md`, not `.anon_report.md`. The source
  extension is KEPT because a real `inputs/` directory holds `Pricing.pdf`
  and `Pricing.xlsx` side by side, and dropping the extension would have one
  anonymised artifact silently overwrite the other. Keeping it also makes the
  mapping raw-path <-> anon-path a pure, invertible function of the
  filename, which is what the guard rewrite (#164) needs to answer "is this
  document scrubbed?" without a directory scan.

  The same rule gives an image's redacted copy its name: `shot.png` ->
  `.anon_shot.png.png` and `shot.jpg` -> `.anon_shot.jpg.png`, so two
  screenshots that differ only in extension cannot collide. The redacted copy
  is ALWAYS written as PNG regardless of the source format — PNG is lossless
  (a JPEG re-encode would smear the boundary of every box) and its bytes are
  deterministic. An embedded image takes its host document's name:
  `.anon_deck.pptx.image1.png`.

  Plain-text inputs (`.md`, `.txt`, `.vtt`, `.srt`, `.log`, `.json`) are NOT
  handled here — they keep going through
  `scripts/anonymize_transcript.anonymize_transcript_file`, which writes
  `.anon_<name>` (no added `.md`). Routing them through both would produce
  two differently-named anonymised artifacts for one input.

NEVER SILENT, NEVER DESTRUCTIVE
  - An unsupported format raises `UnsupportedFormatError` naming the format.
  - A missing local OCR binary raises `OCRUnavailableError` and the image is
    REFUSED — never passed through unread (#163).
  - OCR output the engine cannot trust (mean per-word confidence below
    `_MIN_OCR_CONFIDENCE` — see "NON-LATIN SCRIPT IN IMAGES" below) raises
    `OCRLowConfidenceError` for a direct image and REFUSES it the same way —
    no sidecar, no redacted copy (#173). An embedded picture degrades to
    `IMAGE_SEAM_MARKER` instead of refusing its host document, matching how
    a missing OCR binary is already handled for embedded pictures.
  - A password-protected or corrupt document raises `ExtractionFailedError`.
  - A document that yields NO text raises `EmptyExtractionError` — a scanned
    PDF must not be reported as "successfully anonymised, 0 bytes".
  Nothing is ever skipped silently and no empty `.anon_` file is ever
  written. Every error carries a plain-language `.message` the guard can
  print verbatim (ux-design-v6.md copy rules: consequence before
  instruction, no tool names in the consequence).
  - Extraction is READ-ONLY. The original file is opened for reading and
    never written, moved or re-saved; `.anon_` is a sibling artifact.

DETERMINISTIC
  The same bytes in produce byte-identical text out: extractors walk
  documents in document order, no dict-ordering, no timestamps, no absolute
  paths, and line endings are normalised. `document_formats_converted_and_scrubbed`
  asserts this — a second ingest of the same file must match the first byte
  for byte.

IMAGES — LOCAL OCR, TWO ARTIFACTS, NO CLASSIFIER (#163, design D5, Flow C)
  An image produces BOTH artifacts, always:

      shot.png ──▶ .anon_shot.png.md    OCR text through the same PIISession
                                        (placeholders — carries the round-trip)
               └─▶ .anon_shot.png.png   redacted copy, PII regions filled

  There is deliberately no classifier deciding which branch an image takes.
  Both branches need the same OCR pass anyway, and a classifier would
  sometimes route a UI screenshot to text-only and destroy the reason it was
  added. Images embedded in a DOCX or PPTX go through this identical path
  during document ingest — the `IMAGE_SEAM_MARKER` #162 left is now filled in
  with the picture's OCR text, in document order.

  **THE ROUND-TRIP IS CARRIED BY THE TEXT SIDECAR, NOT THE PIXELS.** Redaction
  of the image copy is destructive by design and that is accepted: the copy is
  never restored and never shipped. Do not attempt reversible pixel redaction.

  **LOCAL ONLY.** `pytesseract` drives a local `tesseract` binary. Presidio's
  Azure Document Intelligence OCR backend is never used — it uploads the image
  to an external service, which defeats the entire purpose of this module.

  ⚠️ **ACCEPTED LIMITATION — LOGOS.** Redaction blanks only what OCR reads as
  TEXT. A client logo is graphics; it is not detected, it is not blanked, and
  it reaches the model. `LOGO_NOTICE` is emitted on the FIRST image ingest of
  an engagement (once per engagement, not once per image) so this surfaces
  rather than being buried. Same for any text OCR simply fails to read: a word
  tesseract does not see is neither in the sidecar nor behind a box.

⚠️ OCR TEXT DOES NOT BEHAVE LIKE PROSE — THE REFLOW IS A DETECTION DECISION
  Read this before "simplifying" `_render_ocr_rows` into `image_to_string()`.

  #162 established that the SHAPE of extracted text decides whether a person's
  name is detected at all. That applies with more force here, because tesseract
  does not hand back the shape the screenshot shows. On the default page
  segmentation it detects COLUMNS as separate blocks and emits each column
  whole, one after another, so a stakeholder table comes out as every name,
  then every role, then every email — the row association is destroyed before
  the detector ever runs, and a form's labels are separated from their values.

  Measured (30 synthetic names x 2 independent name sets x 3 screenshot
  layouts, `en_core_web_lg`; "row-fid" = the name and its own row's other
  fields landing in one output line):

      layout   pipeline                                   PERSON   row-fid
      -------- ------------------------------------------ -------  -------
      table    raw image_to_string (tesseract's own)        27/29      0/29
      table    visual rows, joined with spaces              20/29     29/29
      table    visual rows, cells comma-joined + terminator 29/29     29/29
      form     raw image_to_string (tesseract's own)        27/29         -
      form     visual rows, joined with spaces              21/29         -
      form     visual rows, cells comma-joined + terminator 29/29         -
      cards    raw image_to_string (tesseract's own)        24/29         -
      cards    visual rows, cells comma-joined + terminator 29/29         -
      prose    raw image_to_string (tesseract's own)        28/28     28/28
      prose    visual rows, cells comma-joined + terminator 28/28     28/28

  The middle row of each block is the trap: rebuilding the visual row recovers
  the MEANING (row-fid 0/29 -> 29/29) and simultaneously loses SEVEN names,
  because a rebuilt table row is exactly the shape #159 measured spaCy tagging
  ORGANIZATION. Reflowing without also fixing the shape would have been a
  strict privacy regression sold as a readability win.

  CHOSEN: rebuild visual rows from word geometry, split each row into cells on
  a column-sized horizontal gap, join the cells with ", " and terminate the
  line (`_ocr_rows` + `_row_cells` + `_render_ocr_rows`). It reads as
  apposition — `Aisha Rahman, Chief Financial Officer, a.rahman@...` — which
  is both what the screenshot means and a shape the detector handles: 29/29,
  30/30, 29/29 and 28/28 across the four measurements above, never below the
  raw baseline on any layout. Nothing is invented: no header inference, no
  field labels, no reordering. Label-line shapes (`Full name` / value) fall
  out of the same rule as `Full name, Aisha Rahman.`

  Rejected alternatives, all measured: `Field N:` positional labels (28/29 —
  worse, and invents labels), header-row inference into #162's record shape
  (29/29 on a table but scrambles a form, pairing `Full name: Job title`),
  `- ` bullets (23/29), space-joined rows (20/29).

  RESIDUAL GAP, NOT HIDDEN: a name tesseract MISREADS ("Vikram lyer" for
  "Vikram Iyer") is a name the sidecar carries in corrupted form. It is still
  recognisable to a human. Detection usually still fires on the misread form —
  in which case the pixels are blanked too — but this is OCR recall, and no
  reflow fixes it.

  KNOWN OVER-REDACTION, ACCEPTED: comma-joining a header row of short
  capitalised nouns can read as a list of people — `Name, Role, Email,
  Account.` has spaCy tagging "Email" (and, in context, "Account") PERSON, so
  a column HEADER is replaced by a placeholder and boxed in the image. This
  fails SAFE — a label is lost, never a value — and the entity type in each
  value's own placeholder still tells a reader which column was which. The
  real fix is adding generic column-header words to the engine's allow-list,
  which lives in `engine.py`/`denylist.py`; #163 is forbidden from touching
  either, so it is backlogged rather than bodged here.

⚠️ NON-LATIN SCRIPT IN IMAGES — MEASURED, AND FIXED BY REFUSAL, NOT TRANSLATION
  Only the `eng` language pack is installed (plus `osd`/`snum`, which are
  orientation and number models, not languages). Measured on rendered
  screenshots:

    script                        OCR read as                    detected
    ----------------------------- ------------------------------ ---------
    Sinhala  "නිමල් පෙරේරා"          "HO: 8OE Gs5ebd)"              nothing
    Tamil    "ரமேஷ் குமார்"           "Quwj: ona Gwny"               garbage
    Devanagari "प्रिया शर्मा"            "ava: ftrar rat"               garbage
    Filipino "Maria Clara Santos"  read exactly                   PERSON ✓
    Latin name in an English UI    read exactly                   PERSON ✓

  Latin script is fine, including local names written in Latin — which is what
  a banking admin console actually shows in Colombo, Mumbai and Manila, and
  emails/account numbers stay Latin even inside a non-Latin UI (all three
  non-Latin cases above still had their email read and detected correctly).

  A name in a NON-LATIN script is transliterated into noise. It is therefore
  not in the sidecar in any readable form — no leak there — but until #173 it
  was also NOT DETECTED, which meant its pixels were never blanked and it
  stayed fully legible in the redacted image copy: an artifact an agent may
  open, asserting a scrub it had not performed.

  ⚠️ DO NOT "FIX" THIS BY INSTALLING LANGUAGE PACKS. It looks like the fix —
  `brew install tesseract-lang` — and it is exactly backwards. Today the
  redacted copy is accidentally safe only because OCR yields garbage that the
  detector also cannot read as noise: no leak in the sidecar, but the pixels
  stay unblanked. Transcribe the name CORRECTLY and it lands in the sidecar in
  CLEARTEXT, where `en_core_web_lg` still will not detect it — non-English NER
  is explicitly out of scope for v6 (PRD §6) — and the image is STILL
  unredacted, because detection (not OCR) is what drives which pixels get
  boxed. That trade makes a hidden leak into a plainer one. Both artifacts
  need fixing together, later, or not at all.

  #173 FIX — REFUSE THE WHOLE IMAGE when OCR is not confident, recognisable
  text, rather than pass it through unredacted. The signal is tesseract's own
  per-word confidence from `pytesseract.image_to_data` (already parsed for the
  bounding boxes — see `_ocr_words`): `_mean_confidence` averages it across
  every recognised word, and `_MIN_OCR_CONFIDENCE` (60.0, on tesseract's own
  0-100 scale) is the refusal floor. A CHARACTER-RANGE heuristic ("proportion
  of non-Latin characters tesseract could not map") was considered and
  REJECTED as a primary signal: `eng.traineddata` always emits from ITS OWN
  ASCII-only alphabet regardless of what script it was shown — the Sinhala
  example above OCRs to `"HO: 8OE Gs5ebd)"`, not to any non-Latin codepoint —
  so a character-range check on the OUTPUT text is blind to the exact failure
  this exists to catch. Confidence is not blind to it, and one clean signal
  beats two where the second contributes nothing measured.

  CALIBRATED, BOTH DIRECTIONS, on the fixtures this ticket was required not to
  break (`_image_font(20)` — the same Pillow-bundled, OS-independent face
  every OCR fixture in this module already uses, so these numbers reproduce
  identically on a laptop and in CI):

    fixture                                            mean confidence
    --------------------------------------------------- ---------------
    table layout (#163's own screenshot fixture)              91.5
    form layout (label: value, stacked)                       95.4
    cards layout (3 boxed records side by side)                89.1
    prose layout (wrapped paragraph)                           91.8
    small-font (13px) table                                    82.6
    small-font (13px) prose                                    91.6
    "Nimal Perera" alone, Latin, English UI                    93.0
    "Maria Clara Santos" alone, Latin, English UI               95.3
    ---------------------------------------------------  ---- REFUSE FLOOR (60.0) ----
    non-Latin UI labels + Latin email/account (mixed)          38.6
    Devanagari name alone                                       20.0
    Sinhala name alone                                           7.0
    Tamil name alone                                             7.0

  60.0 sits ~23 points below the WORST passing fixture (82.6) and ~22 points
  above the WORST refusing fixture (38.6) — comfortable margin on both sides,
  not a hairline. Every fixture #163 proved working (four layouts, two Latin
  names, two font sizes) stays above the floor; every non-Latin fixture
  measured falls below it.

  THE MIXED CASE, NAMED EXPLICITLY: a screenshot with non-Latin UI labels
  around a Latin email and account number scores 38.6 and is REFUSED, even
  though the email (89 conf) and account number (91 conf) were individually
  read correctly — per-word confidence on those two values is fine; it is the
  surrounding label noise (confidence 0-40) that drags the image mean below
  the floor. This is deliberate, not a false positive: the fixture that
  produces this shape is exactly a non-Latin label next to a Latin-script
  NAME (see `_build_mixed`-shaped fixtures) — refusing errs safe because nothing
  in this module can tell "harmless non-Latin caption" apart from "a name that
  transliterates into noise" without reading the caption, which is the exact
  thing it just failed to do with confidence. An artifact that cannot prove a
  region has no PII does not get to keep it.

  A GENUINELY TEXTLESS image (a chart, a logo, a photo) is not affected: zero
  recognised words means `_mean_confidence` returns -1 rather than a low
  score, the refusal check is skipped, and the existing `EmptyExtractionError`
  path handles it — an ordinary empty result, not a refusal. The two must stay
  distinguishable: one says "there was nothing to read here", the other says
  "there was something here and it could not be trusted."

D5b DECISION — `presidio-image-redactor` 0.0.60 REJECTED
  PRD §4 and the ticket asked for this to be settled on evidence. Installed
  and run head-to-head, not judged on its "beta, not production ready" label:

  1. **Its analysis text is one flat space-joined line.** `OCR.get_text_from_
     ocr_dict` is literally `separator.join(ocr_result["text"])` with
     `separator=" "`, and `ImageAnalyzerEngine.analyze` calls it with the
     default — no newlines, no rows, no terminators, and no injection point.
     That is the single worst shape for the one detector that matters here.
     Measured on the same fixture, one session, one deny-list:

         presidio-image-redactor's shape (flat space-join) : PERSON 6/8
         this module's reflow                              : PERSON 8/8

     Both scored EMAIL 8/8 and CLIENT 1/1 — the gap is entirely PERSON, i.e.
     entirely the shape, exactly as #162 predicted.
  2. **It cannot produce the sidecar.** It draws boxes and returns an image;
     the placeholder text that carries the round-trip would still have to come
     from our own OCR pass — so the boxes and the sidecar would come from two
     independent OCR runs that can disagree about what the image says.
  3. **It hard-depends on `azure-ai-formrecognizer`** (not an extra — a plain
     `Requires-Dist`), the Azure Document Intelligence client this module is
     forbidden to use, plus `opencv-python`, `matplotlib`, `pydicom` and
     `python-gdcm`. Adding a cloud OCR client to the dependency tree of a
     local-only privacy control is not a trade worth making for a component we
     would be fighting on shape anyway.

  PATH TAKEN (as the ticket prescribes): boxes are drawn directly from
  tesseract's own word bounding boxes intersected with the engine's character
  spans — we already have both, from one pass. `presidio_image_redactor` is
  not imported by this module and is not in `requirements.txt`.

INTERPRETER
  This module is importable and its EXTRACTORS are runnable on plain Python
  3.9 (matching `scripts/artifact_boundary.py`'s importability contract):
  every optional library is imported LAZILY inside its own extractor branch,
  and `scripts/pii/engine.py` is imported lazily too, only on the
  anonymisation path. `import scripts.pii.ingest` pulls in nothing but the
  standard library. Anonymisation itself needs `.venv` (Presidio, 3.10-3.13).
"""
from __future__ import annotations

import csv
import io
import re
import sys
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

__all__ = [
    "SUPPORTED_SUFFIXES",
    "DOCUMENT_SUFFIXES",
    "IMAGE_SUFFIXES",
    "IMAGE_SEAM_MARKER",
    "LOGO_NOTICE",
    "IngestError",
    "UnsupportedFormatError",
    "ExtractorUnavailableError",
    "OCRUnavailableError",
    "OCRLowConfidenceError",
    "ExtractionFailedError",
    "EmptyExtractionError",
    "IngestResult",
    "anon_path_for",
    "redacted_image_path_for",
    "extract_text",
    "ingest_file",
    "is_supported",
    "reset_engagement_notices",
]

# Document formats this module converts. Keep in sync with the extractor
# dispatch in `extract_text` and with the eval's per-format fixture set.
DOCUMENT_SUFFIXES: Tuple[str, ...] = (".pdf", ".docx", ".pptx", ".xlsx", ".csv")

# Image formats. Every one of these goes through the OCR path (#163) and
# produces BOTH a text sidecar and a redacted copy.
IMAGE_SUFFIXES: Tuple[str, ...] = (
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic",
)

SUPPORTED_SUFFIXES: Tuple[str, ...] = DOCUMENT_SUFFIXES + IMAGE_SUFFIXES

# Left in the text where an embedded image sits IN DOCUMENT ORDER when OCR
# read nothing out of it — a photo, a chart, a logo, or a picture the OCR
# could not be run on. Visible on purpose: a silent drop would make an
# image-only slide indistinguishable from an empty one. When OCR DOES read
# text, `_IMAGE_TEXT_OPEN`/`_IMAGE_TEXT_CLOSE` bracket it instead.
IMAGE_SEAM_MARKER = "[embedded image — no text could be read from this picture]"

_IMAGE_TEXT_OPEN = "[embedded image %d — text read from the picture]"
_IMAGE_TEXT_CLOSE = "[end of embedded image %d]"

# ux-design-v6.md Flow C, VERBATIM. Emitted on the FIRST image ingest of an
# engagement — once per engagement, not once per image (see
# `_emit_logo_notice`). This is the accepted-limitation disclosure the PRD
# requires to surface rather than be buried: redaction blanks what OCR reads
# as text, and a logo is not text.
LOGO_NOTICE = """ℹ️  About screenshots

   Cortex blanks out any client details it can *read* in an image —
   names, emails, account numbers.

   It cannot blank a logo. A logo is a picture, not text, so it stays
   visible and reaches Claude.

   If a screenshot shows the client's logo, crop it out before adding it."""


# --- typed errors ----------------------------------------------------------

class IngestError(RuntimeError):
    """Base for every ingest failure.

    Carries the offending `path` and `format` so a caller — the guard, most
    of all — can render a plain-language message without re-deriving them
    from the exception text.
    """

    def __init__(self, message: str, *, path=None, fmt: Optional[str] = None):
        super().__init__(message)
        self.message = message
        self.path = Path(path) if path is not None else None
        self.format = fmt

    def __str__(self) -> str:  # keep the plain-language text as the str form
        return self.message


class UnsupportedFormatError(IngestError):
    """The file's format has no extractor. Names the format in the message."""


class ExtractorUnavailableError(IngestError):
    """The extractor's library is not installed under this interpreter."""


class OCRUnavailableError(ExtractorUnavailableError):
    """The local OCR binary is missing, so an image cannot be read.

    REFUSES. The raw image is never passed through — an unreadable image is
    an unscrubbable image, and the whole point of this path is that nothing
    unscrubbed reaches the model. `.message` is the plain-language text the
    session preflight (`.claude/hooks/pii-preflight.sh`) prints for the same
    condition; the two must not drift.

    A DOCUMENT that merely CONTAINS a picture is treated differently and on
    purpose: it degrades to `IMAGE_SEAM_MARKER` rather than refusing the whole
    file. The picture itself still never reaches the model — the agent reads
    the `.anon_` text, never the source deck — so refusing an entire annual
    report over one logo would cost the consultant everything and protect
    nothing.
    """


class OCRLowConfidenceError(IngestError):
    """The local OCR ran but could not read the image's text with enough
    confidence to trust it — refuses. (#173, closing the #163 non-Latin-script
    leak: see the module docstring's "NON-LATIN SCRIPT IN IMAGES" section.)

    THE LEAK THIS CLOSES: a non-Latin-script name (Sinhala, Tamil,
    Devanagari — the only scripts measured, not necessarily the only ones
    affected) transliterates into ASCII noise the detector cannot recognise
    as a name. Before this error existed, that meant no box was drawn, so the
    real name stayed fully legible in the redacted image copy — an artifact
    whose `.anon_` prefix asserts it has been scrubbed. Refusing the whole
    image is the fix: no sidecar and no redacted copy are written, so nothing
    claims a protection it did not provide.

    A DIRECT image ingest raises this and REFUSES — see `_extract_image`.
    An image EMBEDDED in a DOCX/PPTX does NOT raise it: `_ocr_embedded`
    degrades to `IMAGE_SEAM_MARKER` instead, for the same reason a missing
    OCR binary or a corrupt embedded picture already degrades rather than
    refusing the whole host document (see `OCRUnavailableError`'s own note) —
    but it also withholds that picture's redacted copy, for the identical
    reason: an embedded picture gets its own `.anon_...imageN.png`, which an
    agent can open directly, and that copy would carry the same unblanked
    leak if it were written from low-confidence OCR.

    Distinct from `OCRUnavailableError`: that is raised when the local OCR
    tool cannot run at all. This is raised when it DID run, on real image
    data, and produced output not trustworthy enough to detect against.
    """


class ExtractionFailedError(IngestError):
    """The document could not be read — encrypted, password-protected or
    corrupt. Distinct from `EmptyExtractionError`: the file was rejected,
    not merely textless."""


class EmptyExtractionError(IngestError):
    """The document was read successfully but yielded no text at all (a
    scanned PDF, an image-only deck). Never reported as a successful scrub —
    an empty `.anon_` file would be indistinguishable from a clean one."""


# --- rendering: the detection-driven shape ---------------------------------
#
# See the module docstring's measurement table. Do not replace this with a
# markdown pipe table.

def _clean_cell(value) -> str:
    """One cell as a single-line string. Newlines inside a cell would break
    the one-field-per-line contract the detector (and denylist.py's
    label-line extractor) relies on."""
    if value is None:
        return ""
    text = value if isinstance(value, str) else str(value)
    return re.sub(r"\s+", " ", text).strip()


def _terminate(value: str) -> str:
    """End a field value with sentence punctuation.

    Not cosmetic — measured. Without a terminator, spaCy runs the PERSON span
    on past the newline and swallows the NEXT field's label: `Name: Aisha
    Rahman` / `Role: CFO` anonymises to `Name: <PERSON_2>: Chief Financial
    Officer`, because the detected entity was literally `"Aisha Rahman\\nRole"`.
    Detection is still correct and the round-trip is still exact, but the
    record loses a label, and the placeholder is bound to a value that would
    not match the same person written anywhere else — breaking the cross-file
    placeholder identity this module exists to preserve.

    With the terminator: 30/30 names detected, 0 run-on spans, and every
    entity key is the clean value (`'Aisha Rahman'`, not `'Aisha Rahman\\nRole'`)
    — verified to leave email/phone/account spans byte-identical to the
    unterminated rendering, so a value seen in a spreadsheet and in a
    transcript still shares one placeholder.
    """
    value = value.rstrip(" ,;:")
    if not value:
        return value
    if value[-1] in ".!?":
        return value
    return value + "."


def _render_records(
    headers: Sequence[str],
    rows: Iterable[Sequence[object]],
    *,
    label: str = "Record",
) -> List[str]:
    """Tabular data as record-per-row label lines — the measured-best shape.

        ### Record 1
        Name: Aisha Rahman.
        Role: CFO.

        ### Record 2
        ...

    Empty rows are dropped; empty cells are dropped from their record (a
    `Role:` line with nothing after it is noise, and a bare label is not a
    detection surface). A column with no header takes a positional name so
    every value stays attributable.
    """
    safe_headers = [
        _clean_cell(h) or "Column %d" % (i + 1) for i, h in enumerate(headers)
    ]
    lines: List[str] = []
    index = 0
    for row in rows:
        cells = [_clean_cell(c) for c in row]
        if not any(cells):
            continue
        index += 1
        lines.append("### %s %d" % (label, index))
        for i, cell in enumerate(cells):
            if not cell:
                continue
            header = safe_headers[i] if i < len(safe_headers) else "Column %d" % (i + 1)
            lines.append("%s: %s" % (header, _terminate(cell)))
        lines.append("")
    return lines


def _normalise(lines: Iterable[str]) -> str:
    """Join rendered lines into the final text: CRLF normalised, trailing
    whitespace stripped per line, at most one trailing newline. Determinism
    depends on this — same bytes in, same bytes out."""
    out: List[str] = []
    for line in lines:
        for part in str(line).replace("\r\n", "\n").replace("\r", "\n").split("\n"):
            out.append(part.rstrip())
    # Collapse 3+ blank lines to 2; a document's own spacing is not signal.
    text = "\n".join(out)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip("\n") + "\n"


# --- OCR: words, geometry, and the measured reflow -------------------------
#
# See the module docstring's OCR measurement table. Do not replace any of this
# with `pytesseract.image_to_string()` — that is the 20-27/29 baseline, and it
# also throws away the bounding boxes the redacted copy is drawn from.

# Two words on the same visual row when their vertical centres are within this
# fraction of a word's own height. Loose enough to survive the baseline jitter
# tesseract reports for mixed font sizes in one row, tight enough not to merge
# adjacent rows of body text.
_ROW_TOLERANCE = 0.6

# A horizontal gap wider than this many word-heights is a COLUMN boundary, not
# a word space. A space at 20px type is ~7px and a column gap is ~100px, so the
# threshold sits in a wide empty band; it is expressed in word-heights so it
# scales with the screenshot's resolution instead of assuming one.
_COLUMN_GAP_RATIO = 1.2


class _Word(object):
    """One OCR'd word and where it sits in the image."""

    __slots__ = ("text", "left", "top", "width", "height", "confidence", "start", "end")

    def __init__(self, text: str, left: int, top: int, width: int, height: int,
                 confidence: float):
        self.text = text
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        # tesseract's own 0-100 confidence for THIS word. Never used to drop
        # or filter the word itself (see `_ocr_words`'s own note: dropping a
        # low-confidence word would un-blank it) — carried only so the WHOLE
        # image's confidence can be judged before any word is rendered or
        # redacted at all. See `_mean_confidence` / `_MIN_OCR_CONFIDENCE`.
        self.confidence = confidence
        # Filled in by `_render_ocr_rows`: this word's character span in the
        # rendered text. That span is what turns an engine detection back into
        # a rectangle of pixels.
        self.start = -1
        self.end = -1


def _require_ocr(path, fmt: str):
    """Import the OCR binding and confirm the local binary answers.

    Raises `OCRUnavailableError` — never returns a half-working engine, and
    never falls back to any hosted OCR service.
    """
    try:
        import pytesseract  # noqa: PLC0415 - lazy by contract
        from PIL import Image  # noqa: PLC0415
    except ImportError as exc:
        raise OCRUnavailableError(_OCR_MISSING_MESSAGE, path=path, fmt=fmt) from exc
    try:
        pytesseract.get_tesseract_version()
    except Exception as exc:  # noqa: BLE001 - pytesseract raises its own TesseractNotFoundError
        raise OCRUnavailableError(_OCR_MISSING_MESSAGE, path=path, fmt=fmt) from exc
    return pytesseract, Image


# Mirrors .claude/hooks/pii-preflight.sh's "can't process screenshots yet"
# notice — same consequence, same one command. Copy rules (ux-design-v6.md):
# consequence before instruction, say whether they're blocked, name the file.
_OCR_MISSING_MESSAGE = (
    "This image has NOT been prepared for review and must not be opened — the "
    "tool that reads text out of pictures is not set up on this computer, so "
    "nothing in it could be blanked out. Documents still work normally. "
    "Install it once with `brew install tesseract` (macOS) or "
    "`sudo apt install tesseract-ocr` (Linux), then try again."
)

# Refusal floor for `_mean_confidence` — tesseract's own 0-100 word-confidence
# scale from `image_to_data`. Below this, the image is refused outright rather
# than passed through with an unredacted leak (#173; see the module
# docstring's "NON-LATIN SCRIPT IN IMAGES" for the full calibration table).
# 60.0 sits ~23 points below the worst-scoring fixture that MUST keep working
# (82.6, a small-font table) and ~22 points above the worst-scoring fixture
# that MUST refuse (38.6, a non-Latin UI with a Latin email/account still
# individually legible) — comfortable margin on both sides, not a hairline.
_MIN_OCR_CONFIDENCE = 60.0

# Copy rules (ux-design-v6.md): consequence before instruction, say whether
# they're blocked, no tool names. Mirrors `_OCR_MISSING_MESSAGE`'s shape for
# the sibling failure mode: that one is the tool being absent; this one is the
# tool running and not being trustworthy.
_OCR_LOW_CONFIDENCE_MESSAGE = (
    "The text in this image could not be read clearly enough to check it for "
    "client details, so this image has NOT been prepared for review and must "
    "not be opened. This usually happens when the image contains text in a "
    "script the reader does not recognise. Describe what the image shows in "
    "a sentence instead, or provide the same information as text or a "
    "spreadsheet, and use that."
)


def _mean_confidence(words: Sequence[_Word]) -> float:
    """Mean per-word confidence across every word THIS image's OCR pass
    recognised, on tesseract's own 0-100 scale. -1.0 when no words were
    recognised at all — deliberately NOT a low score: a genuinely textless
    image (a chart, a logo, a photo) must fall through to the ordinary
    `EmptyExtractionError` path, not be refused as unreadable. See
    `_MIN_OCR_CONFIDENCE` and the module docstring's calibration table.
    """
    if not words:
        return -1.0
    return sum(w.confidence for w in words) / len(words)


def _ocr_words(pytesseract, image) -> List[_Word]:
    """Every word tesseract read, with its box, in tesseract's own order.

    NO per-word confidence threshold is applied to WHICH words are kept —
    dropping a low-confidence word would drop it from the sidecar AND from
    the redaction pass — the box is only drawn for text the detector saw —
    so a half-read name would go from "garbled but blanked" to "garbled and
    visible". Words with conf == -1 are tesseract's structural rows
    (page/block/paragraph/line), not text, and are the only thing skipped.

    Each word's confidence IS kept (`_Word.confidence`), read by
    `_mean_confidence` to judge the WHOLE image before anything is rendered —
    a different question from which words to keep.
    """
    from pytesseract import Output  # noqa: PLC0415 - lazy by contract

    data = pytesseract.image_to_data(image, output_type=Output.DICT)
    words: List[_Word] = []
    for i in range(len(data["text"])):
        text = (data["text"][i] or "").strip()
        if not text:
            continue
        try:
            confidence = float(data["conf"][i])
        except (TypeError, ValueError):
            confidence = -1.0
        if confidence < 0:
            continue
        words.append(_Word(
            text,
            int(data["left"][i]), int(data["top"][i]),
            int(data["width"][i]), int(data["height"][i]),
            confidence,
        ))
    return words


def _ocr_rows(words: Sequence[_Word]) -> List[List[_Word]]:
    """Rebuild the image's VISUAL rows from word geometry.

    Not tesseract's own line grouping — that is per-block, and on the default
    page segmentation tesseract treats each COLUMN of a table as its own
    block, so its "lines" run down a column instead of across a row. Grouping
    on the y-centre across the whole page puts a row back together.

    Trade-off, stated: on a genuinely multi-column page (a two-column report
    scan) this interleaves the columns, where tesseract's block order would
    have kept them apart. Screenshots and forms — what actually lands in
    `inputs/` — are the case this optimises for, and the interleaved text is
    still fully scrubbed; only its reading order suffers.
    """
    rows: List[Dict] = []
    for word in sorted(words, key=lambda w: (w.top + w.height / 2.0, w.left)):
        centre = word.top + word.height / 2.0
        for row in rows:
            if abs(centre - row["centre"]) <= _ROW_TOLERANCE * max(word.height, 1):
                row["words"].append(word)
                row["centre"] = sum(
                    w.top + w.height / 2.0 for w in row["words"]
                ) / len(row["words"])
                break
        else:
            rows.append({"centre": centre, "words": [word]})
    rows.sort(key=lambda r: (r["centre"], min(w.left for w in r["words"])))
    return [sorted(row["words"], key=lambda w: w.left) for row in rows]


def _row_cells(row: Sequence[_Word]) -> List[List[_Word]]:
    """Split one visual row into cells wherever a column-sized gap appears."""
    if not row:
        return []
    heights = sorted(w.height for w in row)
    median_height = heights[len(heights) // 2] or 1
    cells: List[List[_Word]] = [[row[0]]]
    for previous, word in zip(row, row[1:]):
        gap = word.left - (previous.left + previous.width)
        if gap > _COLUMN_GAP_RATIO * median_height:
            cells.append([word])
        else:
            cells[-1].append(word)
    return cells


def _render_ocr_rows(rows: Sequence[Sequence[_Word]]) -> str:
    """The measured-best shape: one line per visual row, cells comma-joined,
    line terminated. Records each word's character span on the word itself.

    `Aisha Rahman, Chief Financial Officer, a.rahman@example.test.`

    That reads as apposition, which is what the row means AND the shape the
    detector handles (29/29 vs 20/29 for a space-joined row and 27/29 for
    tesseract's own output — module docstring). The terminator is the same
    fix `_terminate` documents for #162: without it spaCy runs the PERSON span
    past the newline and binds the placeholder to a value that would never
    match the same person written anywhere else.
    """
    lines: List[str] = []
    offset = 0
    for row in rows:
        cells = _row_cells(row)
        if not cells:
            continue
        parts: List[str] = []
        for cell in cells:
            cell_start = offset + sum(len(p) for p in parts) + 2 * len(parts)
            piece_offset = cell_start
            for i, word in enumerate(cell):
                if i:
                    piece_offset += 1  # the joining space
                word.start = piece_offset
                word.end = piece_offset + len(word.text)
                piece_offset = word.end
            parts.append(" ".join(w.text for w in cell))
        line = ", ".join(parts)
        terminated = _terminate(line)
        # `_terminate` only ever strips trailing punctuation/whitespace or
        # appends a full stop, so no word's recorded span can move. Clamp any
        # word the strip truncated rather than letting a stale span point past
        # the end of the line.
        limit = offset + len(terminated)
        for word in row:
            if word.end > limit:
                word.end = max(word.start, limit)
        lines.append(terminated)
        offset += len(terminated) + 1  # + the newline
    return "\n".join(lines) + ("\n" if lines else "")


def _ocr_image(pytesseract, image) -> Tuple[str, List[_Word], float]:
    """One OCR pass -> (rendered text, words carrying their spans and boxes,
    mean word confidence).

    ONE pass, deliberately: the sidecar text and the redaction boxes must
    describe the same reading of the same image. Two independent passes could
    disagree about what the picture says, and the disagreement would show up
    as a name that is in the sidecar's placeholders but still legible in the
    redacted copy (or the reverse).

    The confidence is judged on the WORDS AS TESSERACT READ THEM, before
    `_ocr_rows`/`_render_ocr_rows` reflow them into visual rows — the reflow
    changes how text is JOINED, never what was recognised or with what
    confidence, so judging pre- or post-reflow gives the same number; doing it
    here means a caller can decide whether to trust the text before rendering
    or redacting anything from it (see `_MIN_OCR_CONFIDENCE`).
    """
    words = _ocr_words(pytesseract, image)
    rows = _ocr_rows(words)
    return _render_ocr_rows(rows), words, _mean_confidence(words)


def _open_image(Image, data_or_path, path, fmt: str):
    """Open an image for reading, converting failures into typed errors."""
    try:
        if isinstance(data_or_path, (bytes, bytearray)):
            image = Image.open(io.BytesIO(bytes(data_or_path)))
        else:
            image = Image.open(str(data_or_path))
        image.load()
    except Exception as exc:  # noqa: BLE001 - any decode fault becomes actionable
        raise ExtractionFailedError(
            "This %s image could not be opened — the file looks damaged, or it "
            "is saved in a form Cortex cannot read. It has NOT been prepared "
            "for review and must not be opened. Save it as a PNG or JPEG into "
            "the same folder and use that copy instead." % fmt.upper(),
            path=path, fmt=fmt,
        ) from exc
    return image


def _redact_image(Image, image, words: Sequence[_Word], spans: Sequence[Tuple[int, int]]):
    """A copy of `image` with every word overlapping a detected span filled.

    Boxes come from tesseract's own word geometry; the spans come from the
    engine's character offsets into the very text those words rendered. That
    intersection is the whole redaction algorithm — no second OCR pass, no
    second detector, and nothing that could disagree with the sidecar.

    DESTRUCTIVE AND NOT REVERSIBLE, by design (module docstring): the copy is
    never restored and never shipped. Solid black, not a blur — a blur is
    recoverable and reads as an effect rather than a removal.
    """
    from PIL import ImageDraw  # noqa: PLC0415 - lazy by contract

    redacted = image.convert("RGB")
    draw = ImageDraw.Draw(redacted)
    boxed = 0
    for word in words:
        if word.start < 0:
            continue
        if not any(word.start < end and start < word.end for start, end in spans):
            continue
        # Bleed 2px past the word box on every side. Tesseract's box is tight
        # to the ink it recognised, which can sit a pixel or two inside the
        # glyph's real extent — an ascender, a descender, the tail of a "(",
        # or the antialiased edge of any letter. Measured at a 1px bleed, up
        # to 6% of a word's own background was left unfilled at the boundary;
        # at 2px it is none. Over-covering a neighbouring pixel is the safe
        # direction to err in.
        draw.rectangle(
            [word.left - 2, word.top - 2,
             word.left + word.width + 1, word.top + word.height + 1],
            fill=(0, 0, 0),
        )
        boxed += 1
    return redacted, boxed


def _save_png(image, target: Path) -> None:
    """Write the redacted copy. Always PNG, always deterministic — Pillow
    writes no timestamp and no source path into a PNG, so the same pixels
    produce the same bytes on every run."""
    target.parent.mkdir(parents=True, exist_ok=True)
    image.save(str(target), format="PNG", optimize=False)


# --- the logo notice: once per engagement, not once per image --------------

# Engagements already notified in THIS process. Keyed by the resolved
# engagement directory (or, when no engagement was supplied, by the image's
# own directory — the closest available stand-in for "this engagement").
#
# Process-scoped on purpose. The alternative, a marker file inside the
# engagement, would make a one-off disclosure permanent state on disk that
# nothing ever cleans up, and would mean a consultant who comes back in six
# months and adds a screenshot is never told again. One notice per run of the
# ingest is the behaviour the UX asks for: not once per image.
_NOTIFIED_ENGAGEMENTS = set()


def reset_engagement_notices() -> None:
    """Forget which engagements have seen the logo notice.

    For tests and for a long-lived process that legitimately starts a fresh
    engagement — never call it between images of one engagement, which is
    exactly the per-image spam the notice is scoped to avoid.
    """
    _NOTIFIED_ENGAGEMENTS.clear()


def _emit_logo_notice(engagement_dir, path: Path, stream) -> bool:
    """Print `LOGO_NOTICE` the first time this engagement ingests an image.

    Returns True when it printed. Never raises — a notice that breaks a scrub
    would be worse than a notice nobody saw.
    """
    try:
        key = str(Path(engagement_dir).resolve()) if engagement_dir is not None \
            else str(Path(path).resolve().parent)
    except Exception:  # noqa: BLE001
        key = str(engagement_dir if engagement_dir is not None else Path(path).parent)
    if key in _NOTIFIED_ENGAGEMENTS:
        return False
    _NOTIFIED_ENGAGEMENTS.add(key)
    target = stream if stream is not None else sys.stderr
    try:
        target.write(LOGO_NOTICE.rstrip("\n") + "\n")
        target.flush()
    except Exception:  # noqa: BLE001
        pass
    return True


# --- extractors ------------------------------------------------------------
#
# Each imports its library LAZILY, inside its own branch, so this module stays
# importable on plain Python 3.9 without every optional library present.

class _ImagePayload(object):
    """One picture that was OCR'd, kept so `ingest_file` can write its redacted
    copy from the SAME pass that produced its text.

    `label` is None for an image ingested directly and the 1-based position of
    the picture in document order for one embedded in a DOCX or PPTX. `text` is
    what `_render_ocr_rows` produced for it; `words` carry their spans WITHIN
    that text, which `ingest_file` rebases onto the final document.
    """

    __slots__ = ("label", "image", "words", "text")

    def __init__(self, label, image, words, text):
        self.label = label
        self.image = image
        self.words = words
        self.text = text


def _ocr_embedded(blob: bytes, index: int, path: Path, payloads: Optional[List]) -> List[str]:
    """OCR one embedded picture into lines for the host document's text.

    Degrades to `IMAGE_SEAM_MARKER` — it never raises. A deck is not refused
    because one of its pictures is a logo, an unreadable JPEG, or arrived on a
    machine with no OCR installed: the picture itself never reaches the model
    either way (the agent reads the `.anon_` text, never the source deck), so
    refusing would cost the consultant the whole document and protect nothing.
    The seam says so in the text, where a reader can see it.

    Low-confidence OCR (#173) degrades the SAME way and for the same reason —
    but ALSO withholds the payload, so no redacted copy is written for this
    picture either. An embedded picture's redacted copy is its own file an
    agent can open directly (`.anon_<host>.imageN.png`), so it carries the
    identical leak risk a direct image ingest refuses outright; degrading to
    the seam marker without a payload gets the same "nothing unredacted is
    written" outcome without refusing the whole host document.
    """
    try:
        pytesseract, Image = _require_ocr(path, "image")
        image = _open_image(Image, blob, path, "image")
        text, words, confidence = _ocr_image(pytesseract, image)
    except IngestError:
        return [IMAGE_SEAM_MARKER, ""]
    except Exception:  # noqa: BLE001 - a picture must never break a document
        return [IMAGE_SEAM_MARKER, ""]
    if not text.strip():
        return [IMAGE_SEAM_MARKER, ""]
    if words and confidence < _MIN_OCR_CONFIDENCE:
        return [IMAGE_SEAM_MARKER, ""]
    if payloads is not None:
        payloads.append(_ImagePayload(index, image, words, text))
    return [_IMAGE_TEXT_OPEN % index, "", text.rstrip("\n"), "",
            _IMAGE_TEXT_CLOSE % index, ""]


def _extract_image(path: Path, payloads: Optional[List] = None) -> str:
    """An image -> the OCR text, in the measured shape. Read-only.

    Raises `OCRLowConfidenceError` (#173) when OCR ran but its output is not
    confident enough to trust — BEFORE the payload is appended, so no sidecar
    and no redacted copy is ever written for this image (see `ingest_file`:
    `_extract` runs before any output path is even computed). A genuinely
    textless image (no words at all) is not affected — `_mean_confidence`
    returns -1.0 for it, which is never below `_MIN_OCR_CONFIDENCE`, and it
    falls through to the ordinary `EmptyExtractionError` path instead.
    """
    fmt = path.suffix.lower().lstrip(".") or "image"
    pytesseract, Image = _require_ocr(path, fmt)
    image = _open_image(Image, path, path, fmt)
    text, words, confidence = _ocr_image(pytesseract, image)
    if words and confidence < _MIN_OCR_CONFIDENCE:
        raise OCRLowConfidenceError(_OCR_LOW_CONFIDENCE_MESSAGE, path=path, fmt=fmt)
    if payloads is not None:
        payloads.append(_ImagePayload(None, image, words, text))
    return text


def _extract_pdf(path: Path, payloads: Optional[List] = None) -> str:  # noqa: ARG001
    try:
        from pdfminer.high_level import extract_pages  # noqa: PLC0415 - lazy by contract
        from pdfminer.layout import LAParams, LTTextContainer
        from pdfminer.pdfdocument import PDFPasswordIncorrect
        from pdfminer.pdfparser import PDFSyntaxError
    except ImportError as exc:
        raise ExtractorUnavailableError(
            "This PDF cannot be prepared for review because the PDF reader is "
            "not installed. Run `bash scripts/setup_pii.sh`, then try again.",
            path=path, fmt="pdf",
        ) from exc

    lines: List[str] = []
    try:
        for page_number, page in enumerate(extract_pages(str(path), laparams=LAParams()), 1):
            lines.append("## Page %d" % page_number)
            lines.append("")
            for element in page:
                if isinstance(element, LTTextContainer):
                    lines.append(element.get_text())
            lines.append("")
    except PDFPasswordIncorrect as exc:
        raise ExtractionFailedError(
            "This PDF is password-protected, so nothing could be read from it "
            "and it has not been prepared for review. Save an unprotected copy "
            "into the same folder and use that instead.",
            path=path, fmt="pdf",
        ) from exc
    except PDFSyntaxError as exc:
        raise ExtractionFailedError(
            "This PDF could not be read — the file looks damaged or is not "
            "really a PDF. It has not been prepared for review.",
            path=path, fmt="pdf",
        ) from exc
    except Exception as exc:  # noqa: BLE001 - any reader fault becomes a typed, actionable error
        raise ExtractionFailedError(
            "This PDF could not be read (%s), so it has not been prepared for "
            "review." % type(exc).__name__,
            path=path, fmt="pdf",
        ) from exc
    return _normalise(lines)


def _docx_heading_prefix(style_name: str) -> str:
    """`Heading 2` -> `###` (one deeper than the sheet/page level, so a
    document's own hierarchy nests under the structural markers this module
    emits). Title -> `##`. Anything else -> no prefix."""
    name = (style_name or "").strip().lower()
    if name == "title":
        return "## "
    match = re.match(r"^heading (\d+)$", name)
    if not match:
        return ""
    level = min(int(match.group(1)), 5)
    return "#" * (level + 1) + " "


_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_VML_NS = "{urn:schemas-microsoft-com:vml}"


def _docx_image_blobs(element, document) -> List[bytes]:
    """Bytes of every picture in one Word paragraph, in document order.

    Covers both the modern DrawingML `a:blip` and the legacy VML
    `v:imagedata` — old client decks and RFPs exported from Word 2003
    templates still carry the VML form, and it is the one an agent would
    silently not see.
    """
    blobs: List[bytes] = []
    refs = [(child, _R_NS + "embed") for child in element.iter(_A_NS + "blip")]
    refs += [(child, _R_NS + "id") for child in element.iter(_VML_NS + "imagedata")]
    for child, attribute in refs:
        rel_id = child.get(attribute)
        if not rel_id:
            continue
        try:
            blobs.append(document.part.related_parts[rel_id].blob)
        except Exception:  # noqa: BLE001 - a broken relationship is a seam, not a crash
            continue
    return blobs


def _extract_docx(path: Path, payloads: Optional[List] = None) -> str:
    try:
        import docx  # noqa: PLC0415 - lazy by contract
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise ExtractorUnavailableError(
            "This Word document cannot be prepared for review because the Word "
            "reader is not installed. Run `bash scripts/setup_pii.sh`, then try "
            "again.",
            path=path, fmt="docx",
        ) from exc

    try:
        document = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionFailedError(
            "This Word document could not be read — it may be password-protected "
            "or damaged. It has not been prepared for review.",
            path=path, fmt="docx",
        ) from exc

    lines: List[str] = []
    table_index = 0
    image_index = 0
    # Walk the body in DOCUMENT ORDER. `document.paragraphs` and
    # `document.tables` are separate sequences, so using them would silently
    # reorder a document into "all prose, then all tables" and detach every
    # table from the heading that gives it meaning.
    body = document.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            paragraph = Paragraph(child, document)
            text = _clean_cell(paragraph.text)
            blobs = _docx_image_blobs(child, document)
            if text:
                lines.append(_docx_heading_prefix(paragraph.style.name if paragraph.style else "") + text)
            for blob in blobs:
                image_index += 1
                lines.extend(_ocr_embedded(blob, image_index, path, payloads))
            if text or blobs:
                lines.append("")
        elif tag == "tbl":
            table_index += 1
            table = Table(child, document)
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            if not rows:
                continue
            lines.append("## Table %d" % table_index)
            lines.append("")
            lines.extend(_render_records(rows[0], rows[1:], label="Row"))
    return _normalise(lines)


def _extract_pptx(path: Path, payloads: Optional[List] = None) -> str:
    try:
        from pptx import Presentation  # noqa: PLC0415 - lazy by contract
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise ExtractorUnavailableError(
            "This presentation cannot be prepared for review because the "
            "PowerPoint reader is not installed. Run `bash scripts/setup_pii.sh`, "
            "then try again.",
            path=path, fmt="pptx",
        ) from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionFailedError(
            "This presentation could not be read — it may be password-protected "
            "or damaged. It has not been prepared for review.",
            path=path, fmt="pptx",
        ) from exc

    lines: List[str] = []
    image_index = 0
    for slide_number, slide in enumerate(presentation.slides, 1):
        lines.append("## Slide %d" % slide_number)
        lines.append("")
        table_index = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_index += 1
                try:
                    blob = shape.image.blob
                except Exception:  # noqa: BLE001 - an unreadable picture is a seam
                    blob = None
                if blob is None:
                    lines.append(IMAGE_SEAM_MARKER)
                    lines.append("")
                else:
                    lines.extend(_ocr_embedded(blob, image_index, path, payloads))
                continue
            if getattr(shape, "has_table", False) and shape.has_table:
                table_index += 1
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                if not rows:
                    continue
                lines.append("### Table %d" % table_index)
                lines.append("")
                lines.extend(_render_records(rows[0], rows[1:], label="Row"))
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = _clean_cell("".join(run.text for run in paragraph.runs))
                    if text:
                        lines.append(text)
                lines.append("")
        # Speaker notes carry stakeholder names as often as the slide does.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = _clean_cell(slide.notes_slide.notes_text_frame.text)
            if notes:
                lines.append("### Speaker notes")
                lines.append("")
                lines.append(notes)
                lines.append("")
    return _normalise(lines)


def _extract_xlsx(path: Path, payloads: Optional[List] = None) -> str:  # noqa: ARG001
    try:
        import openpyxl  # noqa: PLC0415 - lazy by contract
    except ImportError as exc:
        raise ExtractorUnavailableError(
            "This spreadsheet cannot be prepared for review because the "
            "spreadsheet reader is not installed. Run `bash scripts/setup_pii.sh`, "
            "then try again.",
            path=path, fmt="xlsx",
        ) from exc

    try:
        # read_only: never touches the original beyond reading it.
        # data_only: formula RESULTS, not formula text — a formula string is
        # not the value a reader (or the detector) needs to see.
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionFailedError(
            "This spreadsheet could not be read — it may be password-protected "
            "or damaged. It has not been prepared for review.",
            path=path, fmt="xlsx",
        ) from exc

    lines: List[str] = []
    try:
        for sheet in workbook.worksheets:
            lines.append("## Sheet: %s" % _clean_cell(sheet.title))
            lines.append("")
            rows = [
                [_clean_cell(c) for c in row]
                for row in sheet.iter_rows(values_only=True)
            ]
            rows = [r for r in rows if any(r)]
            if not rows:
                lines.append("(no data)")
                lines.append("")
                continue
            # First non-empty row is the header row when it looks like one
            # (>=2 labelled columns, none of them purely numeric). Otherwise
            # every row is data and columns take positional names.
            header = rows[0]
            looks_like_header = (
                sum(1 for c in header if c) >= 2
                and not any(re.fullmatch(r"-?\d+(\.\d+)?", c) for c in header if c)
            )
            if looks_like_header:
                lines.extend(_render_records(header, rows[1:], label="Row"))
            else:
                width = max(len(r) for r in rows)
                lines.extend(_render_records([""] * width, rows, label="Row"))
    finally:
        workbook.close()
    return _normalise(lines)


def _extract_csv(path: Path, payloads: Optional[List] = None) -> str:  # noqa: ARG001
    # stdlib only — no lazy import needed, and no pandas (see D11 DECISION).
    try:
        raw = path.read_bytes()
    except OSError as exc:
        raise ExtractionFailedError(
            "This file could not be read, so it has not been prepared for "
            "review.",
            path=path, fmt="csv",
        ) from exc
    text = raw.decode("utf-8-sig", errors="replace")

    delimiter = ","
    try:
        # Deterministic: the sniff reads a fixed prefix of the same bytes, and
        # falls back to comma on any doubt.
        delimiter = csv.Sniffer().sniff(text[:8192], delimiters=",;\t|").delimiter
    except Exception:  # noqa: BLE001 - a sniff failure is normal, not an error
        delimiter = ","

    rows = [row for row in csv.reader(io.StringIO(text), delimiter=delimiter)]
    rows = [[_clean_cell(c) for c in row] for row in rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return ""

    lines: List[str] = ["## %s" % _clean_cell(path.name), ""]
    header = rows[0]
    looks_like_header = (
        sum(1 for c in header if c) >= 2
        and not any(re.fullmatch(r"-?\d+(\.\d+)?", c) for c in header if c)
    )
    if looks_like_header:
        lines.extend(_render_records(header, rows[1:], label="Row"))
    else:
        width = max(len(r) for r in rows)
        lines.extend(_render_records([""] * width, rows, label="Row"))
    return _normalise(lines)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
}
# Every image format goes to the one OCR extractor.
for _suffix in IMAGE_SUFFIXES:
    _EXTRACTORS[_suffix] = _extract_image
del _suffix


# --- public API ------------------------------------------------------------

def is_supported(path) -> bool:
    """True when `extract_text` has an extractor for this file's format."""
    return Path(path).suffix.lower() in _EXTRACTORS


def is_image(path) -> bool:
    """True when this file goes down the OCR path (#163)."""
    return Path(path).suffix.lower() in IMAGE_SUFFIXES


def anon_path_for(path, output_dir=None) -> Path:
    """The `.anon_` text sibling for `path` — `report.pdf` -> `.anon_report.pdf.md`.

    A pure, invertible function of the filename; see the module docstring's
    OUTPUT NAMING note for why the source extension is kept. Images get one of
    these too: the sidecar is what carries the round-trip, not the pixels.
    """
    path = Path(path)
    directory = Path(output_dir) if output_dir is not None else path.parent
    return directory / (".anon_%s.md" % path.name)


def redacted_image_path_for(path, output_dir=None, *, index: Optional[int] = None) -> Path:
    """The redacted image copy for `path` — `shot.png` -> `.anon_shot.png.png`.

    Same invertible rule as `anon_path_for`, so `shot.png` and `shot.jpg`
    cannot collide. `index` names a picture embedded in a document:
    `deck.pptx` picture 1 -> `.anon_deck.pptx.image1.png`. Always PNG (module
    docstring: lossless, and deterministic bytes).
    """
    path = Path(path)
    directory = Path(output_dir) if output_dir is not None else path.parent
    if index is None:
        return directory / (".anon_%s.png" % path.name)
    return directory / (".anon_%s.image%d.png" % (path.name, index))


def extract_text(path) -> str:
    """Document -> plain text, with structure preserved and NO anonymisation.

    Read-only: the original file is never modified. Raises
    `UnsupportedFormatError`, `ExtractorUnavailableError` (including
    `OCRUnavailableError`), `ExtractionFailedError` or `EmptyExtractionError`
    — never returns empty, never silently skips.

    For an image this is the OCR text. `ingest_file` is what also writes the
    redacted copy; extraction stays a pure read.
    """
    return _extract(path)[0]


def _extract(path, payloads: Optional[List] = None):
    """`extract_text`, plus the OCR'd pictures for `ingest_file` to redact.

    Split out so the redacted copy is drawn from the SAME OCR pass that
    produced the text — see `_ocr_image`.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    fmt = suffix.lstrip(".") or "(no extension)"

    if not path.is_file():
        raise ExtractionFailedError(
            "There is no file at %s, so nothing could be prepared for review."
            % path.name,
            path=path, fmt=fmt,
        )

    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise UnsupportedFormatError(
            "%s files cannot be prepared for review, so this one has been left "
            "alone and must not be opened. Save it as PDF, Word, PowerPoint, "
            "Excel, CSV or a picture into the same folder and use that copy "
            "instead." % (fmt.upper() if suffix else "Extensionless"),
            path=path, fmt=fmt,
        )

    text = extractor(path, payloads)
    if not text or not text.strip():
        raise EmptyExtractionError(
            "No text could be read out of this %s file — it may hold only "
            "pictures, scans or empty cells. It has not been prepared for "
            "review and must not be opened." % fmt,
            path=path, fmt=fmt,
        )
    return text, (payloads if payloads is not None else [])


class IngestResult(object):
    """What one ingest produced. `session` is returned so a caller can keep
    feeding the SAME session to the next document — that is what makes
    placeholders consistent across an engagement's whole input set."""

    def __init__(self, source: Path, anon_path: Path, fmt: str,
                 extracted_chars: int, anonymized_chars: int, session,
                 redacted_paths: Optional[List[Path]] = None,
                 regions_redacted: int = 0, logo_notice_shown: bool = False):
        self.source = source
        self.anon_path = anon_path
        self.format = fmt
        self.extracted_chars = extracted_chars
        self.anonymized_chars = anonymized_chars
        self.session = session
        # One entry per picture redacted: the image itself for a direct image
        # ingest, one per embedded picture for a DOCX/PPTX. Empty for a format
        # that carries no pictures.
        self.redacted_paths: List[Path] = list(redacted_paths or ())
        self.regions_redacted = regions_redacted
        self.logo_notice_shown = logo_notice_shown

    @property
    def redacted_path(self) -> Optional[Path]:
        """The redacted copy, for the common case of one image in one file."""
        return self.redacted_paths[0] if self.redacted_paths else None

    @property
    def entity_mapping(self) -> Dict[str, Dict[str, str]]:
        return self.session.entity_mapping

    def mapping_file_dict(self) -> Dict:
        return self.session.mapping_file_dict()

    def __repr__(self) -> str:
        return "IngestResult(%s -> %s, %s, %d chars)" % (
            self.source.name, self.anon_path.name, self.format, self.anonymized_chars,
        )


def ingest_file(
    path,
    *,
    engagement_dir=None,
    session=None,
    entity_mapping: Optional[Dict[str, Dict[str, str]]] = None,
    output_dir=None,
    notice_stream=None,
) -> IngestResult:
    """Extract, anonymise through `scripts/pii/engine.py`, write `.anon_X.md`.

    Exactly one anonymisation scheme: the text goes through a `PIISession`,
    the same one transcripts go through (module docstring). Supply either a
    live `session` (to keep placeholders consistent across documents) or an
    `engagement_dir` (to resolve the deny-list from the engagement's own
    documents, as `anonymize_transcript_file` does).

    The original file is never modified. A failure raises before anything is
    written — no partial or empty `.anon_` artifact is ever left behind.
    """
    path = Path(path)
    payloads: List[_ImagePayload] = []
    text, payloads = _extract(path, payloads)

    # The logo notice fires as soon as an image is actually read — including
    # one embedded in a deck — and once per engagement, not once per image.
    logo_notice_shown = False
    if payloads:
        logo_notice_shown = _emit_logo_notice(engagement_dir, path, notice_stream)

    # The FILENAME is anonymised too, not just the body: `HDFC_Annual_Report.pdf`
    # names the client in the one line a reader is guaranteed to read. Folding
    # it into the same anonymisation pass means the deny-list catches it.
    document = "# Anonymised extract of `%s` (%s)\n\n%s" % (
        path.name, path.suffix.lower().lstrip("."), text,
    )

    session = _resolve_session(session, engagement_dir, entity_mapping)
    # Two analyses of ONE string, deliberately. `session.anonymize` re-runs the
    # analyzer internally and we cannot reach inside it without changing
    # engine.py — but the input is identical and detection is deterministic, so
    # the spans that place the boxes are exactly the spans that placed the
    # placeholders. The alternative (analysing the bare OCR text separately)
    # would let the two drift at the document header's boundary.
    spans = [(r.start, r.end) for r in session.analyze(document)] if payloads else []
    anonymized = session.anonymize(document)

    target = anon_path_for(path, output_dir)
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(anonymized, encoding="utf-8")

    redacted_paths, regions = _write_redacted_copies(
        path, payloads, document, spans, output_dir,
    )

    return IngestResult(
        source=path,
        anon_path=target,
        fmt=path.suffix.lower().lstrip("."),
        extracted_chars=len(text),
        anonymized_chars=len(anonymized),
        session=session,
        redacted_paths=redacted_paths,
        regions_redacted=regions,
        logo_notice_shown=logo_notice_shown,
    )


def _write_redacted_copies(path: Path, payloads, document: str,
                            spans: Sequence[Tuple[int, int]], output_dir):
    """Write one redacted PNG per OCR'd picture. Never touches the original.

    Each picture's words carry spans relative to ITS OWN rendered text, while
    `spans` are offsets into the whole document. Rebasing is done by locating
    that picture's text inside the document rather than by arithmetic on the
    renderer's bookkeeping: `_normalise` rewrites the document after the
    extractor built it (it rstrips lines and collapses blank runs), so any
    offset carried forward from render time would be stale. Searching for the
    text the picture actually contributed cannot go stale.
    """
    if not payloads:
        return [], 0

    from PIL import Image  # noqa: PLC0415 - lazy by contract

    written: List[Path] = []
    total = 0
    for payload in payloads:
        body = payload.text.rstrip("\n")
        base = document.find(body)
        if base < 0:
            # The picture's text did not survive into the document verbatim.
            # Redact nothing rather than guess an offset and black out the
            # wrong pixels — and say so, loudly, rather than write a copy that
            # looks redacted and is not.
            raise ExtractionFailedError(
                "The text read out of a picture in this file could not be "
                "matched back to it, so no safely redacted copy could be made. "
                "It has NOT been prepared for review and must not be opened.",
                path=path, fmt=path.suffix.lower().lstrip("."),
            )
        local = [(max(0, s - base), e - base) for s, e in spans
                 if e > base and s < base + len(body)]
        redacted, boxed = _redact_image(Image, payload.image, payload.words, local)
        target = redacted_image_path_for(path, output_dir, index=payload.label)
        _save_png(redacted, target)
        written.append(target)
        total += boxed
    return written, total


def _resolve_session(session, engagement_dir, entity_mapping):
    """Open (or reuse) the ONE session this text is anonymised through.

    `engine` is imported lazily so this module stays importable — and its
    extractors runnable — on plain Python 3.9, where Presidio cannot import
    at all (engine.py's INTERPRETER SPLIT note).
    """
    if session is not None:
        return session
    from . import engine as _engine  # noqa: PLC0415 - lazy by contract
    if engagement_dir is not None:
        return _engine.PIISession.for_engagement(
            engagement_dir, entity_mapping=entity_mapping
        )
    raise ValueError(
        "ingest_file needs either a `session` or an `engagement_dir` — "
        "anonymisation without a resolved deny-list would produce a vacuous "
        "scrub (see engine.py's EMPTY_DENY_LIST_WARNING)."
    )


# --- CLI -------------------------------------------------------------------

def _main(argv: Optional[Sequence[str]] = None) -> int:
    import argparse  # noqa: PLC0415 - CLI only

    parser = argparse.ArgumentParser(
        description=("Convert a document or image to anonymised text "
                      "(.anon_<name>.md), plus a redacted copy for every image "
                      "(.anon_<name>.png)."),
    )
    parser.add_argument("--file", required=True, help="Path to the document or image")
    parser.add_argument("--engagement-dir", required=True,
                         help="Engagement directory (resolves the client deny-list)")
    parser.add_argument("--output-dir", default=None,
                         help="Where to write the .anon_ file (default: alongside the source)")
    args = parser.parse_args(list(argv) if argv is not None else None)

    try:
        result = ingest_file(
            Path(args.file),
            engagement_dir=Path(args.engagement_dir),
            output_dir=Path(args.output_dir) if args.output_dir else None,
        )
    except IngestError as exc:
        sys.stderr.write("%s\n" % exc.message)
        return 1
    print("Anonymised: %s" % result.anon_path)
    for redacted in result.redacted_paths:
        print("Redacted copy: %s (%d region(s) blanked)"
              % (redacted, result.regions_redacted))
    return 0


if __name__ == "__main__":
    sys.exit(_main())
